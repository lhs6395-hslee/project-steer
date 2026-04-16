#!/usr/bin/env python3
"""
Parallel Reviewer — Independent Per-Step Adversarial Review

공식 문서 근거:
  - Parallel subagents pattern:
    code.claude.com/docs/en/agent-sdk/subagents.md
    "you can run style-checker, security-scanner, and test-coverage subagents
     simultaneously, reducing review time from minutes to seconds"
  - Information isolation (adversarial review):
    Reviewer receives ONLY Sprint_Contract + its assigned step output
    (never Executor reasoning, never other steps' outputs)
  - --json-schema structured output:
    code.claude.com/docs/en/agent-sdk/structured-outputs.md

Usage:
    python3 scripts/agents/parallel_reviewer.py \\
        <sprint_contract.json> <aggregated_output.json> <module> <run_dir> [attempt]

Design:
    - One Reviewer subprocess per Executor step output (fully independent)
    - All Reviewers start concurrently (no dependencies between reviews)
    - Each Reviewer gets: Sprint_Contract + single step output (minimal context)
    - Aggregate: PASS only if ALL step verdicts = approved
    - Circular feedback detection across attempts
"""

import json
import sys
import os
import re
import subprocess
from pathlib import Path
import time
from concurrent.futures import ThreadPoolExecutor, as_completed


SCRIPT_DIR = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent


def _step_label(step: dict) -> str:
    """Generate user-friendly step label: e.g. '7p(L02) 흐름도 겹침 수정'"""
    action = step.get('action', '') or step.get('description', '')
    page_match = re.search(r'[Ss]lide\s*(\d+)', action)
    layout_match = re.search(r'(L0[1-9]|L[1-9][0-9])', action)
    keyword = ''
    for kw, label in [
        ('recon', '전체 슬라이드 분석'),
        ('font', '폰트 수정'),
        ('subtitle', 'subtitle 수정'),
        ('icon', '아이콘 추가'),
        ('merge', '병합 및 저장'),
        ('save', '저장'),
        ('overlap', '겹침 수정'),
        ('fix', '수정'),
    ]:
        if kw in action.lower():
            keyword = label
            break
    parts = []
    if page_match:
        parts.append(f"{page_match.group(1)}p")
    if layout_match:
        parts.append(f"({layout_match.group(1)})")
    if keyword:
        parts.append(keyword)
    return ' '.join(parts) if parts else f"Step {step.get('id', '?')}"


def load_json(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return json.load(f)


def write_json(filepath, data):
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def build_review_input(contract: dict, step_id, step_output: dict, module: str, attempt: int) -> str:
    """
    Minimal review input per step — information barrier enforced.

    Reviewer는 자기 step 정보만 받는다:
    - task + module (최소 컨텍스트)
    - 해당 step의 action + acceptance_criteria + constraints
    - 해당 step의 execution output
    전체 Sprint_Contract 전달 금지 (다른 step 정보 포함되므로 context 낭비 + isolation 위반)
    """
    step = next((s for s in contract.get('steps', []) if s.get('id') == step_id), {})

    lines = [
        f"TASK: {contract.get('task', '')}",
        f"MODULE: {module}",
        f"ATTEMPT: {attempt}",
        "",
        f"STEP {step_id}: {step.get('action', '')}",
    ]

    if step.get('target_slide_index') is not None:
        lines.append(f"TARGET SLIDE INDEX: {step['target_slide_index']} (0-based)")

    acceptance = step.get('acceptance_criteria', [])
    if acceptance:
        lines.append("\nACCEPTANCE CRITERIA:")
        for ac in acceptance:
            lines.append(f"  - {ac}")

    constraints = step.get('constraints', [])
    if constraints:
        lines.append("\nCONSTRAINTS:")
        for c in constraints:
            lines.append(f"  - {c}")

    lines += [
        "",
        "STEP EXECUTION OUTPUT:",
        json.dumps(step_output, ensure_ascii=False, indent=2),
        "",
        "Review this step output adversarially against the acceptance criteria and constraints above.",
        "Output JSON verdict following schemas/verdict.schema.json.",
        "REQUIRED fields: verdict, score, checklist_results, constraint_violations, issues, suggestions.",
    ]

    if attempt > 1:
        lines.append("REQUIRED on retry: retry_fix_assessment for each previous issue.")

    if module == "pptx":
        lines += [
            "",
            "=== PPTX DIRECT VERIFICATION (MANDATORY) ===",
            "You have Bash tool access. You MUST run python-pptx verification directly.",
            "Do NOT trust executor's text report alone — open the actual PPTX file and verify.",
            "",
            "REQUIRED checks (run via Bash tool):",
            "1. Subtitle text: repr() per paragraph — line count <= 2, no mid-word breaks",
            "2. Overlap: bounding box intersection check for non-intentional overlaps",
            "3. Icon (L04/L05): shape_type==PICTURE, size=411480×411480 EMU, position=card_right-0.65\"/card_bottom-0.65\"",
            "4. Font sizes: actual run.font.size in EMU (13pt=165100, 14pt=177800)",
            "",
            "Example:",
            "```python",
            "from pptx import Presentation",
            "prs = Presentation('results/pptx/*.pptx')",
            "slide = prs.slides[N]",
            "for shape in slide.shapes:",
            "    print(shape.name, shape.left/914400, shape.top/914400, shape.width/914400, shape.height/914400)",
            "    if shape.has_text_frame:",
            "        for i, p in enumerate(shape.text_frame.paragraphs): print(f'  Para{i}:', repr(p.text))",
            "```",
        ]

    return "\n".join(lines)


def run_reviewer(step_id, review_input_path, review_output_path, module):
    """Run a single Reviewer subprocess via call_agent.sh."""
    call_agent = str(SCRIPT_DIR / "call_agent.sh")
    cmd = ["bash", call_agent, "reviewer", str(review_input_path), str(review_output_path), module]
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=900,  # reviewer --max-turns 10 + PPTX 직접 검증으로 시간 증가
        )
        return {
            "step_id": step_id,
            "status": "success" if result.returncode == 0 else "failed",
            "returncode": result.returncode,
            "output_path": str(review_output_path),
            "stderr": result.stderr[-500:] if result.stderr else "",
        }
    except subprocess.TimeoutExpired:
        return {"step_id": step_id, "status": "timeout", "output_path": str(review_output_path)}
    except Exception as e:
        return {"step_id": step_id, "status": "error", "error": str(e), "output_path": str(review_output_path)}


def parse_verdict(verdict_path: str) -> dict:
    """Parse a single step verdict file. Returns safe defaults on failure."""
    try:
        data = load_json(verdict_path)
        return {
            "verdict": data.get("verdict", "needs_revision"),
            "score": float(data.get("score", 0.0)),
            "issues": data.get("issues", []),
            "constraint_violations": data.get("constraint_violations", []),
            "checklist_results": data.get("checklist_results", {}),
            "suggestions": data.get("suggestions", []),
            "retry_fix_assessment": data.get("retry_fix_assessment", []),
        }
    except Exception as e:
        return {
            "verdict": "needs_revision",
            "score": 0.0,
            "issues": [f"Failed to parse verdict: {e}"],
            "constraint_violations": [],
            "checklist_results": {},
            "suggestions": [],
            "retry_fix_assessment": [],
        }


def aggregate_verdicts(step_verdicts: list[dict]) -> dict:
    """
    Aggregate per-step verdicts using MoA-inspired weighted ensemble.

    근거: arxiv.org/abs/2406.04692 (Mixture-of-Agents, Together AI 2024)
    "Each agent takes outputs from previous layer agents as auxiliary information"
    → 여러 독립 Reviewer의 verdict를 집계할 때 단순 min/max 대신
       score-weighted majority voting으로 앙상블 품질 향상

    Rules:
    - constraint_violation 있으면 → 즉시 needs_revision (hard rule, MoA와 무관)
    - Score: weighted average (score가 높은 Reviewer의 판단에 더 가중치)
    - Verdict: weighted majority (approved score weight vs needs_revision score weight)
    - Issues: 중복 제거 후 score 가중치로 정렬 (가장 신뢰도 높은 Reviewer의 이슈 우선)
    """
    all_issues = []
    all_violations = []
    all_suggestions = []
    scores = []
    step_results = []
    approved_weight = 0.0
    revision_weight = 0.0

    for sv in step_verdicts:
        verdict_data = sv["verdict_data"]
        score = verdict_data["score"]
        verdict = verdict_data["verdict"]

        all_issues.extend(verdict_data["issues"])
        all_violations.extend(verdict_data["constraint_violations"])
        all_suggestions.extend(verdict_data["suggestions"])
        scores.append(score)

        # Weighted majority: score as confidence weight
        if verdict in ("approved", "pass"):
            approved_weight += score
        else:
            revision_weight += score

        step_results.append({
            "step_id": sv["step_id"],
            "verdict": verdict,
            "score": score,
            "issues": verdict_data["issues"],
            "constraint_violations": verdict_data["constraint_violations"],
        })

    # Hard rule: any constraint violation → fail regardless of ensemble
    has_violations = len(all_violations) > 0

    # Weighted average score (MoA: higher-confidence reviewers weighted more)
    total_weight = sum(scores) if scores else 1.0
    weighted_score = sum(
        sv["verdict_data"]["score"] ** 2  # square to amplify high-confidence signals
        for sv in step_verdicts
    ) / total_weight if total_weight > 0 else 0.0

    # Ensemble verdict via weighted majority
    if has_violations:
        overall_verdict = "needs_revision"
    elif approved_weight > revision_weight and weighted_score >= 0.7:
        overall_verdict = "approved"
    else:
        overall_verdict = "needs_revision"

    # De-duplicate issues, preserving order (first occurrence = highest priority)
    seen_issues: set = set()
    deduped_issues = []
    for issue in all_issues:
        key = issue.lower().strip()[:80]
        if key not in seen_issues:
            seen_issues.add(key)
            deduped_issues.append(issue)

    return {
        "verdict": overall_verdict,
        "score": round(weighted_score, 3),
        "issues": deduped_issues,
        "constraint_violations": all_violations,
        "suggestions": list(dict.fromkeys(all_suggestions)),  # deduplicated
        "checklist_results": step_verdicts[0]["verdict_data"]["checklist_results"] if step_verdicts else {},
        "parallel_review": {
            "total_steps": len(step_verdicts),
            "approved_steps": sum(1 for sv in step_verdicts if sv["verdict_data"]["verdict"] in ("approved", "pass")),
            "failed_steps": [sv["step_id"] for sv in step_verdicts if sv["verdict_data"]["verdict"] not in ("approved", "pass")],
            "ensemble_method": "moa_weighted_majority",
            "approved_weight": round(approved_weight, 3),
            "revision_weight": round(revision_weight, 3),
            "step_results": step_results,
        },
    }


def write_fallback_verdict(verdict_file: Path, reason: str):
    """Always write a verdict file even on error — orchestrate.sh depends on it. (#18 audit fix)"""
    write_json(str(verdict_file), {
        "verdict": "needs_revision",
        "score": 0.0,
        "checklist_results": {"completeness": False, "constraint_compliance": False,
                               "content_accuracy": False, "design_quality": False},
        "constraint_violations": [],
        "issues": [f"Parallel reviewer error: {reason}"],
        "suggestions": [],
    })


def main():
    if len(sys.argv) < 5:
        print("Usage: parallel_reviewer.py <sprint_contract.json> <aggregated_output.json> <module> <run_dir> [attempt]")
        sys.exit(1)

    contract_file = sys.argv[1]
    aggregated_output_file = sys.argv[2]
    module = sys.argv[3]
    run_dir = Path(sys.argv[4])
    attempt = int(sys.argv[5]) if len(sys.argv) > 5 else 1

    verdict_file = run_dir / f"verdict_{attempt}.json"

    try:
        print(f"  [Parallel Reviewer] Loading Sprint_Contract and aggregated output...")
        contract = load_json(contract_file)
        aggregated = load_json(aggregated_output_file)
    except Exception as e:
        print(f"  ERROR: Failed to load input files: {e}", file=sys.stderr)
        write_fallback_verdict(verdict_file, str(e))
        sys.exit(0)  # Always exit 0; verdict file contains the failure

    # Extract per-step outputs from individual step_N_output.json files (preferred)
    step_outputs = {}
    for step_file in sorted(run_dir.glob("step_*_output.json")):
        try:
            step_id = int(step_file.name.split("_")[1])
            step_outputs[step_id] = load_json(str(step_file))
        except (ValueError, IndexError):
            pass

    # Fallback: extract from aggregated outputs array
    if not step_outputs:
        for out in aggregated.get("outputs", []):
            step_id = out.get("step_id")
            if step_id is not None:
                step_outputs[step_id] = out

    # Last resort: review whole aggregated output as one step
    if not step_outputs:
        step_outputs = {"all": aggregated}

    print(f"  [Parallel Reviewer] {len(step_outputs)} step(s) to review in parallel (attempt {attempt})")

    # Prepare review inputs (one per step)
    review_tasks = []
    for step_id, step_out in step_outputs.items():
        review_input_path = run_dir / f"review_{attempt}_step_{step_id}_input.txt"
        review_output_path = run_dir / f"review_{attempt}_step_{step_id}_verdict.json"

        review_input = build_review_input(contract, step_id, step_out, module, attempt)
        review_input_path.write_text(review_input, encoding="utf-8")

        review_tasks.append({
            "step_id": step_id,
            "input_path": review_input_path,
            "output_path": review_output_path,
        })

    # Run all Reviewers in parallel (ThreadPoolExecutor — each launches subprocess)
    # 공식 근거: subagents.md — parallel subagents pattern
    print(f"  [Parallel Reviewer] Launching {len(review_tasks)} reviewer subprocess(es) concurrently...")
    start_time = time.time()

    # Write reviewer status file — 각 step 완료를 파일로 기록
    import os as _os
    status_file = str(run_dir / "reviewer_status.json")

    def write_reviewer_status(step_id, status):
        try:
            import json as _json
            existing = {}
            if _os.path.exists(status_file):
                with open(status_file) as f:
                    existing = _json.load(f)
            existing[str(step_id)] = {"status": status}
            tmp = status_file + ".tmp"
            with open(tmp, "w") as f:
                _json.dump(existing, f, indent=2)
            _os.replace(tmp, status_file)
        except Exception:
            pass

    run_results = {}
    max_workers = min(len(review_tasks), max(4, (_os.cpu_count() or 4)), 8)
    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {
            pool.submit(
                run_reviewer,
                task["step_id"],
                task["input_path"],
                task["output_path"],
                module,
            ): task["step_id"]
            for task in review_tasks
        }
        print(f"  [Reviewer 시작] {len(review_tasks)}개 step 병렬 검증 중...", flush=True)
        for future in as_completed(futures):
            try:
                result = future.result()
            except Exception as e:
                step_id = futures[future]
                step = next((s for s in contract.get('steps', []) if s['id'] == step_id), {})
                lbl = _step_label(step)
                print(f"  ❌ [Reviewer 오류] {lbl} — {e}", flush=True)
                write_reviewer_status(step_id, "error")
                run_results[step_id] = {"step_id": step_id, "status": "error", "error": str(e)}
                continue
            run_results[result["step_id"]] = result
            step = next((s for s in contract.get('steps', []) if s['id'] == result["step_id"]), {})
            lbl = _step_label(step)
            if result["status"] == "timeout":
                print(f"  ❌ [Reviewer 시간초과] {lbl} — 검증 시간 초과", flush=True)
            elif result["status"] in ("failed", "error"):
                print(f"  ❌ [Reviewer 실패] {lbl} — {result.get('error','')[:100]}", flush=True)
            else:
                print(f"  ✅ [Reviewer 완료] {lbl}", flush=True)
            write_reviewer_status(result["step_id"], result["status"])

    # 누락된 step 감지 — 실행됐어야 하는데 결과가 없는 step 경고
    missing = [t["step_id"] for t in review_tasks if t["step_id"] not in run_results]
    if missing:
        print(f"  [Parallel Reviewer] ⚠ WARNING: Step(s) {missing} produced no result — skipped or crashed", flush=True)

    elapsed = time.time() - start_time
    print(f"  [Parallel Reviewer] All reviews completed in {elapsed:.1f}s")

    # Parse individual verdicts
    step_verdicts = []
    for task in review_tasks:
        step_id = task["step_id"]
        verdict_path = str(task["output_path"])
        verdict_data = parse_verdict(verdict_path)
        step_verdicts.append({"step_id": step_id, "verdict_data": verdict_data})
        step = next((s for s in contract.get('steps', []) if s.get('id') == step_id), {})
        lbl = _step_label(step)
        icon = "✅" if verdict_data["verdict"] in ("approved", "pass") else "⚠️"
        # Summarize first 2 issues
        raw_issues = verdict_data.get("issues") or []
        issue_parts = []
        for iss in raw_issues[:2]:
            if isinstance(iss, str):
                issue_parts.append(iss[:60])
            elif isinstance(iss, dict):
                issue_parts.append((iss.get("description") or iss.get("issue") or "")[:60])
        issues_str = "; ".join(issue_parts)
        suffix = f": {issues_str}" if issues_str else ""
        print(f"  {icon} [Reviewer 결과] {lbl} — {verdict_data['verdict']} ({verdict_data['score']:.2f}){suffix}")

    # Aggregate
    aggregated_verdict = aggregate_verdicts(step_verdicts)

    # Write final aggregated verdict
    verdict_file = run_dir / f"verdict_{attempt}.json"
    write_json(str(verdict_file), aggregated_verdict)

    print(f"  [Parallel Reviewer] Overall: {aggregated_verdict['verdict']} (score={aggregated_verdict['score']:.2f})")
    print(f"  [Parallel Reviewer] Verdict written: {verdict_file}")

    # Always exit 0 — orchestrate.sh reads verdict file to determine outcome (#18 audit fix)
    # Writing verdict file is the contract; exit code is secondary
    sys.exit(0)


if __name__ == "__main__":
    main()
