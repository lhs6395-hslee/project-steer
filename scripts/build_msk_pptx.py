#!/usr/bin/env python3
"""Build AWS MSK Expert Intro presentation from template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import os

# ── Constants ──
PRIMARY = RGBColor(0, 67, 218)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(33, 33, 33)
GRAY = RGBColor(80, 80, 80)
WHITE = RGBColor(255, 255, 255)
BG_BOX_RGB = RGBColor(248, 249, 250)
BORDER_RGB = RGBColor(220, 220, 220)

FONT_TITLE = "프리젠테이션 7 Bold"
FONT_DESC = "프리젠테이션 5 Medium"
FONT_BODY = "Freesentation"

TEMPLATE = "templates/pptx_template.pptx"
OUTPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"


def set_text(shape, text, font_name=None, font_size=None, color=None, bold=None, alignment=None):
    """Replace all text in a shape's text_frame (creates new run, loses original formatting)."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    fmt = run.font
    if font_name:
        fmt.name = font_name
    if font_size:
        fmt.size = Pt(font_size)
    if color:
        fmt.color.rgb = color
    if bold is not None:
        fmt.bold = bold


def set_multiline_text(shape, lines, font_name=None, font_size=None, color=None, bold=None, alignment=None, line_spacing=None):
    """Replace text with multiple lines, each as a separate paragraph (creates new runs)."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if alignment:
            p.alignment = alignment
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        fmt = run.font
        if font_name:
            fmt.name = font_name
        if font_size:
            fmt.size = Pt(font_size)
        if color:
            fmt.color.rgb = color
        if bold is not None:
            fmt.bold = bold


def replace_text_preserve_format(shape, new_text):
    """Replace text in shape while preserving original formatting (font, size, color, scheme).
    
    For single-paragraph text: replaces first run's text, removes extra runs.
    Preserves scheme colors, font names, sizes, bold settings from template.
    """
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p = tf.paragraphs[0]
    if p.runs:
        # Set all text into first run, preserving its formatting
        p.runs[0].text = new_text
        # Remove extra runs by clearing their text
        for run in p.runs[1:]:
            run.text = ""
    # Remove extra paragraphs
    from pptx.oxml.ns import qn
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def replace_multiline_preserve_format(shape, lines, font_size_override=None):
    """Replace text with multiple lines, preserving first run's formatting per paragraph.
    
    Copies formatting from the first paragraph's first run to new paragraphs.
    Preserves scheme colors. Optionally overrides font size.
    """
    import copy as _copy
    from lxml import etree
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        return
    
    # Get reference paragraph XML for cloning
    ref_p = tf.paragraphs[0]._p
    
    # Clear all existing paragraphs
    txBody = tf._txBody
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for p_elem in list(txBody.iterchildren(f'{ns}p')):
        txBody.remove(p_elem)
    
    # Create new paragraphs by cloning reference
    for i, line in enumerate(lines):
        new_p = _copy.deepcopy(ref_p)
        # Set text in first run only
        runs = new_p.findall(f'.//{ns}r')
        if runs:
            # Set first run text
            t_elem = runs[0].find(f'{ns}t')
            if t_elem is not None:
                t_elem.text = line
            # Override font size if requested
            if font_size_override is not None:
                rPr = runs[0].find(f'{ns}rPr')
                if rPr is not None:
                    rPr.set('sz', str(int(font_size_override * 100)))  # pt to hundredths
            # Clear extra runs
            for extra_run in runs[1:]:
                new_p.remove(extra_run)
        txBody.append(new_p)


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY, font_size=12, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    fmt = run.font
    fmt.name = font_name
    fmt.size = Pt(font_size)
    fmt.color.rgb = color
    fmt.bold = bold
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_name=FONT_BODY, font_size=12, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_before = Pt(line_spacing) if i > 0 else Pt(0)
        run = p.add_run()
        run.text = line
        fmt = run.font
        fmt.name = font_name
        fmt.size = Pt(font_size)
        fmt.color.rgb = color
        fmt.bold = bold
    return txBox


def add_rich_textbox(slide, left, top, width, height, paragraphs_data, alignment=PP_ALIGN.LEFT):
    """Add a text box with mixed formatting per paragraph.
    paragraphs_data: list of list of dicts with keys: text, font_name, font_size, color, bold
    """
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, runs_data in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        for rd in runs_data:
            run = p.add_run()
            run.text = rd.get("text", "")
            fmt = run.font
            fmt.name = rd.get("font_name", FONT_BODY)
            fmt.size = Pt(rd.get("font_size", 12))
            fmt.color.rgb = rd.get("color", BLACK)
            if rd.get("bold"):
                fmt.bold = True
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None,
                     text=None, font_name=FONT_BODY, font_size=13, text_color=DARK_GRAY, bold=False):
    """Add a rounded rectangle shape with optional centered text inside."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        # 수직 중앙 (Middle)
        from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR
        tf.paragraphs[0].space_before = Pt(0)
        from pptx.oxml.ns import qn
        txBody = tf._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.bold = bold
    return shape


def add_shape_with_text(slide, shape_type, left, top, width, height, fill_color=PRIMARY,
                        text="", font_name=FONT_TITLE, font_size=14, text_color=WHITE, bold=True):
    """Add any shape with centered text (horizontal + vertical Middle)."""
    shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.bold = bold
    return shape


def add_oval(slide, left, top, width, height, fill_color=None):
    """Add an oval shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, fill_color=PRIMARY):
    """Add a right arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_body_header(slide, title_text, desc_text, body_title=None, body_desc=None):
    """Add standard body slide header (title left + description right) + optional body title.
    
    Returns content_top_y (EMU) — 콘텐츠 시작 Y 좌표.
    body_title이 있으면 구분선 + 본문 제목 영역이 추가되어 콘텐츠가 아래로 밀림.
    """
    # 중제목 (좌): left=354806, top=557906
    add_textbox(slide,
        left=354806, top=557906, width=2720317, height=932614,
        text=title_text, font_name=FONT_TITLE, font_size=28,
        color=PRIMARY, bold=True)
    # 중제목 설명글 (우): left=3763700, top=558800
    add_textbox(slide,
        left=3763700, top=558800, width=7551483, height=937918,
        text=desc_text, font_name=FONT_DESC, font_size=12,
        color=GRAY, bold=False)
    
    if body_title:
        # 본문 제목 (좌): Freesentation 16pt, PRIMARY
        bt_top = inches(1.75)
        add_textbox(slide,
            left=SLIDE_MARGIN, top=bt_top, width=12192000 - 2 * SLIDE_MARGIN, height=inches(0.35),
            text=body_title, font_name=FONT_BODY, font_size=16, color=PRIMARY, bold=True)
        
        # 본문 설명글 (아래): Freesentation 13pt, DARK_GRAY
        if body_desc:
            add_textbox(slide,
                left=SLIDE_MARGIN, top=bt_top + inches(0.4),
                width=12192000 - 2 * SLIDE_MARGIN, height=inches(0.3),
                text=body_desc, font_name=FONT_BODY, font_size=13, color=DARK_GRAY)
            return bt_top + inches(0.85)  # 콘텐츠 시작: 설명글 아래 + 여유
        return bt_top + inches(0.55)
    
    return inches(2.0)  # 본문 제목 없으면 기존대로


# ── Inch/EMU helpers ──
def inches(val):
    return int(val * 914400)

SLIDE_W = 12192000  # 13.333"
SLIDE_MARGIN = inches(0.5)  # 좌우 최소 여백


def estimate_text_width_pt(text, font_size_pt):
    """텍스트의 대략적인 렌더링 폭을 pt 단위로 추정.
    한글/CJK는 font_size와 동일, 영문/숫자는 font_size * 0.55, 공백은 font_size * 0.3.
    """
    width = 0
    for ch in text:
        if ord(ch) > 0x2E80:  # CJK/한글
            width += font_size_pt
        elif ch == ' ':
            width += font_size_pt * 0.3
        elif ch in '.:/-()':
            width += font_size_pt * 0.4
        else:  # 영문/숫자
            width += font_size_pt * 0.55
    return width


def auto_fit_textbox_width(shape, lines, font_size_pt, margin_left=None, center=False):
    """텍스트 줄 중 가장 긴 줄에 맞춰 텍스트박스 너비를 자동 조정.
    
    Args:
        shape: 텍스트박스 shape
        lines: 텍스트 줄 리스트
        font_size_pt: 폰트 크기 (pt)
        margin_left: 고정 left 위치 (None이면 현재 유지)
        center: True면 슬라이드 중앙 배치 (left 자동 계산)
    """
    # 가장 긴 줄의 폭 계산
    max_width_pt = max(estimate_text_width_pt(line, font_size_pt) for line in lines if line)
    # pt → EMU (1pt = 12700 EMU) + 패딩 10%
    needed_width = int(max_width_pt * 12700 * 1.1)
    
    # 슬라이드 경계 내로 제한
    if center:
        max_allowed = SLIDE_W - 2 * SLIDE_MARGIN
        needed_width = min(needed_width, max_allowed)
        shape.width = needed_width
        shape.left = (SLIDE_W - needed_width) // 2
    else:
        current_left = margin_left if margin_left is not None else shape.left
        max_allowed = SLIDE_W - current_left - SLIDE_MARGIN
        needed_width = min(needed_width, max_allowed)
        shape.width = needed_width
        if margin_left is not None:
            shape.left = margin_left


def main():
    prs = Presentation(TEMPLATE)

    # ════════════════════════════════════════
    # Step 2: Cover Slide (index 0)
    # 제목 텍스트박스 너비 확장 (50pt 유지, 2줄로 수용)
    # 부제 너비도 동일하게 확장
    # ════════════════════════════════════════
    cover = prs.slides[0]

    # shape[2] = Title: 자동 너비 + 중앙 배치, 3줄 수용, 수직 중앙
    title_shape = cover.shapes[2]
    title_lines = ["Amazon Managed Streaming", "for Apache Kafka", "(MSK)"]
    title_shape.height = inches(2.8)
    title_shape.top = inches(1.63)
    auto_fit_textbox_width(title_shape, title_lines, font_size_pt=50, center=True)
    replace_multiline_preserve_format(title_shape, title_lines)

    # shape[5] = Subtitle: 자동 너비 + 중앙 배치, 제목 아래 0.15" 간격
    subtitle_shape = cover.shapes[5]
    subtitle_text = "Expert Introduction — 아키텍처 · 핵심 기능 · 운영 전략"
    subtitle_shape.top = title_shape.top + title_shape.height + inches(0.15)
    auto_fit_textbox_width(subtitle_shape, [subtitle_text], font_size_pt=24, center=True)
    replace_text_preserve_format(subtitle_shape, subtitle_text)

    # shape[7] = Year
    replace_text_preserve_format(cover.shapes[7], "2026")
    # shape[8] = Date — 오늘 기준 04/13
    replace_text_preserve_format(cover.shapes[8], "04/13")

    # ════════════════════════════════════════
    # Step 3: TOC Slide (index 1~)
    # 서식 보존: 번호 열/섹션명 열 모두 기존 run XML 복제 방식
    # 페이징: 5개씩, 6개 이상이면 새 목차 페이지 추가
    # ════════════════════════════════════════
    toc_sections_all = [
        "MSK 아키텍처 및 핵심 개념",
        "MSK 클러스터 구성과 핵심 기능",
        "MSK 운영 전략 및 Best Practices",
    ]

    # 5개씩 페이징
    pages = []
    for i in range(0, len(toc_sections_all), 5):
        pages.append(toc_sections_all[i:i+5])

    # 첫 번째 목차 페이지: 기존 슬라이드(index 1) 수정
    toc_first = prs.slides[1]
    page1 = pages[0]
    nums1 = [str(i+1) for i in range(len(page1))]
    while len(nums1) < 5:
        nums1.append("")
    while len(page1) < 5:
        page1.append("")
    replace_multiline_preserve_format(toc_first.shapes[1], nums1)
    auto_fit_textbox_width(toc_first.shapes[2], page1, font_size_pt=24)
    replace_multiline_preserve_format(toc_first.shapes[2], page1)

    # 추가 목차 페이지: 첫 번째 목차 슬라이드의 shape들을 통째로 복제
    added_toc_count = 0
    for page_idx in range(1, len(pages)):
        page = pages[page_idx]
        start_num = page_idx * 5 + 1
        nums = [str(start_num + i) for i in range(len(page))]
        while len(nums) < 5:
            nums.append("")
        while len(page) < 5:
            page.append("")

        import copy as _copy_mod
        ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        # 빈 목차 슬라이드 추가
        new_slide = prs.slides.add_slide(prs.slide_layouts[4])
        added_toc_count += 1

        # 첫 번째 목차의 모든 shape XML을 새 슬라이드의 spTree에 복제
        src_spTree = toc_first._element.find(f'{ns_p}cSld').find(f'{ns_p}spTree')
        dst_spTree = new_slide._element.find(f'{ns_p}cSld').find(f'{ns_p}spTree')

        # 기존 dst shape 제거 (placeholder 포함)
        for child in list(dst_spTree):
            if child.tag != f'{ns_a}grpSpPr' and child.tag != f'{ns_p}nvGrpSpPr':
                # grpSpPr는 spTree 자체의 속성이므로 유지
                if 'grpSpPr' not in child.tag:
                    dst_spTree.remove(child)

        # 원본 shape들 복제하여 추가
        for child in src_spTree:
            if 'grpSpPr' not in child.tag:
                dst_spTree.append(_copy_mod.deepcopy(child))

        # 이제 new_slide.shapes로 접근 가능
        # shape[0]=CONTENTS, [1]=번호열, [2]=섹션명열, [3]=장식그룹, [4]=placeholder
        replace_multiline_preserve_format(new_slide.shapes[1], nums)
        auto_fit_textbox_width(new_slide.shapes[2], page, font_size_pt=24)
        replace_multiline_preserve_format(new_slide.shapes[2], page)

    # ════════════════════════════════════════
    # Step 7: Ending Slide (index 40) — 템플릿 그대로 사용 (수정 없음)
    # "Thank You" + 원본 태그라인 유지
    # ════════════════════════════════════════

    # ════════════════════════════════════════
    # Step 4: Add 3 body slides (layout_index=1)
    # ════════════════════════════════════════
    body_layout = prs.slide_layouts[1]

    # ── Body Slide 1: MSK Architecture Overview ──
    slide1 = prs.slides.add_slide(body_layout)
    content_y1 = add_body_header(slide1, "1-1. MSK 아키텍처 개요",
        "Apache Kafka 완전관리형 서비스의 구성 요소와 데이터 흐름",
        body_title="Producer → MSK Cluster → Consumer 데이터 흐름",
        body_desc="메시지 생산부터 소비까지의 전체 아키텍처 구성도")

    # Architecture diagram: Producer → MSK Cluster → Consumer
    # 중앙 정렬: 전체 너비를 먼저 계산하고 시작 left를 (SLIDE_W - total) / 2로 배치
    body_top = content_y1 + inches(0.2)
    box_w = inches(2.0)
    box_h = inches(1.3)
    arrow_w = inches(0.5)
    arrow_h = inches(0.4)
    msk_w = inches(5.6)
    gap = inches(0.2)  # 요소 간 간격

    total_diagram_w = box_w + gap + arrow_w + gap + msk_w + gap + arrow_w + gap + box_w
    # = 2.0 + 0.2 + 0.5 + 0.2 + 5.6 + 0.2 + 0.5 + 0.2 + 2.0 = 11.4"
    diagram_left = (SLIDE_W - total_diagram_w) // 2  # 중앙 정렬

    # Producer box
    prod_left = diagram_left
    add_rounded_rect(slide1, prod_left, body_top, box_w, box_h, fill_color=PRIMARY,
                     text="Producer", font_name=FONT_TITLE, font_size=14, text_color=WHITE, bold=True)

    # Arrow 1
    arr1_left = prod_left + box_w + gap
    arr1_top = body_top + (box_h - arrow_h) // 2
    add_arrow(slide1, arr1_left, arr1_top, arrow_w, arrow_h, fill_color=PRIMARY)

    # MSK Cluster box
    msk_left = arr1_left + arrow_w + gap
    add_rounded_rect(slide1, msk_left, body_top, msk_w, box_h,
        fill_color=RGBColor(230, 240, 255), line_color=PRIMARY, line_width=1.5)
    add_textbox(slide1, msk_left, body_top + inches(0.1), msk_w, inches(0.35),
        "MSK Cluster (Multi-AZ Brokers)", font_name=FONT_TITLE, font_size=14,
        color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    # Broker nodes inside — 3개 균등 배치
    broker_count = 3
    broker_margin = inches(0.2)
    broker_gap = inches(0.15)
    broker_total_w = msk_w - 2 * broker_margin
    broker_w = (broker_total_w - (broker_count - 1) * broker_gap) // broker_count
    broker_h = inches(0.45)
    broker_top = body_top + inches(0.65)
    for i, az in enumerate(["AZ-a", "AZ-b", "AZ-c"]):
        b_left = msk_left + broker_margin + i * (broker_w + broker_gap)
        add_rounded_rect(slide1, b_left, broker_top, broker_w, broker_h, fill_color=PRIMARY,
                         text=f"Broker ({az})", font_name=FONT_TITLE, font_size=12, text_color=WHITE, bold=True)

    # Arrow 2
    arr2_left = msk_left + msk_w + gap
    add_arrow(slide1, arr2_left, arr1_top, arrow_w, arrow_h, fill_color=PRIMARY)

    # Consumer box
    cons_left = arr2_left + arrow_w + gap
    add_rounded_rect(slide1, cons_left, body_top, box_w, box_h, fill_color=PRIMARY,
                     text="Consumer", font_name=FONT_TITLE, font_size=14, text_color=WHITE, bold=True)

    # Bottom 3-column explanation — 좌우 여백 동일 (0.5")
    col_count = 3
    col_gap_val = inches(0.25)
    usable_w = SLIDE_W - 2 * SLIDE_MARGIN  # 12.333"
    col_w = (usable_w - (col_count - 1) * col_gap_val) // col_count
    col_h = inches(2.5)
    col_top = inches(4.4)
    cols_data = [
        ("핵심 구성요소", [
            "• Broker: 메시지 저장/전달 노드",
            "• Topic: 메시지 카테고리 단위",
            "• Partition: 병렬 처리 단위",
            "• Replication Factor: 데이터 복제 수",
            "• ZooKeeper / KRaft: 메타데이터 관리",
        ]),
        ("네트워크 구성", [
            "• VPC 내 Private Subnet 배치",
            "• Security Group 기반 접근 제어",
            "• ENI: 브로커별 네트워크 인터페이스",
            "• PrivateLink: Cross-VPC 연결",
            "• Multi-AZ 자동 분산 배치",
        ]),
        ("스토리지", [
            "• EBS gp3: 기본 브로커 스토리지",
            "• Provisioned IOPS: 고성능 워크로드",
            "• Tiered Storage: S3 기반 장기 보관",
            "• 자동 확장: 디스크 용량 자동 증설",
            "• 보존 기간별 계층화 전략",
        ]),
    ]
    for i, (title, bullets) in enumerate(cols_data):
        c_left = SLIDE_MARGIN + i * (col_w + col_gap_val)
        add_rounded_rect(slide1, c_left, col_top, col_w, col_h,
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        add_textbox(slide1, c_left + inches(0.2), col_top + inches(0.15),
            col_w - inches(0.4), inches(0.35),
            title, font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
        add_multiline_textbox(slide1,
            c_left + inches(0.2), col_top + inches(0.55),
            col_w - inches(0.4), col_h - inches(0.7),
            bullets, font_name=FONT_BODY, font_size=13, color=DARK_GRAY,
            line_spacing=2)

    # ── Body Slide 2: Cluster Config & Core Features (2x2 grid) ──
    slide2 = prs.slides.add_slide(body_layout)
    content_y2 = add_body_header(slide2, "2-1. 클러스터 구성과 핵심 기능",
        "Provisioned vs Serverless, 보안, 모니터링, Connect 통합",
        body_title="MSK 핵심 기능 4가지",
        body_desc="클러스터 유형, 보안, 모니터링, Connect 통합 비교")

    card_data = [
        ("클러스터 유형", [
            "• Provisioned: 브로커 수/인스턴스 타입 직접 지정",
            "• Serverless: 자동 스케일링, 용량 관리 불필요",
            "• Express: 고성능 + 간소화된 운영",
        ]),
        ("보안 체계", [
            "• 인증: IAM / SASL-SCRAM / mTLS",
            "• 암호화: TLS in-transit, KMS at-rest",
            "• 네트워크: VPC 내 배치, Private Link",
        ]),
        ("모니터링", [
            "• CloudWatch 메트릭 (Basic/Enhanced/Topic)",
            "• Open Monitoring (Prometheus 호환)",
            "• Broker 로그 → CloudWatch/S3/Firehose",
        ]),
        ("MSK Connect", [
            "• Kafka Connect 완전관리형 서비스",
            "• Source/Sink Connector 플러그인",
            "• Auto Scaling, 커스텀 플러그인 지원",
        ]),
    ]
    card_w = inches(5.9)
    card_h = inches(2.0)
    card_gap_x = inches(0.3)
    card_gap_y = inches(0.3)
    # 중앙 정렬: 전체 그리드 너비 계산 후 시작 left 결정
    grid_total_w = 2 * card_w + card_gap_x  # 5.9*2 + 0.3 = 12.1"
    grid_left = (SLIDE_W - grid_total_w) // 2
    grid_top = content_y2 + inches(0.2)

    card_icons = ["icons/cluster.png", "icons/security.png", "icons/monitoring.png", "icons/connect.png"]
    for idx, (title, bullets) in enumerate(card_data):
        row = idx // 2
        col = idx % 2
        c_left = grid_left + col * (card_w + card_gap_x)
        c_top = grid_top + row * (card_h + card_gap_y)

        add_rounded_rect(slide2, c_left, c_top, card_w, card_h,
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        # 아이콘 이미지 (우상단)
        icon_size = inches(0.45)
        icon_path = card_icons[idx]
        import os
        if os.path.exists(icon_path):
            slide2.shapes.add_picture(icon_path,
                Emu(c_left + card_w - icon_size - inches(0.2)),
                Emu(c_top + inches(0.15)),
                Emu(icon_size), Emu(icon_size))
        # Card title — 14pt
        add_textbox(slide2, c_left + inches(0.25), c_top + inches(0.15),
            card_w - inches(1.0), inches(0.4),
            title, font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
        # Card body — 13pt
        add_multiline_textbox(slide2,
            c_left + inches(0.25), c_top + inches(0.6),
            card_w - inches(0.5), card_h - inches(0.8),
            bullets, font_name=FONT_BODY, font_size=13, color=DARK_GRAY,
            line_spacing=3)

    # ── Body Slide 3: Operations & Best Practices (numbered list) ──
    slide3 = prs.slides.add_slide(body_layout)
    content_y3 = add_body_header(slide3, "3-1. 운영 전략 및 Best Practices",
        "클러스터 사이징, 파티션 설계, 장애 대응, 비용 최적화",
        body_title="운영 핵심 체크리스트",
        body_desc="프로덕션 환경에서 반드시 확인해야 할 4가지 항목")

    items = [
        ("1", "클러스터 사이징",
         "브로커 수 = (피크 처리량 ÷ 브로커당 처리량) × RF\nkafka.m5.large 이상 권장, EBS gp3 기본"),
        ("2", "파티션 설계",
         "파티션 수 = max(Producer 처리량 ÷ 파티션당 쓰기, Consumer 처리량 ÷ 파티션당 읽기)\n브로커당 4,000개 이하 권장"),
        ("3", "장애 대응",
         "Multi-AZ 배포 필수, min.insync.replicas=2\nunclean.leader.election=false, 브로커 장애 시 자동 복구"),
        ("4", "비용 최적화",
         "Tiered Storage로 장기 보관 비용 절감\nServerless로 변동 워크로드 대응, 불필요 토픽/파티션 정리"),
    ]

    # Body area: 2.0"–7.2" (body_limit_y=6583680)
    # 4 items: 2.0 + 4*1.0 + 3*0.12 = 2.0 + 4.0 + 0.36 = 6.36" ✓ (< 7.2", margin 0.84")
    item_top_start = content_y3 + inches(0.15)
    item_h = inches(0.95)
    item_gap = inches(0.1)
    badge_size = inches(0.5)
    # 중앙 정렬: badge + gap + content_box 전체를 중앙 배치
    content_w = inches(10.5)
    total_row_w = badge_size + inches(0.2) + content_w  # 0.55 + 0.2 + 10.5 = 11.25"
    row_left = (SLIDE_W - total_row_w) // 2
    badge_left = row_left
    content_left = row_left + badge_size + inches(0.2)

    for i, (num, title, body) in enumerate(items):
        y = item_top_start + i * (item_h + item_gap)

        # Content card — 도형 내부에 제목+본문 텍스트 직접 삽입, 내부 여백 균등
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(content_left), Emu(y), Emu(content_w), Emu(item_h))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_BOX_RGB
        card.line.color.rgb = BORDER_RGB
        card.line.width = Pt(0.5)

        # 내부 여백 설정 (상하좌우 동일)
        from pptx.oxml.ns import qn
        bodyPr = card.text_frame._txBody.find(qn('a:bodyPr'))
        pad = int(0.12 * 914400)  # 0.12" 내부 여백
        if bodyPr is not None:
            bodyPr.set('lIns', str(pad))
            bodyPr.set('rIns', str(pad))
            bodyPr.set('tIns', str(pad))
            bodyPr.set('bIns', str(pad))

        tf = card.text_frame
        tf.word_wrap = True
        # 제목 paragraph — 좌측 정렬
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        run_t = p_title.add_run()
        run_t.text = title
        run_t.font.name = FONT_TITLE
        run_t.font.size = Pt(14)
        run_t.font.color.rgb = PRIMARY
        run_t.font.bold = True
        # 본문 paragraphs
        for line in body.split("\n"):
            p_body = tf.add_paragraph()
            p_body.alignment = PP_ALIGN.LEFT
            p_body.space_before = Pt(2)
            run_b = p_body.add_run()
            run_b.text = line
            run_b.font.name = FONT_BODY
            run_b.font.size = Pt(13)
            run_b.font.color.rgb = DARK_GRAY

        # Number badge (circle) — 수직/수평 중앙
        badge_y = y + (item_h - badge_size) // 2
        add_shape_with_text(slide3, MSO_SHAPE.OVAL, badge_left, badge_y, badge_size, badge_size,
            fill_color=PRIMARY, text=num, font_name=FONT_TITLE, font_size=14, text_color=WHITE, bold=True)

    # ════════════════════════════════════════
    # Step 5: Layout Test Slides (5종)
    # L01 Bento Grid, L02 Three Cards, L03 Grid 2x2,
    # L04 Process Arrow, L05 Phased Columns
    # ════════════════════════════════════════

    # ── L01 Bento Grid ──
    sl = prs.slides.add_slide(body_layout)
    add_body_header(sl, "L01. Bento Grid", "좌 50% + 우 2분할 레이아웃 테스트")
    add_rounded_rect(sl, SLIDE_MARGIN, inches(2), inches(5.9), inches(4.8),
        fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
    add_textbox(sl, SLIDE_MARGIN+inches(0.15), inches(2.12), inches(5.6), inches(0.35),
        "MSK 아키텍처 개요", font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
    add_multiline_textbox(sl, SLIDE_MARGIN+inches(0.15), inches(2.55), inches(5.6), inches(3.5),
        ["• Apache Kafka 완전관리형 서비스", "• Multi-AZ 브로커 자동 분산 배치",
         "• VPC 내 Private Subnet 배치", "• ZooKeeper/KRaft 메타데이터 관리",
         "• EBS gp3 기본 스토리지"], font_name=FONT_BODY, font_size=13, color=DARK_GRAY)
    rl = SLIDE_MARGIN + inches(6.15)
    rw = inches(6.18)
    add_rounded_rect(sl, rl, inches(2), rw, inches(2.25),
        fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
    add_textbox(sl, rl+inches(0.15), inches(2.12), rw-inches(0.3), inches(0.35),
        "클러스터 유형", font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
    add_multiline_textbox(sl, rl+inches(0.15), inches(2.55), rw-inches(0.3), inches(1.5),
        ["• Provisioned: 직접 지정", "• Serverless: 자동 스케일링", "• Express: 고성능"],
        font_name=FONT_BODY, font_size=13, color=DARK_GRAY)
    add_rounded_rect(sl, rl, inches(4.55), rw, inches(2.25),
        fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
    add_textbox(sl, rl+inches(0.15), inches(4.67), rw-inches(0.3), inches(0.35),
        "보안 체계", font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
    add_multiline_textbox(sl, rl+inches(0.15), inches(5.1), rw-inches(0.3), inches(1.5),
        ["• IAM / SASL-SCRAM / mTLS", "• TLS + KMS 암호화", "• Security Group 접근 제어"],
        font_name=FONT_BODY, font_size=13, color=DARK_GRAY)

    # ── L02 Three Cards ──
    sl = prs.slides.add_slide(body_layout)
    add_body_header(sl, "L02. Three Cards", "3열 카드 레이아웃 테스트")
    usable = 12192000 - 2 * SLIDE_MARGIN
    cgap = inches(0.25)
    cw = (usable - 2 * cgap) // 3
    l02_cards = [
        ("Provisioned", ["• 브로커 수 직접 지정", "• 세밀한 성능 튜닝", "• 예측 가능한 비용", "• 대규모 프로덕션"]),
        ("Serverless", ["• 자동 스케일링", "• 사용한 만큼 과금", "• 운영 부담 최소", "• 변동 워크로드"]),
        ("Express", ["• 고성능 + 간소화", "• 1-Click 생성", "• 자동 리밸런싱", "• 최신 기능 우선"]),
    ]
    for ci, (ct, cb) in enumerate(l02_cards):
        cl = SLIDE_MARGIN + ci * (cw + cgap)
        add_rounded_rect(sl, cl, inches(2), cw, inches(4.8),
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        add_textbox(sl, cl+inches(0.15), inches(2.12), cw-inches(0.3), inches(0.35),
            ct, font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
        add_multiline_textbox(sl, cl+inches(0.15), inches(2.55), cw-inches(0.3), inches(3.5),
            cb, font_name=FONT_BODY, font_size=13, color=DARK_GRAY)

    # ── L03 Grid 2x2 ──
    sl = prs.slides.add_slide(body_layout)
    add_body_header(sl, "L03. Grid 2x2", "2x2 카드 그리드 레이아웃 테스트")
    gw = (usable - cgap) // 2
    gh = inches(2.3)
    l03_grid = [
        ("브로커 관리", ["• Multi-AZ 자동 분산", "• 장애 시 자동 복구", "• 인스턴스 타입 변경"]),
        ("스토리지", ["• EBS gp3 기본", "• Tiered Storage (S3)", "• 자동 디스크 확장"]),
        ("모니터링", ["• CloudWatch 3단계", "• Prometheus 호환", "• Broker 로그 전송"]),
        ("MSK Connect", ["• Kafka Connect 관리형", "• Source/Sink 플러그인", "• Auto Scaling"]),
    ]
    for gi, (gt, gb) in enumerate(l03_grid):
        gr, gc = gi // 2, gi % 2
        gl = SLIDE_MARGIN + gc * (gw + cgap)
        gtp = inches(2) + gr * (gh + cgap)
        add_rounded_rect(sl, gl, gtp, gw, gh,
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        add_textbox(sl, gl+inches(0.15), gtp+inches(0.12), gw-inches(0.3), inches(0.35),
            gt, font_name=FONT_TITLE, font_size=14, color=PRIMARY, bold=True)
        add_multiline_textbox(sl, gl+inches(0.15), gtp+inches(0.5), gw-inches(0.3), gh-inches(0.6),
            gb, font_name=FONT_BODY, font_size=13, color=DARK_GRAY)

    # ── L04 Process Arrow ──
    sl = prs.slides.add_slide(body_layout)
    add_body_header(sl, "L04. Process Arrow", "4단계 프로세스 화살표 레이아웃 테스트")
    pw = (usable - 3 * cgap) // 4
    l04_steps = [
        ("1. 설계", ["• VPC/Subnet 설계", "• 브로커 수 결정", "• 보안 그룹 설정"]),
        ("2. 생성", ["• 클러스터 생성", "• 인증 방식 설정", "• 암호화 활성화"]),
        ("3. 연결", ["• Producer 연결", "• MSK Connect", "• PrivateLink"]),
        ("4. 운영", ["• 모니터링 구성", "• 파티션 리밸런싱", "• 비용 최적화"]),
    ]
    for pi, (pt, pb) in enumerate(l04_steps):
        pl = SLIDE_MARGIN + pi * (pw + cgap)
        chev = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Emu(pl), Emu(inches(2)), Emu(pw), Emu(inches(0.6)))
        chev.fill.solid(); chev.fill.fore_color.rgb = PRIMARY; chev.line.fill.background()
        add_textbox(sl, pl, inches(2), pw, inches(0.6),
            pt, font_name=FONT_TITLE, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_rounded_rect(sl, pl, inches(2.9), pw, inches(3.9),
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        add_multiline_textbox(sl, pl+inches(0.15), inches(3.05), pw-inches(0.3), inches(3.5),
            pb, font_name=FONT_BODY, font_size=13, color=DARK_GRAY)

    # ── L05 Phased Columns ──
    sl = prs.slides.add_slide(body_layout)
    add_body_header(sl, "L05. Phased Columns", "단계별 그라데이션 컬럼 레이아웃 테스트")
    grad_colors = [RGBColor(0,27,94), RGBColor(0,67,218), RGBColor(59,122,237), RGBColor(160,195,250)]
    grad_txt = [WHITE, WHITE, WHITE, RGBColor(0,27,94)]
    l05_phases = [
        ("Phase 1: 평가", ["• 현재 환경 분석", "• 워크로드 파악", "• 비용/성능 비교"]),
        ("Phase 2: 설계", ["• 클러스터 유형 선택", "• 네트워크 설계", "• 보안 아키텍처"]),
        ("Phase 3: 구축", ["• 클러스터 프로비저닝", "• MirrorMaker 설정", "• 데이터 동기화"]),
        ("Phase 4: 운영", ["• 모니터링 대시보드", "• 알림 정책 설정", "• 성능 튜닝"]),
    ]
    for pi, (pt, pb) in enumerate(l05_phases):
        pl = SLIDE_MARGIN + pi * (pw + cgap)
        ph = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(pl), Emu(inches(2)), Emu(pw), Emu(inches(0.5)))
        ph.fill.solid(); ph.fill.fore_color.rgb = grad_colors[pi]; ph.line.fill.background()
        add_textbox(sl, pl, inches(2), pw, inches(0.5),
            pt, font_name=FONT_TITLE, font_size=13, color=grad_txt[pi], bold=True, alignment=PP_ALIGN.CENTER)
        add_rounded_rect(sl, pl, inches(2.7), pw, inches(4.1),
            fill_color=BG_BOX_RGB, line_color=BORDER_RGB, line_width=0.5)
        add_multiline_textbox(sl, pl+inches(0.15), inches(2.85), pw-inches(0.3), inches(3.5),
            pb, font_name=FONT_BODY, font_size=13, color=DARK_GRAY)

    # ════════════════════════════════════════
    # Step 6: Delete unnecessary slides (index 2~39, reverse order)
    # 템플릿 41장 + 추가 본문 3장 + 추가 목차 N장
    # 보존: 0=cover, 1=toc, 40=ending, 41+=added slides
    # 삭제: 2~39 (템플릿 가이드 슬라이드 38장)
    # ════════════════════════════════════════
    xml_slides = prs.slides._sldIdLst

    # Delete template slides 2~39 in reverse order
    for idx in range(39, 1, -1):
        slide_id_elem = xml_slides[idx]
        xml_slides.remove(slide_id_elem)

    # Now: 0=cover, 1=toc, 2=ending(was 40), 3+=added slides (toc pages + body slides)
    # Find ending slide and move to last
    for i, slide in enumerate(prs.slides):
        if slide.slide_layout.name == "끝맺음":
            elem = xml_slides[i]
            xml_slides.remove(elem)
            xml_slides.append(elem)
            break

    # ════════════════════════════════════════
    # Step 8: 섹션(Section) 재구성 — 산출물 구조에 맞게
    # ════════════════════════════════════════
    from lxml import etree as _etree
    import uuid as _uuid
    _prs_xml = prs.part._element
    _ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    _ns_p14 = '{http://schemas.microsoft.com/office/powerpoint/2010/main}'

    # 최종 슬라이드 ID 수집
    _sldIdLst = _prs_xml.find(f'{_ns_p}sldIdLst')
    _slide_ids = [elem.get('id') for elem in _sldIdLst]
    # _slide_ids: [표지, 목차, 본문1, 본문2, 본문3, 끝맺음]

    # 기존 sectionLst 찾아서 교체
    for _extLst in _prs_xml.findall(f'.//{_ns_p}extLst'):
        for _ext in list(_extLst):
            _old = _ext.find(f'{_ns_p14}sectionLst')
            if _old is not None:
                _ext.remove(_old)
                _new_sectionLst = _etree.SubElement(_ext, f'{_ns_p14}sectionLst')

                # 슬라이드 구조: 표지(1) + 목차(1+N) + 본문(M) + 끝맺음(1)
                toc_page_count = 1 + added_toc_count  # 첫 번째 + 추가분
                body_start = 1 + toc_page_count
                body_end = len(_slide_ids) - 1  # 끝맺음 제외

                sections_def = [
                    ("표지", [_slide_ids[0]]),
                    ("목차", _slide_ids[1:1+toc_page_count]),
                    ("본문", _slide_ids[body_start:body_end]),
                    ("끝맺음", [_slide_ids[-1]]),
                ]
                for sec_name, sec_ids in sections_def:
                    sec_elem = _etree.SubElement(_new_sectionLst, f'{_ns_p14}section')
                    sec_elem.set('name', sec_name)
                    sec_elem.set('id', '{' + str(_uuid.uuid4()).upper() + '}')
                    sldIdLst_elem = _etree.SubElement(sec_elem, f'{_ns_p14}sldIdLst')
                    for sid in sec_ids:
                        sldId_elem = _etree.SubElement(sldIdLst_elem, f'{_ns_p14}sldId')
                        sldId_elem.set('id', sid)
                break

    # ════════════════════════════════════════
    # Step 9: Save
    # ════════════════════════════════════════
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved to {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")

    # Verify
    verify = Presentation(OUTPUT)
    print(f"Verification - Slide count: {len(verify.slides)}")
    for i, slide in enumerate(verify.slides):
        layout_name = slide.slide_layout.name
        shapes_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text[:60].replace('\n', ' | ')
                if txt.strip():
                    shapes_text.append(txt)
        print(f"  Slide {i}: [{layout_name}] {'; '.join(shapes_text[:3])}")


if __name__ == "__main__":
    main()
