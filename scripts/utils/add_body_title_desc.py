"""Add body title/desc textboxes to L01(idx5), L04(idx8), L05(idx9)
and shift/shrink content shapes to fit within 7.0\" limit.

Layout-spec EMU constants:
  SUMMARY_TITLE: left=457200, top=1600200, w=11277600, h=320040 (16pt Freesentation PRIMARY Bold)
  SUMMARY_DESC:  left=457200, top=1965960, w=11277600, h=274320 (13pt Freesentation DARK_GRAY)
  CONTENT_NEW_START: top = 2560320 EMU (2.80") — same as Grid 2x2 card start
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INPUT_PATH = "results/pptx/AWS_MSK_Expert_Intro.pptx"
OUTPUT_PATH = "results/pptx/AWS_MSK_Expert_Intro.pptx"

PRIMARY   = RGBColor(0, 67, 218)
DARK_GRAY = RGBColor(33, 33, 33)

# layout-spec.md 스펙
SUMMARY_TITLE_LEFT = 457200
SUMMARY_TITLE_TOP  = 1600200
SUMMARY_TITLE_W    = 11277600
SUMMARY_TITLE_H    = 320040

SUMMARY_DESC_LEFT  = 457200
SUMMARY_DESC_TOP   = 1965960
SUMMARY_DESC_W     = 11277600
SUMMARY_DESC_H     = 274320

# 콘텐츠 새 시작 top (요약 설명 bottom + gap)
CONTENT_NEW_TOP = 2560320   # 2.80" — 요약 설명 bottom(2240280) + 0.35" gap
BODY_LIMIT      = 6401280   # 7.0" (하단 여백 0.3" 이상)
SLIDE_H         = 6858000   # 7.5"


def add_textbox(slide, left, top, width, height, text, font_name, font_size_pt,
                font_color, bold=False):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    bodyPr = txBox._element.find(".//a:bodyPr", nsmap)
    if bodyPr is not None:
        bodyPr.set("lIns", "109728")
        bodyPr.set("rIns", "109728")
        bodyPr.set("tIns", "109728")
        bodyPr.set("bIns", "109728")
        bodyPr.set("anchor", "t")
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def shift_and_resize_content(slide, header_names, content_new_top, max_bottom):
    """Shift all non-header shapes down to content_new_top, shrink if overflow."""
    content_shapes = [s for s in slide.shapes if s.name not in header_names]
    if not content_shapes:
        return

    # Find current content top (minimum top among content shapes)
    current_top = min(s.top for s in content_shapes)
    shift = content_new_top - current_top

    for shape in content_shapes:
        new_top = shape.top + shift
        shape.top = new_top

        # If bottom exceeds limit, shrink height
        new_bottom = new_top + shape.height
        if new_bottom > max_bottom:
            new_height = max_bottom - new_top
            if new_height > 0:
                shape.height = new_height


def process_slide(slide, idx, title_text, desc_text):
    print(f"\n--- Processing slide idx {idx} ---")
    header_names = {"Text Placeholder 1", "TextBox 2", "TextBox 3"}

    # Check if summary already exists
    existing_tops = [s.top for s in slide.shapes if s.name not in header_names]
    print(f"  Content shapes: {len(existing_tops)}, min top: {min(existing_tops)/914400:.2f}\"")

    # Shift content shapes
    shift_and_resize_content(slide, header_names, CONTENT_NEW_TOP, BODY_LIMIT)

    # Verify after shift
    content_shapes = [s for s in slide.shapes if s.name not in header_names]
    bottoms = [s.top + s.height for s in content_shapes]
    print(f"  After shift — min top: {min(s.top for s in content_shapes)/914400:.2f}\", max bottom: {max(bottoms)/914400:.2f}\"")

    # Add summary title
    add_textbox(slide,
                SUMMARY_TITLE_LEFT, SUMMARY_TITLE_TOP,
                SUMMARY_TITLE_W, SUMMARY_TITLE_H,
                title_text, "Freesentation", 16, PRIMARY, bold=True)
    print(f"  Added summary title: \"{title_text}\"")

    # Add summary desc
    add_textbox(slide,
                SUMMARY_DESC_LEFT, SUMMARY_DESC_TOP,
                SUMMARY_DESC_W, SUMMARY_DESC_H,
                desc_text, "Freesentation", 13, DARK_GRAY, bold=False)
    print(f"  Added summary desc: \"{desc_text}\"")

    print(f"  Total shapes: {len(slide.shapes)}")


def main():
    prs = Presentation(INPUT_PATH)

    slides_config = [
        {
            "idx": 5,  # L01 Bento Grid
            "title": "MSK 핵심 아키텍처와 클러스터 유형별 특성",
            "desc": "아키텍처 개요, 클러스터 유형(Provisioned/Serverless/Express), 보안 체계 비교",
        },
        {
            "idx": 8,  # L04 Process Arrow
            "title": "MSK 도입 4단계 프로세스",
            "desc": "설계 → 생성 → 연결 → 운영 단계별 핵심 작업 및 체크포인트",
        },
        {
            "idx": 9,  # L05 Phased Columns
            "title": "MSK 전환 4단계 로드맵",
            "desc": "평가 → 설계 → 구축 → 운영 단계별 주요 액션과 산출물",
        },
    ]

    for cfg in slides_config:
        slide = prs.slides[cfg["idx"]]
        process_slide(slide, cfg["idx"], cfg["title"], cfg["desc"])

    prs.save(OUTPUT_PATH)
    print(f"\nSaved to {OUTPUT_PATH}")

    # Verify
    prs2 = Presentation(OUTPUT_PATH)
    for cfg in slides_config:
        slide = prs2.slides[cfg["idx"]]
        shapes = slide.shapes
        content = [s for s in shapes if s.name not in {"Text Placeholder 1","TextBox 2","TextBox 3"}]
        bottoms = [s.top + s.height for s in content]
        min_top = min(s.top for s in content) / 914400
        max_bottom = max(bottoms) / 914400
        margin = (6858000 - max(bottoms)) / 914400
        print(f"[idx {cfg['idx']}] shapes={len(shapes)}, content_top={min_top:.2f}\", max_bottom={max_bottom:.2f}\", margin={margin:.2f}\"")


if __name__ == "__main__":
    main()
