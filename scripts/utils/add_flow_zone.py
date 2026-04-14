from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import shutil, copy

INPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
TEMP  = "results/pptx/AWS_MSK_Expert_Intro_tmp.pptx"

PRIMARY   = RGBColor(0, 67, 218)
LIGHT_BG  = RGBColor(248, 249, 250)  # #F8F9FA
DARK_GRAY = RGBColor(33, 33, 33)

# Flow diagram zone EMU
FLOW_LEFT   = 457200
FLOW_TOP    = 2332260   # 2.55"
FLOW_WIDTH  = 11277600
FLOW_HEIGHT = 822960    # 0.90"

# Content shift: 1.10" = 1005840 EMU
SHIFT = 1005840
MAX_BOTTOM = 6401280    # 7.0"

EXEMPT = {"Text Placeholder 1", "TextBox 2", "TextBox 3", "TextBox 20", "TextBox 21"}

SLIDE_CONFIGS = {
    8: {
        "flow_text": "1. 설계    →    2. 생성    →    3. 연결    →    4. 운영",
        "content_shapes": [
            "Chevron 4","TextBox 5","Chevron 8","TextBox 9",
            "Chevron 12","TextBox 13","Chevron 16","TextBox 17",
            "Rounded Rectangle 6","TextBox 7",
            "Rounded Rectangle 10","TextBox 11",
            "Rounded Rectangle 14","TextBox 15",
            "Rounded Rectangle 18","TextBox 19",
        ],
        "card_shapes": ["Rounded Rectangle 6","Rounded Rectangle 10","Rounded Rectangle 14","Rounded Rectangle 18"],
        "card_tbs":    ["TextBox 7","TextBox 11","TextBox 15","TextBox 19"],
    },
    9: {
        "flow_text": "1. 평가    →    2. 설계    →    3. 구축    →    4. 운영",
        "content_shapes": [
            "Rounded Rectangle 4","TextBox 5","Rounded Rectangle 8","TextBox 9",
            "Rounded Rectangle 12","TextBox 13","Rounded Rectangle 16","TextBox 17",
            "Rounded Rectangle 6","TextBox 7",
            "Rounded Rectangle 10","TextBox 11",
            "Rounded Rectangle 14","TextBox 15",
            "Rounded Rectangle 18","TextBox 19",
        ],
        "card_shapes": ["Rounded Rectangle 6","Rounded Rectangle 10","Rounded Rectangle 14","Rounded Rectangle 18"],
        "card_tbs":    ["TextBox 7","TextBox 11","TextBox 15","TextBox 19"],
    },
}

nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

def remove_picture(slide):
    to_remove = [s for s in slide.shapes if s.shape_type == 13]
    for pic in to_remove:
        sp = pic._element
        sp.getparent().remove(sp)
        print(f"  Removed picture: {pic.name}")

def add_flow_zone(slide, flow_text):
    # Add rounded rectangle
    rr = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = freeform; use 5 for RoundedRectangle
        Emu(FLOW_LEFT), Emu(FLOW_TOP), Emu(FLOW_WIDTH), Emu(FLOW_HEIGHT)
    )
    # Set fill
    rr.fill.solid()
    rr.fill.fore_color.rgb = LIGHT_BG
    # Set border PRIMARY
    rr.line.color.rgb = PRIMARY
    rr.line.width = Pt(1.0)
    # Set corner radius via XML
    spPr = rr._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr")
    prstGeom = spPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
    if prstGeom is not None:
        avLst = prstGeom.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", "val 16667")  # ~moderate corner radius

    # Add text
    tf = rr.text_frame
    tf.word_wrap = False
    bodyPr = rr._element.find(".//a:bodyPr", nsmap)
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")
        bodyPr.set("lIns", "182880")
        bodyPr.set("rIns", "182880")
        bodyPr.set("tIns", "91440")
        bodyPr.set("bIns", "91440")

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = flow_text
    run.font.name = "Freesentation"
    run.font.size = Pt(14)
    run.font.color.rgb = PRIMARY
    run.font.bold = True
    print(f"  Flow zone added: '{flow_text}'")
    return rr

def shift_and_shrink(slide, cfg):
    names = set(cfg["content_shapes"])
    for shape in slide.shapes:
        if shape.name in names:
            shape.top += SHIFT
    print(f"  Shifted {len(names)} shapes by {SHIFT/914400:.2f}\"")

    # Shrink cards if overflow
    card_map = {s.name: s for s in slide.shapes if s.name in set(cfg["card_shapes"])}
    tb_map   = {s.name: s for s in slide.shapes if s.name in set(cfg["card_tbs"])}

    for card_name, card in card_map.items():
        bottom = card.top + card.height
        if bottom > MAX_BOTTOM:
            excess = bottom - MAX_BOTTOM
            card.height -= excess
            # find matching TB
            idx = cfg["card_shapes"].index(card_name)
            tb_name = cfg["card_tbs"][idx]
            if tb_name in tb_map:
                tb = tb_map[tb_name]
                tb.height = max(tb.height - excess, Emu(457200))
            print(f"  Shrunk {card_name} by {excess/914400:.2f}\" → bottom={card.top/914400+card.height/914400:.2f}\"")

prs = Presentation(INPUT)

for slide_idx, cfg in SLIDE_CONFIGS.items():
    slide = prs.slides[slide_idx]
    print(f"\n=== Slide idx {slide_idx} ===")
    remove_picture(slide)
    shift_and_shrink(slide, cfg)
    add_flow_zone(slide, cfg["flow_text"])

prs.save(TEMP)
shutil.move(TEMP, INPUT)
print(f"\nSaved.")

# Verify
prs2 = Presentation(INPUT)
for slide_idx in [8, 9]:
    slide = prs2.slides[slide_idx]
    print(f"\n[verify idx {slide_idx}]")
    for shape in slide.shapes:
        if shape.name in {"Text Placeholder 1","TextBox 2","TextBox 3"}:
            continue
        top = shape.top/914400
        bot = (shape.top+shape.height)/914400
        print(f"  {shape.name}: top={top:.2f}\" bot={bot:.2f}\" {'⚠ OVERFLOW' if bot > 7.01 else ''}")
