#!/usr/bin/env python3
"""
Comprehensive PPTX analysis with WCAG calculations, whitespace ratios, 
dimension verification, and detailed measurements.
"""
import json
import sys
from pathlib import Path
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Constants
EMU_PER_INCH = 914400
SLIDE_WIDTH = 13.333  # inches
SLIDE_HEIGHT = 7.500  # inches
CONTENT_LIMIT_Y = 7.0  # inches

# Style guide colors (RGB)
COLORS = {
    'PRIMARY': (0, 67, 218),
    'DARK_GRAY': (33, 33, 33),
    'WHITE': (255, 255, 255),
    'ORANGE': (255, 140, 0),
    'GREEN': (34, 139, 34),
    'BG_BOX': (248, 249, 250),
}

def emu_to_inches(emu):
    """Convert EMU to inches."""
    return emu / EMU_PER_INCH

def get_rgb_from_fill(fill):
    """Extract RGB from fill, return None if not solid or no fill."""
    try:
        if fill.type == 1:  # SOLID
            color = fill.fore_color
            if color.type == 1:  # RGB
                return (color.rgb[0], color.rgb[1], color.rgb[2])
    except:
        pass
    return None

def get_text_rgb(run):
    """Extract RGB from text run."""
    try:
        if run.font.color.type == 1:  # RGB
            rgb = run.font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except:
        pass
    return None

def calculate_luminance(rgb):
    """Calculate relative luminance per WCAG 2.1."""
    r, g, b = [x / 255.0 for x in rgb]
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * r + 0.7152 * g + 0.0722 * b

def calculate_contrast_ratio(rgb1, rgb2):
    """Calculate contrast ratio between two RGB colors."""
    l1 = calculate_luminance(rgb1)
    l2 = calculate_luminance(rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def wcag_check(ratio, font_size_pt):
    """Check if contrast ratio meets WCAG AA standard."""
    if font_size_pt >= 18 or (font_size_pt >= 14):  # Large text threshold
        return ratio >= 3.0
    return ratio >= 4.5

def analyze_slide(slide, slide_idx):
    """Analyze a single slide comprehensively."""
    result = {
        'slide_index': slide_idx,
        'shapes': [],
        'contrast_checks': [],
        'overlaps': [],
        'whitespace_ratio': 0.0,
        'shapes_exceeding_7_0': [],
        'subtitle_violations': []
    }
    
    # Content area (excluding title area ~1.5")
    content_area = SLIDE_WIDTH * (SLIDE_HEIGHT - 1.5)
    total_shape_area = 0.0
    
    shapes_data = []
    
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            'index': idx,
            'name': shape.name,
            'type': str(shape.shape_type),
            'left_inches': emu_to_inches(shape.left),
            'top_inches': emu_to_inches(shape.top),
            'width_inches': emu_to_inches(shape.width),
            'height_inches': emu_to_inches(shape.height),
            'bottom_inches': emu_to_inches(shape.top + shape.height),
            'left_emu': shape.left,
            'top_emu': shape.top,
            'width_emu': shape.width,
            'height_emu': shape.height,
        }
        
        # Check if exceeds 7.0"
        if shape_info['bottom_inches'] > CONTENT_LIMIT_Y:
            result['shapes_exceeding_7_0'].append({
                'shape_index': idx,
                'shape_name': shape.name,
                'bottom_inches': round(shape_info['bottom_inches'], 3),
                'overflow': round(shape_info['bottom_inches'] - CONTENT_LIMIT_Y, 3)
            })
        
        # Calculate shape area for whitespace
        if shape_info['top_inches'] > 1.5:  # Below title area
            shape_area = shape_info['width_inches'] * shape_info['height_inches']
            total_shape_area += shape_area
        
        # Extract fill color
        fill_rgb = None
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                fill_rgb = get_rgb_from_fill(shape.fill)
        except:
            pass
        shape_info['fill_rgb'] = fill_rgb
        
        # Extract text and formatting
        if shape.has_text_frame:
            text_frame = shape.text_frame
            shape_info['text_content'] = text_frame.text
            shape_info['text_runs'] = []
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        text_rgb = get_text_rgb(run)
                        font_size = run.font.size.pt if run.font.size else None
                        font_name = run.font.name
                        bold = run.font.bold
                        
                        run_info = {
                            'text': run.text,
                            'font_name': font_name,
                            'font_size_pt': font_size,
                            'bold': bold,
                            'rgb': text_rgb
                        }
                        shape_info['text_runs'].append(run_info)
                        
                        # Contrast check
                        if text_rgb and fill_rgb and font_size:
                            ratio = calculate_contrast_ratio(text_rgb, fill_rgb)
                            passes = wcag_check(ratio, font_size)
                            result['contrast_checks'].append({
                                'slide': slide_idx,
                                'shape_index': idx,
                                'shape_name': shape.name,
                                'text_rgb': text_rgb,
                                'bg_rgb': fill_rgb,
                                'text_luminance': round(calculate_luminance(text_rgb), 3),
                                'bg_luminance': round(calculate_luminance(fill_rgb), 3),
                                'contrast_ratio': round(ratio, 2),
                                'font_size_pt': font_size,
                                'wcag_aa_pass': passes,
                                'required_ratio': 3.0 if font_size >= 18 else 4.5
                            })
        
        shapes_data.append(shape_info)
        result['shapes'].append(shape_info)
    
    # Calculate whitespace ratio
    result['whitespace_ratio'] = round(1 - (total_shape_area / content_area), 3) if content_area > 0 else 0.0
    result['whitespace_pass'] = result['whitespace_ratio'] >= 0.20
    
    # Detect overlaps
    for i, s1 in enumerate(shapes_data):
        for j, s2 in enumerate(shapes_data[i+1:], i+1):
            # Check bounding box overlap
            left1, right1 = s1['left_emu'], s1['left_emu'] + s1['width_emu']
            top1, bottom1 = s1['top_emu'], s1['top_emu'] + s1['height_emu']
            left2, right2 = s2['left_emu'], s2['left_emu'] + s2['width_emu']
            top2, bottom2 = s2['top_emu'], s2['top_emu'] + s2['height_emu']
            
            if (left1 < right2 and right1 > left2 and top1 < bottom2 and bottom1 > top2):
                result['overlaps'].append({
                    'shape1_index': i,
                    'shape1_name': s1['name'],
                    'shape2_index': j,
                    'shape2_name': s2['name'],
                    'shape1_bounds': f"({s1['left_inches']:.2f}, {s1['top_inches']:.2f}, {s1['width_inches']:.2f}, {s1['height_inches']:.2f})",
                    'shape2_bounds': f"({s2['left_inches']:.2f}, {s2['top_inches']:.2f}, {s2['width_inches']:.2f}, {s2['height_inches']:.2f})"
                })
    
    return result

def main():
    pptx_path = 'results/pptx/AWS_MSK_Expert_Intro.pptx'
    output_path = 'results/pptx/measurement_data.json'
    
    if not Path(pptx_path).exists():
        print(f"ERROR: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)
    
    prs = Presentation(pptx_path)
    
    analysis = {
        'file': pptx_path,
        'slide_count': len(prs.slides),
        'slide_dimensions': {
            'width_inches': SLIDE_WIDTH,
            'height_inches': SLIDE_HEIGHT,
            'width_emu': prs.slide_width,
            'height_emu': prs.slide_height
        },
        'slides': []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_analysis = analyze_slide(slide, slide_idx)
        analysis['slides'].append(slide_analysis)
    
    # Write output
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(analysis, f, indent=2, ensure_ascii=False)
    
    print(f"Analysis complete. Output saved to {output_path}")
    
    # Summary stats
    total_contrast_checks = sum(len(s['contrast_checks']) for s in analysis['slides'])
    failed_contrast = sum(1 for s in analysis['slides'] for c in s['contrast_checks'] if not c['wcag_aa_pass'])
    total_exceeding = sum(len(s['shapes_exceeding_7_0']) for s in analysis['slides'])
    total_overlaps = sum(len(s['overlaps']) for s in analysis['slides'])
    
    print(f"\nSummary:")
    print(f"  Total contrast checks: {total_contrast_checks}")
    print(f"  Failed WCAG AA: {failed_contrast}")
    print(f"  Shapes exceeding 7.0\": {total_exceeding}")
    print(f"  Total overlaps detected: {total_overlaps}")

if __name__ == '__main__':
    main()
