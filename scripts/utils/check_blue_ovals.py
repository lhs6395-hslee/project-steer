#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Check Slide 5 for blue oval fallback icons"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE

EMU_TO_INCHES = 914400

def emu_to_inches(emu):
    return emu / EMU_TO_INCHES

def rgb_to_tuple(rgb_color):
    try:
        if hasattr(rgb_color, 'rgb'):
            return (rgb_color.rgb[0], rgb_color.rgb[1], rgb_color.rgb[2])
    except:
        pass
    return None

prs = Presentation('results/pptx/AWS_MSK_Expert_Intro.pptx')

# Slide 5 is index 4 (0-based)
slide = prs.slides[4]

print("SLIDE 5 (3-1. 운영 전략) - Checking for blue oval fallback icons")
print("="*80)

oval_count = 0

for idx, shape in enumerate(slide.shapes):
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            # Check if it's an oval
            is_oval = False
            if hasattr(shape, 'auto_shape_type'):
                shape_type_str = str(shape.auto_shape_type)
                if 'OVAL' in shape_type_str or shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    is_oval = True
            
            if is_oval:
                top = emu_to_inches(shape.top)
                left = emu_to_inches(shape.left)
                width = emu_to_inches(shape.width)
                height = emu_to_inches(shape.height)
                
                # Check fill color
                fill_color = None
                if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL_TYPE.SOLID:
                    fill_color = rgb_to_tuple(shape.fill.fore_color)
                
                # Check if has text
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                
                print(f"\nOval shape {idx} at ({left:.2f}\", {top:.2f}\"):")
                print(f"  Size: {width:.2f}\" × {height:.2f}\"")
                print(f"  Fill color: {fill_color}")
                print(f"  Text: '{text}'")
                
                # Check if it's a blue fallback icon
                if fill_color:
                    r, g, b = fill_color
                    if b > 200 and r < 100 and g < 150:  # Blue-ish
                        oval_count += 1
                        print(f"  ⚠️  DETECTED: Blue oval fallback icon!")
                        if text and len(text) <= 3:
                            print(f"  ⚠️  CRITICAL: Contains text symbol '{text}' - this is a fallback icon")
        except Exception as e:
            pass

print("\n" + "="*80)
print(f"Total blue oval fallback icons: {oval_count}")
if oval_count > 0:
    print("❌ CRITICAL FAIL: Blue oval fallback icons are explicitly FORBIDDEN")
    print("   Per constraint: 'No blue oval + text fallback icons allowed (CRITICAL FAIL)'")
    print("   These must be replaced with actual PNG icons from icons/ folder")
