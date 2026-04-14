#!/usr/bin/env python3
"""Check bottom margins for all slides."""

import sys
from pptx import Presentation
import json

def emu_to_inches(emu):
    return emu / 914400

def check_bottom_margins(prs_path):
    prs = Presentation(prs_path)
    results = []
    
    SLIDE_H = 7.500
    MIN_MARGIN = 0.3
    
    for slide_idx, slide in enumerate(prs.slides):
        max_bottom = 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'height'):
                bottom = emu_to_inches(shape.top + shape.height)
                if bottom > max_bottom:
                    max_bottom = bottom
        
        margin = SLIDE_H - max_bottom
        
        results.append({
            "slide_index": slide_idx,
            "max_bottom": f"{max_bottom:.3f}\"",
            "bottom_margin": f"{margin:.3f}\"",
            "min_required": f"{MIN_MARGIN}\"",
            "pass": margin >= MIN_MARGIN
        })
    
    return results

if __name__ == "__main__":
    results = check_bottom_margins(sys.argv[1])
    print(json.dumps(results, indent=2))
