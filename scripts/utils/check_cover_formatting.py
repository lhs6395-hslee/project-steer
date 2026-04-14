#!/usr/bin/env python3
"""
Check cover slide formatting details including:
- Title/subtitle positioning and dimensions
- Date format
- Scheme color preservation
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def emu_to_inches(emu):
    """Convert EMU to inches"""
    return emu / 914400

def check_cover_formatting(pptx_path):
    """Check cover slide formatting against acceptance criteria"""
    prs = Presentation(pptx_path)
    
    if len(prs.slides) == 0:
        return {"error": "No slides found"}
    
    cover_slide = prs.slides[0]
    results = {
        "title": {},
        "subtitle": {},
        "date": {},
        "scheme_colors": {},
        "issues": []
    }
    
    # Find title and subtitle textboxes
    title_shape = None
    subtitle_shape = None
    date_shapes = []
    
    for idx, shape in enumerate(cover_slide.shapes):
        if shape.shape_type == 17:  # TEXT_BOX
            # Check by position - title should be around top=1.63", subtitle around top=4.18"
            top_inches = emu_to_inches(shape.top)
            left_inches = emu_to_inches(shape.left)
            width_inches = emu_to_inches(shape.width)
            height_inches = emu_to_inches(shape.height)
            
            # Get text content
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text
            
            # Identify shape by position and size
            if 1.4 < top_inches < 2.0 and width_inches > 8.0:
                # Likely title
                title_shape = shape
                results["title"] = {
                    "index": idx,
                    "left_inches": left_inches,
                    "top_inches": top_inches,
                    "width_inches": width_inches,
                    "height_inches": height_inches,
                    "text": text,
                    "left_expected": 1.67,
                    "width_expected": 10.0,
                    "left_match": abs(left_inches - 1.67) < 0.2,
                    "width_match": abs(width_inches - 10.0) < 0.5
                }
            elif 3.8 < top_inches < 5.0 and width_inches > 8.0:
                # Likely subtitle
                subtitle_shape = shape
                results["subtitle"] = {
                    "index": idx,
                    "left_inches": left_inches,
                    "top_inches": top_inches,
                    "width_inches": width_inches,
                    "height_inches": height_inches,
                    "text": text
                }
            elif height_inches < 0.5 and width_inches < 2.5:
                # Likely date element
                date_shapes.append({
                    "index": idx,
                    "text": text,
                    "left_inches": left_inches,
                    "top_inches": top_inches
                })
    
    # Check title-subtitle gap
    if title_shape and subtitle_shape:
        title_bottom = emu_to_inches(title_shape.top + title_shape.height)
        subtitle_top = emu_to_inches(subtitle_shape.top)
        gap = subtitle_top - title_bottom
        results["title_subtitle_gap"] = {
            "gap_inches": gap,
            "expected": 0.15,
            "match": abs(gap - 0.15) < 0.05
        }
        
        if abs(gap - 0.15) > 0.05:
            results["issues"].append(f"Title-subtitle gap is {gap:.3f}\" (expected 0.15\")")
    
    # Check title positioning
    if "title" in results and results["title"]:
        if not results["title"]["left_match"]:
            results["issues"].append(
                f"Title left position is {results['title']['left_inches']:.3f}\" (expected 1.67\")"
            )
        if not results["title"]["width_match"]:
            results["issues"].append(
                f"Title width is {results['title']['width_inches']:.3f}\" (expected 10.0\")"
            )
    
    # Check date format
    results["date"]["shapes"] = date_shapes
    date_texts = [d["text"] for d in date_shapes]
    
    # Check if date format is MM/DD
    date_issues = []
    for d in date_shapes:
        text = d["text"]
        # Check if it contains year as separate line
        if "2026" in text or "\n" in text:
            date_issues.append(f"Date contains year or newline: '{text}'")
        # Check if it matches MM/DD format
        if "/" in text and len(text.strip()) <= 5:
            # Looks like MM/DD format
            pass
        elif text.strip().isdigit() and len(text.strip()) == 4:
            # Year only
            date_issues.append(f"Date shape contains year only: '{text}'")
    
    if date_issues:
        results["date"]["issues"] = date_issues
        results["issues"].extend(date_issues)
    
    # Check scheme colors (need to check text runs)
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        color_info = []
        for para in tf.paragraphs:
            for run in para.runs:
                font = run.font
                if font.color.type:
                    color_type = str(font.color.type)
                    color_info.append({
                        "text": run.text[:30],
                        "color_type": color_type,
                        "rgb": str(font.color.rgb) if hasattr(font.color, 'rgb') else None
                    })
        results["scheme_colors"]["title"] = color_info
        
        # Check if scheme colors are preserved (should not be RGB, should be scheme)
        if color_info:
            for ci in color_info:
                if "SCHEME" not in ci["color_type"]:
                    results["issues"].append(
                        f"Title text '{ci['text']}...' uses {ci['color_type']} instead of SCHEME color"
                    )
    
    # Check subtitle lines
    if subtitle_shape and subtitle_shape.has_text_frame:
        subtitle_text = subtitle_shape.text_frame.text
        lines = subtitle_text.split('\n')
        results["subtitle"]["line_count"] = len(lines)
        results["subtitle"]["lines"] = lines
        
        if len(lines) > 2:
            results["issues"].append(
                f"Subtitle has {len(lines)} lines (max 2 allowed): {lines}"
            )
    
    return results

if __name__ == "__main__":
    import sys
    import json
    
    if len(sys.argv) < 2:
        print("Usage: python check_cover_formatting.py <pptx_path>")
        sys.exit(1)
    
    results = check_cover_formatting(sys.argv[1])
    print(json.dumps(results, indent=2, ensure_ascii=False))
