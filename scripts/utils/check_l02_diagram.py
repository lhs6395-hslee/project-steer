#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Check L02 Three Cards for required top diagram"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_TO_INCHES = 914400

def emu_to_inches(emu):
    return emu / EMU_TO_INCHES

prs = Presentation('results/pptx/AWS_MSK_Expert_Intro.pptx')

# Slide 7 is index 6 (0-based)
slide = prs.slides[6]

print("SLIDE 7 (L02 Three Cards) - Checking for top diagram zone")
print("="*80)

for idx, shape in enumerate(slide.shapes):
    top = emu_to_inches(shape.top)
    left = emu_to_inches(shape.left)
    width = emu_to_inches(shape.width)
    height = emu_to_inches(shape.height)
    bottom = top + height
    
    # Check if in potential diagram zone (top of body area)
    if 2.0 <= top <= 3.5:
        shape_type = str(shape.shape_type)
        print(f"Shape {idx}: {shape_type}")
        print(f"  Position: ({left:.2f}\", {top:.2f}\") Size: {width:.2f}\" × {height:.2f}\"")
        print(f"  Bottom: {bottom:.2f}\"")
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"  ✓ GROUP detected - likely diagram")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if width > 2.0:
                print(f"  ✓ PICTURE detected - likely diagram/flowchart image")
        
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"  Text: {text[:80]}")
        print()

# Check for cards in lower area
print("\nChecking for three cards in lower area (4.0\" - 7.0\"):")
cards = []
for shape in slide.shapes:
    top = emu_to_inches(shape.top)
    if 4.0 <= top <= 7.0 and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        cards.append(shape)

print(f"Found {len(cards)} card-like shapes")

# Final verdict
print("\n" + "="*80)
if any(s.shape_type == MSO_SHAPE_TYPE.GROUP and 2.0 <= emu_to_inches(s.top) <= 3.5 for s in slide.shapes):
    print("✓ L02 PASS: Top diagram zone present (GROUP)")
elif any(s.shape_type == MSO_SHAPE_TYPE.PICTURE and 2.0 <= emu_to_inches(s.top) <= 3.5 and emu_to_inches(s.width) > 2.0 for s in slide.shapes):
    print("✓ L02 PASS: Top diagram zone present (PICTURE)")
else:
    print("❌ L02 FAIL: Top diagram zone REQUIRED but not found")
