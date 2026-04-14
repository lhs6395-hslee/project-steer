#!/usr/bin/env python3
"""
Text overflow and boundary checker for PPTX presentations.
Validates body content text overflow and shape boundaries.
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import json

# Constants (from pptx_style_guide.md)
SLIDE_W = 13.333  # inches
SLIDE_H = 7.500   # inches
BODY_START_Y = 2.0  # inches
BODY_LIMIT_Y = 7.0  # inches
MIN_BOTTOM_MARGIN = 0.3  # inches

# Font size and line height estimates
FONT_SIZES = {
    'title': 28,  # pt
    'subtitle': 28,  # pt
    'card_title': 14,  # pt
    'card_body': 13,  # pt
    'body_title': 16,  # pt
    'body_text': 13,  # pt
}

def emu_to_inches(emu):
    """Convert EMU to inches."""
    return emu / 914400

def check_text_overflow(prs_path):
    """Check for text overflow and boundary violations."""
    prs = Presentation(prs_path)
    
    issues = []
    all_checks_passed = True
    
    # Skip title slide (0) and TOC (1) and thank you slide (10)
    content_slides = range(2, 10)
    
    for slide_idx in content_slides:
        slide = prs.slides[slide_idx]
        slide_issues = []
        
        # Track lowest point on slide
        max_bottom = 0
        
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
                
            # Get shape boundaries
            top_inches = emu_to_inches(shape.top)
            height_inches = emu_to_inches(shape.height)
            bottom_inches = top_inches + height_inches
            
            # Update max bottom
            if bottom_inches > max_bottom:
                max_bottom = bottom_inches
            
            # Check if shape is in body area (skip title/subtitle)
            if top_inches < 1.5:  # Skip header area shapes
                continue
            
            # Check vertical bounds (2.0" to 7.0")
            if top_inches < BODY_START_Y:
                slide_issues.append({
                    "shape_idx": shape_idx,
                    "shape_name": shape.name,
                    "issue": "shape_above_body_start",
                    "top": f"{top_inches:.3f}\"",
                    "body_start": f"{BODY_START_Y}\"",
                    "severity": "CRITICAL"
                })
                all_checks_passed = False
            
            if bottom_inches > BODY_LIMIT_Y:
                slide_issues.append({
                    "shape_idx": shape_idx,
                    "shape_name": shape.name,
                    "issue": "shape_exceeds_body_limit",
                    "bottom": f"{bottom_inches:.3f}\"",
                    "body_limit": f"{BODY_LIMIT_Y}\"",
                    "overflow": f"{bottom_inches - BODY_LIMIT_Y:.3f}\"",
                    "severity": "CRITICAL"
                })
                all_checks_passed = False
            
            # Check text overflow within shape
            text_frame = shape.text_frame
            if text_frame.text.strip():
                # Count actual line breaks in text
                text_lines = text_frame.text.count('\n') + 1
                
                # Estimate line height based on font size
                # Get font size from first paragraph
                if text_frame.paragraphs:
                    first_run = None
                    for para in text_frame.paragraphs:
                        if para.runs:
                            first_run = para.runs[0]
                            break
                    
                    if first_run and first_run.font.size:
                        font_size_pt = first_run.font.size.pt
                    else:
                        # Use default based on shape type/position
                        if top_inches < 1.8:
                            font_size_pt = FONT_SIZES['subtitle']
                        elif bottom_inches - top_inches < 0.6:
                            font_size_pt = FONT_SIZES['card_title']
                        else:
                            font_size_pt = FONT_SIZES['card_body']
                    
                    # Line height is typically 1.2x font size
                    line_height_inches = (font_size_pt * 1.2) / 72
                    estimated_text_height = text_lines * line_height_inches
                    
                    # Add padding (typical textbox inset)
                    padding_inches = 0.12 * 2  # top + bottom
                    
                    if estimated_text_height + padding_inches > height_inches:
                        overflow = estimated_text_height + padding_inches - height_inches
                        slide_issues.append({
                            "shape_idx": shape_idx,
                            "shape_name": shape.name,
                            "issue": "text_overflow",
                            "text_lines": text_lines,
                            "estimated_height": f"{estimated_text_height:.3f}\"",
                            "shape_height": f"{height_inches:.3f}\"",
                            "overflow": f"{overflow:.3f}\"",
                            "font_size": f"{font_size_pt}pt",
                            "severity": "CRITICAL",
                            "text_preview": text_frame.text[:100]
                        })
                        all_checks_passed = False
            
            # Check subtitle constraints (shape in subtitle area: 0.5" to 1.5")
            if 0.5 <= top_inches <= 1.5 and height_inches < 2.0:
                text = text_frame.text.strip()
                line_count = text.count('\n') + 1
                
                # Check max 2 lines
                if line_count > 2:
                    slide_issues.append({
                        "shape_idx": shape_idx,
                        "shape_name": shape.name,
                        "issue": "subtitle_too_many_lines",
                        "line_count": line_count,
                        "max_allowed": 2,
                        "severity": "CRITICAL",
                        "text": text
                    })
                    all_checks_passed = False
                
                # Check for mid-word line breaks
                if '\n' in text:
                    lines = text.split('\n')
                    for i, line in enumerate(lines[:-1]):
                        # Check if line ends with incomplete word
                        if line and not line[-1].isspace():
                            # Check if next line starts without space
                            if i + 1 < len(lines) and lines[i + 1] and not lines[i + 1][0].isspace():
                                slide_issues.append({
                                    "shape_idx": shape_idx,
                                    "shape_name": shape.name,
                                    "issue": "subtitle_mid_word_break",
                                    "line_end": line[-20:],
                                    "next_line_start": lines[i + 1][:20],
                                    "severity": "CRITICAL"
                                })
                                all_checks_passed = False
        
        # Check bottom margin
        bottom_margin = SLIDE_H - max_bottom
        if bottom_margin < MIN_BOTTOM_MARGIN:
            slide_issues.append({
                "issue": "insufficient_bottom_margin",
                "margin": f"{bottom_margin:.3f}\"",
                "required": f"{MIN_BOTTOM_MARGIN}\"",
                "max_bottom": f"{max_bottom:.3f}\"",
                "severity": "CRITICAL"
            })
            all_checks_passed = False
        
        if slide_issues:
            issues.append({
                "slide_index": slide_idx,
                "slide_name": f"Slide {slide_idx + 1}",
                "issues": slide_issues
            })
    
    return {
        "all_checks_passed": all_checks_passed,
        "total_issues": sum(len(s["issues"]) for s in issues),
        "slides_with_issues": len(issues),
        "details": issues
    }

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: check_text_overflow.py <pptx_file>")
        sys.exit(1)
    
    result = check_text_overflow(sys.argv[1])
    print(json.dumps(result, indent=2, ensure_ascii=False))
