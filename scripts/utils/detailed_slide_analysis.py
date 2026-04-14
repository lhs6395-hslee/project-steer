#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Detailed slide-by-slide analysis with layout identification"""

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

EMU_TO_INCHES = 914400

def emu_to_inches(emu):
    return emu / EMU_TO_INCHES

def analyze_slide_structure(slide_idx, slide):
    """Analyze slide structure in detail"""
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*80}")
    
    # Find title and subtitle
    title_text = None
    subtitle_text = None
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text = shape.text_frame.text.strip()
        top = emu_to_inches(shape.top)
        left = emu_to_inches(shape.left)
        
        if 0.5 < top < 1.5:
            if left < 3 and text:
                title_text = text
            elif left > 5 and text:
                subtitle_text = text
    
    print(f"Title: {title_text}")
    print(f"Subtitle: {subtitle_text}")
    
    # Count shapes by type
    pictures = []
    auto_shapes = []
    text_boxes = []
    groups = []
    
    for shape in slide.shapes:
        top = emu_to_inches(shape.top)
        left = emu_to_inches(shape.left)
        width = emu_to_inches(shape.width)
        height = emu_to_inches(shape.height)
        
        # Skip header area
        if top + height < 2.0:
            continue
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append({
                'left': left,
                'top': top,
                'width': width,
                'height': height
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append({
                'left': left,
                'top': top,
                'width': width,
                'height': height
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if shape.has_text_frame:
                text_boxes.append({
                    'left': left,
                    'top': top,
                    'width': width,
                    'height': height,
                    'text': shape.text_frame.text.strip()[:50]
                })
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            groups.append({
                'left': left,
                'top': top,
                'width': width,
                'height': height
            })
    
    print(f"\nContent shapes (body area 2.0\"-7.5\"):")
    print(f"  Pictures: {len(pictures)}")
    print(f"  Auto shapes: {len(auto_shapes)}")
    print(f"  Text boxes: {len(text_boxes)}")
    print(f"  Groups: {len(groups)}")
    
    # Analyze layout pattern
    layout_guess = "UNKNOWN"
    
    if len(auto_shapes) == 3 and len(pictures) <= 3:
        layout_guess = "L02 Three Cards (likely)"
    elif len(auto_shapes) == 4 and len(pictures) >= 4:
        layout_guess = "L03 Grid 2x2 (likely)"
    elif len(auto_shapes) >= 4:
        if subtitle_text and ('Arrow' in subtitle_text or 'Process' in subtitle_text):
            layout_guess = "L04 Process Arrow (likely)"
        elif subtitle_text and ('Column' in subtitle_text or 'Phased' in subtitle_text):
            layout_guess = "L05 Phased Columns (likely)"
    elif len(auto_shapes) >= 2:
        if subtitle_text and 'Bento' in subtitle_text:
            layout_guess = "L01 Bento (likely)"
    
    print(f"\nLayout identification: {layout_guess}")
    
    # Check for specific issues
    issues = []
    
    # Check if pictures are in bottom-right position (icon position)
    for pic in pictures:
        if 10.0 <= pic['left'] <= 11.0 and 6.0 <= pic['top'] <= 6.5:
            if 0.5 <= pic['width'] <= 0.7 and 0.5 <= pic['height'] <= 0.7:
                print(f"  ✓ Bottom-right icon present at ({pic['left']:.1f}\", {pic['top']:.1f}\")")
    
    # Check for top diagram zone (L02 requirement)
    has_top_diagram = any(
        g['top'] >= 2.0 and g['top'] <= 3.5 for g in groups
    ) or any(
        p['top'] >= 2.0 and p['top'] <= 3.5 and p['width'] > 2.0 for p in pictures
    )
    
    if has_top_diagram:
        print(f"  ✓ Top diagram zone present")
    
    # Check content bounds
    all_shapes = pictures + auto_shapes + text_boxes + groups
    if all_shapes:
        max_bottom = max(s['top'] + s['height'] for s in all_shapes)
        print(f"\nContent bounds:")
        print(f"  Max bottom: {max_bottom:.2f}\"")
        print(f"  Bottom margin: {7.5 - max_bottom:.2f}\"")
        
        if max_bottom > 7.0:
            issues.append(f"Content exceeds 7.0\" limit (at {max_bottom:.2f}\")")
        if 7.5 - max_bottom < 0.3:
            issues.append(f"Insufficient bottom margin ({7.5 - max_bottom:.2f}\" < 0.3\")")
    
    if issues:
        print(f"\n⚠️  Issues found:")
        for issue in issues:
            print(f"  - {issue}")
    
    return layout_guess

def main():
    pptx_path = 'results/pptx/AWS_MSK_Expert_Intro.pptx'
    
    try:
        prs = Presentation(pptx_path)
        print(f"Detailed Analysis: {pptx_path}")
        print(f"Total slides: {len(prs.slides)}")
        
        for idx, slide in enumerate(prs.slides):
            analyze_slide_structure(idx, slide)
        
        print(f"\n{'='*80}")
        print("Analysis complete")
        
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == '__main__':
    sys.exit(main())
