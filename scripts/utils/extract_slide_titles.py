#!/usr/bin/env python3
"""Extract slide titles for TOC comparison"""

from pptx import Presentation
import json

def extract_slide_titles(pptx_path):
    prs = Presentation(pptx_path)
    titles = []
    
    for idx, slide in enumerate(prs.slides):
        # Try to find title from placeholders first
        title_text = ""
        
        # Check placeholders
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Look for patterns like "1-1.", "2-1.", "3-1." which indicate section titles
                if text and any(text.startswith(f"{i}-") for i in range(1, 10)):
                    # Extract title after the numbering
                    lines = text.split('\n')
                    if len(lines) >= 2:
                        # Second line is usually the main title
                        title_text = lines[1].strip()
                    else:
                        title_text = text
                    break
        
        titles.append({
            "index": idx,
            "title": title_text if title_text else f"Slide {idx}"
        })
    
    return titles

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python extract_slide_titles.py <pptx_path>")
        sys.exit(1)
    
    titles = extract_slide_titles(sys.argv[1])
    print(json.dumps(titles, indent=2, ensure_ascii=False))
