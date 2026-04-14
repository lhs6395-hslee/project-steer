from pptx import Presentation
from pptx.util import Emu, Inches
from lxml import etree
import shutil, os

INPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
TEMP  = "results/pptx/AWS_MSK_Expert_Intro_tmp.pptx"
ICON  = "icons/process.png"

nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
EXEMPT = {"Text Placeholder 1", "TextBox 2", "TextBox 3"}

def set_anchor(shape, anchor_val):
    if shape.has_text_frame:
        bodyPr = shape._element.find(".//a:bodyPr", nsmap)
        if bodyPr is not None:
            bodyPr.set("anchor", anchor_val)
            return True
    return False

def set_no_fill(shape):
    """Set shape fill to noFill via XML"""
    sp = shape._element
    # Find spPr
    spPr = sp.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr") or \
           sp.find("p:spPr", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
    # Find via tag
    for child in sp:
        if child.tag.endswith("}spPr"):
            spPr = child
            break
    if spPr is None:
        print(f"    WARNING: spPr not found for {shape.name}")
        return False
    # Remove existing fill elements
    for tag in ["{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill",
                "{http://schemas.openxmlformats.org/drawingml/2006/main}noFill"]:
        for elem in spPr.findall(f".//{tag}"):
            spPr.remove(elem)
        for elem in list(spPr):
            if elem.tag == tag:
                spPr.remove(elem)
    # Insert noFill after ln or as first child
    no_fill = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}noFill")
    # Insert at position 1 (after xfrm if present)
    spPr.insert(1, no_fill)
    return True

prs = Presentation(INPUT)

# ── Slide idx 8 (L04 Process Arrow) ──
slide8 = prs.slides[8]
print("=== Slide idx 8: anchor top for content textboxes ===")
CONTENT_TBS_8 = {"TextBox 7", "TextBox 11", "TextBox 15", "TextBox 19"}
for shape in slide8.shapes:
    if shape.name in CONTENT_TBS_8:
        set_anchor(shape, "t")
        print(f"  anchor=t: {shape.name}")

# Add process.png icon bottom-right
print("  Adding process.png icon...")
with open(ICON, "rb") as f:
    icon_data = f.read()
from io import BytesIO
pic = slide8.shapes.add_picture(BytesIO(icon_data),
    Inches(10.5), Inches(6.2), Inches(0.6), Inches(0.6))
print(f"  Icon added: left={pic.left/914400:.2f}\" top={pic.top/914400:.2f}\"")

# ── Slide idx 9 (L05 Phased Columns) ──
slide9 = prs.slides[9]
print("\n=== Slide idx 9: fix fills + anchor top for content textboxes ===")

# Remove bad fills from TextBox 11 and TextBox 21
BAD_FILLS = {"TextBox 11", "TextBox 21"}
CONTENT_TBS_9 = {"TextBox 7", "TextBox 11", "TextBox 15", "TextBox 19"}

for shape in slide9.shapes:
    if shape.name in BAD_FILLS:
        set_no_fill(shape)
        print(f"  noFill set: {shape.name} (was solid)")
    if shape.name in CONTENT_TBS_9:
        set_no_fill(shape)
        set_anchor(shape, "t")
        print(f"  noFill + anchor=t: {shape.name}")

prs.save(TEMP)
shutil.move(TEMP, INPUT)
print(f"\nSaved to {INPUT}")

# Verify
print("\n=== Verification ===")
prs2 = Presentation(INPUT)

slide8v = prs2.slides[8]
print("Slide idx 8 content TBs:")
for shape in slide8v.shapes:
    if shape.name in CONTENT_TBS_8:
        bodyPr = shape._element.find(".//a:bodyPr", nsmap)
        anchor = bodyPr.get("anchor") if bodyPr is not None else "MISSING"
        print(f"  {shape.name}: anchor={anchor}")

slide9v = prs2.slides[9]
print("Slide idx 9 checks:")
for shape in slide9v.shapes:
    if shape.name in {"TextBox 11", "TextBox 21"}:
        try:
            c = shape.fill.fore_color.rgb
            print(f"  {shape.name}: fill=#{c} ← STILL HAS FILL!")
        except:
            print(f"  {shape.name}: fill=noFill ✓")
    if shape.name in CONTENT_TBS_9:
        bodyPr = shape._element.find(".//a:bodyPr", nsmap)
        anchor = bodyPr.get("anchor") if bodyPr is not None else "MISSING"
        print(f"  {shape.name}: anchor={anchor}")
