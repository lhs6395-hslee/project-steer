from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree
import shutil, os

INPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
OUTPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
TEMP = "results/pptx/AWS_MSK_Expert_Intro_tmp.pptx"

EXEMPT = {"Text Placeholder 1", "TextBox 2", "TextBox 3"}
nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

def set_anchor_middle(shape):
    """Set txBody anchor to ctr (middle vertical alignment)"""
    # Try python-pptx text_frame first
    if shape.has_text_frame:
        try:
            bodyPr = shape.text_frame._txBody.find("a:bodyPr", nsmap)
            if bodyPr is None:
                bodyPr = shape._element.find(".//a:bodyPr", nsmap)
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")
                return True
        except:
            pass
    return False

def fix_textbox_width(slide, target_text_partial, new_width_emu):
    """Find textbox containing target_text_partial and fix its width"""
    for shape in slide.shapes:
        if shape.name in EXEMPT:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text
            if target_text_partial in text:
                shape.width = Emu(new_width_emu)
                print(f"  Fixed width of '{shape.name}': text='{text[:40]}', new_width={new_width_emu}")
                return True
    return False

def fix_phase_colors(slide, phase_configs):
    """Fix phase header background colors and text colors.
    phase_configs: list of (shape_name_contains, bg_rgb, text_rgb)
    """
    # Find shapes that are phase headers (colored rectangles with phase text)
    for shape in slide.shapes:
        if shape.name in EXEMPT:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        for (text_match, bg_color, text_color) in phase_configs:
            if text_match in text:
                # Fix background color
                try:
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*bg_color)
                    print(f"  Phase color: '{text}' -> bg={bg_color}")
                except Exception as e:
                    print(f"  Error setting fill for '{text}': {e}")
                # Fix text color
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = RGBColor(*text_color)
                    print(f"  Phase text: '{text}' -> text={text_color}")
                except Exception as e:
                    print(f"  Error setting text color for '{text}': {e}")
                break

def main():
    prs = Presentation(INPUT)

    # --- Slide idx 8 (L04 Process Arrow) ---
    slide8 = prs.slides[8]
    print(f"\n=== Slide idx 8 (L04) ===")

    # 1. Fix anchor=ctr for all non-exempt shapes
    fixed = 0
    for shape in slide8.shapes:
        if shape.name in EXEMPT:
            continue
        if set_anchor_middle(shape):
            fixed += 1
    print(f"  anchor=ctr fixed: {fixed} shapes")

    # 2. Fix body title width (already w=11277600? verify and fix)
    fix_textbox_width(slide8, "MSK 도입", 11277600)

    # Also check ALL textboxes width
    for shape in slide8.shapes:
        if shape.name in EXEMPT:
            continue
        if shape.has_text_frame and shape.width < Emu(11000000):
            text = shape.text_frame.text[:30]
            print(f"  Shape '{shape.name}': width={shape.width/914400:.2f}\", text='{text}'")

    # --- Slide idx 9 (L05 Phased Columns) ---
    slide9 = prs.slides[9]
    print(f"\n=== Slide idx 9 (L05) ===")

    # 1. Fix anchor=ctr for all non-exempt shapes
    fixed = 0
    for shape in slide9.shapes:
        if shape.name in EXEMPT:
            continue
        if set_anchor_middle(shape):
            fixed += 1
    print(f"  anchor=ctr fixed: {fixed} shapes")

    # 2. Fix body title width
    fix_textbox_width(slide9, "MSK 전환", 11277600)

    # 3. Fix phase colors — 4 phases with distinct colors
    # SUB_ORANGE (238,129,80): contrast vs WHITE = 2.55:1 FAIL → use BLACK text
    # SUB_GREEN (76,184,143): contrast vs WHITE = 3.04:1 FAIL → use BLACK text
    phase_configs = [
        ("Phase 1", (0, 27, 94), (255, 255, 255)),     # DARK_NAVY bg, WHITE text 16.75:1 ✓
        ("Phase 2", (0, 67, 218), (255, 255, 255)),    # PRIMARY bg, WHITE text 8.59:1 ✓
        ("Phase 3", (238, 129, 80), (0, 0, 0)),        # SUB_ORANGE bg, BLACK text
        ("Phase 4", (76, 184, 143), (0, 0, 0)),        # SUB_GREEN bg, BLACK text
        # Also try Korean phase text matches
        ("평가", (0, 27, 94), (255, 255, 255)),
        ("설계", (0, 67, 218), (255, 255, 255)),
        ("구축", (238, 129, 80), (0, 0, 0)),
        ("운영", (76, 184, 143), (0, 0, 0)),
    ]
    fix_phase_colors(slide9, phase_configs)

    # Save atomically
    prs.save(TEMP)
    shutil.move(TEMP, OUTPUT)
    print(f"\nSaved to {OUTPUT}")

    # Verification
    prs2 = Presentation(OUTPUT)

    print("\n--- Verification: idx 8 ---")
    slide8v = prs2.slides[8]
    for shape in slide8v.shapes:
        if shape.name in EXEMPT or not shape.has_text_frame:
            continue
        bodyPr = shape._element.find(".//a:bodyPr", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        anchor = bodyPr.get("anchor") if bodyPr is not None else "MISSING"
        text = shape.text_frame.text[:30]
        print(f"  '{shape.name}': anchor={anchor}, text='{text}'")

    print("\n--- Verification: idx 9 ---")
    slide9v = prs2.slides[9]
    for shape in slide9v.shapes:
        if shape.name in EXEMPT or not shape.has_text_frame:
            continue
        bodyPr = shape._element.find(".//a:bodyPr", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        anchor = bodyPr.get("anchor") if bodyPr is not None else "MISSING"
        text = shape.text_frame.text[:30]
        # Check fill color
        try:
            fill_color = shape.fill.fore_color.rgb
            print(f"  '{shape.name}': anchor={anchor}, fill={fill_color}, text='{text}'")
        except:
            print(f"  '{shape.name}': anchor={anchor}, text='{text}'")

if __name__ == "__main__":
    main()
