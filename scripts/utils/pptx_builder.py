"""Reusable PPTX builder utilities for harness pipeline.

Provides helper functions for common PPTX operations:
- Template loading and format-preserving shape text replacement
- Cover slide with centered title (50pt, 3줄), subtitle, date (MM/DD)
- Body slide creation with standard header
- Shape creation (boxes, badges, arrows)
- Slide deletion, reordering, and section rebuild
- Final save with verification

Body Content Rules:
- 콘텐츠 블록은 슬라이드 좌우 중앙 정렬: start_left = (SLIDE_W - total_w) / 2
- 도형/카드 제목: 프리젠테이션 7 Bold, 14pt, PRIMARY
- 도형/카드 내용: Freesentation, 13pt, DARK_GRAY
- 좌우 여백 대칭 (±0.1" 이내)

Usage:
    from scripts.utils.pptx_builder import PptxBuilder
    
    builder = PptxBuilder("templates/pptx_template.pptx")
    builder.set_cover(
        title_lines=["Amazon Managed Streaming", "for Apache Kafka", "(MSK)"],
        subtitle="Expert Introduction — 아키텍처 · 핵심 기능 · 운영 전략",
        year="2026", date="04/13"
    )
    builder.set_toc(sections=["섹션1", "섹션2", "섹션3"])  # 6개 이상이면 자동 페이징
    slide = builder.add_body_slide("1-1. 제목", "설명글")
    builder.add_box(slide, left=0.5, top=2.2, width=5, height=1.5, text="내용")
    builder.set_ending(tagline="태그라인")
    builder.cleanup_template(keep_indices=[0, 1])
    toc_extra = builder.set_toc(...)  # returns added_toc_count
    builder.rebuild_sections(added_toc_count=toc_extra)
    builder.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os


# ── Design Constants ──
PRIMARY = RGBColor(0, 67, 218)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(33, 33, 33)
GRAY = RGBColor(80, 80, 80)
WHITE = RGBColor(255, 255, 255)
BG_BOX = RGBColor(248, 249, 250)
BORDER = RGBColor(220, 220, 220)

FONT_TITLE = "프리젠테이션 7 Bold"
FONT_DESC = "프리젠테이션 5 Medium"
FONT_BODY = "Freesentation"

# Body content font sizes (도형/카드 내 텍스트)
SHAPE_TITLE_SIZE = 14   # 도형/카드 제목
SHAPE_BODY_SIZE = 13    # 도형/카드 내용

# EMU helpers
def inches(val):
    return int(val * 914400)

SLIDE_W = 12192000  # 13.333"
SLIDE_H = 6858000   # 7.500"
BODY_START_Y = inches(2.0)
BODY_LIMIT_Y = inches(7.0)  # 하단 여백 0.5" 확보


class PptxBuilder:
    """High-level PPTX builder wrapping python-pptx."""

    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.added_slide_indices = []

    # ── Shape Text Replacement ──

    @staticmethod
    def _set_text(shape, text, font_name=None, font_size=None, color=None, bold=None, alignment=None):
        """Replace all text in a shape."""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        if alignment:
            p.alignment = alignment
        run = p.add_run()
        run.text = text
        fmt = run.font
        if font_name: fmt.name = font_name
        if font_size: fmt.size = Pt(font_size)
        if color: fmt.color.rgb = color
        if bold is not None: fmt.bold = bold

    @staticmethod
    def _set_multiline(shape, lines, font_name=None, font_size=None, color=None, bold=None, line_spacing=None):
        """Replace text with multiple lines."""
        tf = shape.text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if line_spacing and i > 0:
                p.space_before = Pt(line_spacing)
            run = p.add_run()
            run.text = line
            fmt = run.font
            if font_name: fmt.name = font_name
            if font_size: fmt.size = Pt(font_size)
            if color: fmt.color.rgb = color
            if bold is not None: fmt.bold = bold

    # ── Format-Preserving Text Replacement ──

    @staticmethod
    def _replace_text_preserve_format(shape, new_text):
        """Replace text preserving original formatting (scheme colors, font).
        Uses first run's XML as template, replaces text only.
        """
        tf = shape.text_frame
        if not tf.paragraphs:
            return
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = new_text
            for run in p.runs[1:]:
                run.text = ""
        from pptx.oxml.ns import qn
        for para in tf.paragraphs[1:]:
            para._p.getparent().remove(para._p)

    @staticmethod
    def _replace_multiline_preserve_format(shape, lines):
        """Replace text with multiple lines, preserving first run's formatting.
        Clones first paragraph XML for each new line.
        """
        import copy as _copy
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            return
        ref_p = tf.paragraphs[0]._p
        txBody = tf._txBody
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        for p_elem in list(txBody.iterchildren(f'{ns}p')):
            txBody.remove(p_elem)
        for line in lines:
            new_p = _copy.deepcopy(ref_p)
            runs = new_p.findall(f'.//{ns}r')
            if runs:
                t_elem = runs[0].find(f'{ns}t')
                if t_elem is not None:
                    t_elem.text = line
                for extra_run in runs[1:]:
                    new_p.remove(extra_run)
            txBody.append(new_p)

    # ── Cover ──

    def set_cover(self, title_lines, subtitle, year="2026", date=""):
        """Modify cover slide (index 0).
        
        Args:
            title_lines: list of strings (1~3 lines, 50pt)
            subtitle: string (24pt)
            year: string
            date: string (MM/DD format)
        
        Rules:
            - 서식 보존: scheme color(흰색) 유지를 위해 XML 복제 방식 사용
            - 제목: 좌우 중앙 (left=1.67", width=10.0"), 높이 2.8" (3줄 수용), 수직 중앙 (top=1.63")
            - 부제: 제목 bottom + 0.15" 간격
            - 회사 소개 shape[3],[4]: 수정하지 않음
        """
        cover = self.prs.slides[0]
        
        # shape[2] = Title: 중앙 배치, 3줄 수용
        title_shape = cover.shapes[2]
        title_shape.height = inches(2.8)
        title_shape.top = inches(1.63)
        self._auto_fit_textbox_width(title_shape, title_lines, font_size_pt=50, center=True)
        self._replace_multiline_preserve_format(title_shape, title_lines)
        
        # shape[5] = Subtitle: 제목 아래 0.15" 간격
        subtitle_shape = cover.shapes[5]
        subtitle_shape.top = title_shape.top + title_shape.height + inches(0.15)
        self._auto_fit_textbox_width(subtitle_shape, [subtitle], font_size_pt=24, center=True)
        self._replace_text_preserve_format(subtitle_shape, subtitle)
        
        # shape[7] = Year
        self._replace_text_preserve_format(cover.shapes[7], year)
        # shape[8] = Date (MM/DD)
        if date:
            self._replace_text_preserve_format(cover.shapes[8], date)

    # ── TOC ──

    def set_toc(self, sections):
        """Modify TOC slide(s) with section list. Auto-paginates at 5 sections per page.
        
        Args:
            sections: list of section name strings
        
        Rules:
            - 서식 보존: replace_multiline_preserve_format 사용
            - 5줄 패딩: 5개 미만이면 빈 문자열로 채움
            - 6개 이상: 첫 번째 목차 spTree를 복제하여 추가 페이지 생성
            - 섹션명 자동 너비: auto_fit_textbox_width 적용
        
        Returns:
            added_toc_count: 추가된 목차 페이지 수 (0이면 단일 페이지)
        """
        import copy as _copy
        
        # 5개씩 페이징
        pages = []
        for i in range(0, len(sections), 5):
            pages.append(sections[i:i+5])
        
        toc_first = self.prs.slides[1]
        
        # 첫 번째 페이지
        page1 = list(pages[0])
        nums1 = [str(i+1) for i in range(len(page1))]
        while len(nums1) < 5:
            nums1.append("")
        while len(page1) < 5:
            page1.append("")
        self._replace_multiline_preserve_format(toc_first.shapes[1], nums1)
        self._auto_fit_textbox_width(toc_first.shapes[2], page1, font_size_pt=24)
        self._replace_multiline_preserve_format(toc_first.shapes[2], page1)
        
        # 추가 페이지: spTree 자식 복제 방식
        added_toc_count = 0
        ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        
        for page_idx in range(1, len(pages)):
            page = list(pages[page_idx])
            start_num = page_idx * 5 + 1
            nums = [str(start_num + i) for i in range(len(page))]
            while len(nums) < 5:
                nums.append("")
            while len(page) < 5:
                page.append("")
            
            new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[4])
            added_toc_count += 1
            
            # spTree 자식 복제
            src_spTree = toc_first._element.find(f'{ns_p}cSld').find(f'{ns_p}spTree')
            dst_spTree = new_slide._element.find(f'{ns_p}cSld').find(f'{ns_p}spTree')
            for child in list(dst_spTree):
                if 'grpSpPr' not in child.tag:
                    dst_spTree.remove(child)
            for child in src_spTree:
                if 'grpSpPr' not in child.tag:
                    dst_spTree.append(_copy.deepcopy(child))
            
            self._replace_multiline_preserve_format(new_slide.shapes[1], nums)
            self._auto_fit_textbox_width(new_slide.shapes[2], page, font_size_pt=24)
            self._replace_multiline_preserve_format(new_slide.shapes[2], page)
        
        return added_toc_count

    @staticmethod
    def _estimate_text_width_pt(text, font_size_pt):
        """Estimate text rendering width in pt."""
        width = 0
        for ch in text:
            if ord(ch) > 0x2E80:
                width += font_size_pt
            elif ch == ' ':
                width += font_size_pt * 0.3
            elif ch in '.:/-()':
                width += font_size_pt * 0.4
            else:
                width += font_size_pt * 0.55
        return width

    @staticmethod
    def _auto_fit_textbox_width(shape, lines, font_size_pt, center=False):
        """Auto-fit textbox width to longest text line."""
        non_empty = [l for l in lines if l]
        if not non_empty:
            return
        max_w_pt = max(PptxBuilder._estimate_text_width_pt(l, font_size_pt) for l in non_empty)
        needed = int(max_w_pt * 12700 * 1.1)
        if center:
            max_allowed = SLIDE_W - 2 * inches(0.5)
            needed = min(needed, max_allowed)
            shape.width = needed
            shape.left = (SLIDE_W - needed) // 2
        else:
            max_allowed = SLIDE_W - shape.left - inches(0.5)
            needed = min(needed, max_allowed)
            shape.width = needed

    # ── Ending ──

    def set_ending(self, tagline="", slide_index=None):
        """Modify ending slide tagline with format preservation.
        
        Uses replace_text_preserve_format to keep scheme color (BACKGROUND_1),
        font (프리젠테이션 4 Regular), and size (16pt).
        """
        if slide_index is None:
            for i, slide in enumerate(self.prs.slides):
                if slide.slide_layout.name == "끝맺음":
                    slide_index = i
        if slide_index is not None and tagline:
            self._replace_text_preserve_format(
                self.prs.slides[slide_index].shapes[2], tagline)

    # ── Body Slides ──

    def add_body_slide(self, title_text, desc_text):
        """Add a new body slide with standard header. Returns the slide object."""
        layout = self.prs.slide_layouts[1]  # 본문
        slide = self.prs.slides.add_slide(layout)
        # Title
        self._add_textbox(slide, 354806, 557906, 2720317, 932614,
                         title_text, FONT_TITLE, 28, PRIMARY, True)
        # Description
        self._add_textbox(slide, 3763700, 558800, 7551483, 937918,
                         desc_text, FONT_DESC, 12, GRAY, False)
        return slide

    # ── Shape Helpers ──

    def add_box(self, slide, left, top, width, height, text="", fill_color=BG_BOX, line_color=BORDER,
                font_size=SHAPE_BODY_SIZE, text_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
        """Add a rounded rectangle with optional text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            inches(left), inches(top), inches(width), inches(height))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        if text:
            self._add_textbox(slide, inches(left), inches(top), inches(width), inches(height),
                             text, FONT_BODY, font_size, text_color, bold, alignment)
        return shape

    def add_badge(self, slide, left, top, size, text, fill_color=PRIMARY, text_color=WHITE):
        """Add a circular badge with text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, inches(left), inches(top), inches(size), inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        self._add_textbox(slide, inches(left), inches(top), inches(size), inches(size),
                         text, FONT_BODY, 14, text_color, True, PP_ALIGN.CENTER)
        return shape

    def add_arrow(self, slide, left, top, width, height, fill_color=PRIMARY):
        """Add a right arrow shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, inches(left), inches(top), inches(width), inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_text(self, slide, left, top, width, height, text,
                 font_name=FONT_BODY, font_size=SHAPE_BODY_SIZE, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
        """Add a text box (inches)."""
        return self._add_textbox(slide, inches(left), inches(top), inches(width), inches(height),
                                text, font_name, font_size, color, bold, alignment)

    def add_multiline_text(self, slide, left, top, width, height, lines,
                           font_name=FONT_BODY, font_size=SHAPE_BODY_SIZE, color=DARK_GRAY, bold=False, line_spacing=2):
        """Add a text box with multiple lines (inches)."""
        txBox = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if line_spacing and i > 0:
                p.space_before = Pt(line_spacing)
            run = p.add_run()
            run.text = line
            fmt = run.font
            fmt.name = font_name
            fmt.size = Pt(font_size)
            fmt.color.rgb = color
            fmt.bold = bold
        return txBox

    # ── Cleanup & Save ──

    def cleanup_template(self, keep_indices=None):
        """Delete template slides except specified indices, then move ending to last.
        
        keep_indices: list of original template slide indices to keep (e.g., [0, 1, 40])
        Added body slides are automatically kept.
        """
        if keep_indices is None:
            keep_indices = [0, 1, 40]  # cover, toc, ending

        # Delete template slides (index 2~39) in reverse
        xml_slides = self.prs.slides._sldIdLst
        delete_indices = [i for i in range(2, 40) if i not in keep_indices]
        for idx in sorted(delete_indices, reverse=True):
            if idx < len(xml_slides):
                xml_slides.remove(xml_slides[idx])

        # Move ending to last
        for i, slide in enumerate(self.prs.slides):
            if slide.slide_layout.name == "끝맺음":
                elem = xml_slides[i]
                xml_slides.remove(elem)
                xml_slides.append(elem)
                break

    def rebuild_sections(self, added_toc_count=0):
        """Rebuild presentation sections to match output structure.
        
        Args:
            added_toc_count: number of additional TOC pages (from set_toc return value)
        """
        from lxml import etree
        import uuid
        prs_xml = self.prs.part._element
        ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        ns_p14 = '{http://schemas.microsoft.com/office/powerpoint/2010/main}'

        sldIdLst = prs_xml.find(f'{ns_p}sldIdLst')
        slide_ids = [elem.get('id') for elem in sldIdLst]

        toc_page_count = 1 + added_toc_count
        body_start = 1 + toc_page_count
        body_end = len(slide_ids) - 1

        for extLst in prs_xml.findall(f'.//{ns_p}extLst'):
            for ext in list(extLst):
                old = ext.find(f'{ns_p14}sectionLst')
                if old is not None:
                    ext.remove(old)
                    new_sectionLst = etree.SubElement(ext, f'{ns_p14}sectionLst')
                    sections_def = [
                        ("표지", [slide_ids[0]]),
                        ("목차", slide_ids[1:1+toc_page_count]),
                        ("본문", slide_ids[body_start:body_end]),
                        ("끝맺음", [slide_ids[-1]]),
                    ]
                    for sec_name, sec_ids in sections_def:
                        sec_elem = etree.SubElement(new_sectionLst, f'{ns_p14}section')
                        sec_elem.set('name', sec_name)
                        sec_elem.set('id', '{' + str(uuid.uuid4()).upper() + '}')
                        sldIdLst_elem = etree.SubElement(sec_elem, f'{ns_p14}sldIdLst')
                        for sid in sec_ids:
                            sldId_elem = etree.SubElement(sldIdLst_elem, f'{ns_p14}sldId')
                            sldId_elem.set('id', sid)
                    break

    def save(self, output_path):
        """Save and verify."""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        # Verify
        verify = Presentation(output_path)
        count = len(verify.slides)
        print(f"✅ Saved: {output_path} ({count} slides)")
        for i, slide in enumerate(verify.slides):
            layout = slide.slide_layout.name
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text[:50].replace("\n", " | ")
                    if t.strip():
                        texts.append(t)
            print(f"  [{i}] {layout}: {'; '.join(texts[:2])}")
        return count

    # ── Internal ──

    @staticmethod
    def _add_textbox(slide, left, top, width, height, text, font_name, font_size, color, bold, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        fmt = run.font
        fmt.name = font_name
        fmt.size = Pt(font_size)
        fmt.color.rgb = color
        fmt.bold = bold
        return txBox
