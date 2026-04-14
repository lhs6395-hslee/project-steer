#!/usr/bin/env python3
"""
Comprehensive PPTX Review Script with Theme Color Resolution
Addresses RGB extraction failure on scheme colors by parsing theme XML
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
import json
from typing import Dict, List, Tuple, Optional

# EMU constants
EMU_PER_INCH = 914400

# Expected dimensions from layout spec
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BODY_START_Y = Inches(2.0)
BODY_LIMIT_Y = Inches(7.0)
MIN_BOTTOM_MARGIN = Inches(0.3)

# Color constants (RGB)
PRIMARY = (68, 114, 196)  # #4472C4
DARK_GRAY = (89, 89, 89)  # #595959
WHITE = (255, 255, 255)
LIGHT_BG_1 = (248, 249, 250)  # #F8F9FA
LIGHT_BG_2 = (255, 255, 255)  # #FFFFFF

def parse_theme_colors(prs: Presentation) -> Dict[str, Tuple[int, int, int]]:
    """
    Extract theme color mappings from theme1.xml in PPTX package
    Returns dict mapping scheme color names to RGB tuples
    """
    theme_colors = {}
    
    try:
        # Access the package part for theme
        theme_part = prs.part.package.part_related_by('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
        theme_xml = theme_part.blob
        
        # Parse theme XML
        root = ET.fromstring(theme_xml)
        
        # Namespace for theme elements
        namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # Extract color scheme
        color_scheme = root.find('.//a:clrScheme', namespaces)
        
        if color_scheme is not None:
            # Map scheme color names
            for color_elem in color_scheme:
                tag_name = color_elem.tag.split('}')[-1]  # Remove namespace
                
                # Find srgbClr or sysClr child
                srgb_elem = color_elem.find('a:srgbClr', namespaces)
                sys_elem = color_elem.find('a:sysClr', namespaces)
                
                if srgb_elem is not None and 'val' in srgb_elem.attrib:
                    hex_val = srgb_elem.attrib['val']
                    r = int(hex_val[0:2], 16)
                    g = int(hex_val[2:4], 16)
                    b = int(hex_val[4:6], 16)
                    theme_colors[tag_name] = (r, g, b)
                elif sys_elem is not None and 'lastClr' in sys_elem.attrib:
                    hex_val = sys_elem.attrib['lastClr']
                    r = int(hex_val[0:2], 16)
                    g = int(hex_val[2:4], 16)
                    b = int(hex_val[4:6], 16)
                    theme_colors[tag_name] = (r, g, b)
        
        print(f"[INFO] Extracted {len(theme_colors)} theme colors from theme1.xml")
        
    except Exception as e:
        print(f"[WARNING] Failed to parse theme colors: {e}")
        # Fallback defaults
        theme_colors = {
            'dk1': (0, 0, 0),
            'lt1': (255, 255, 255),
            'dk2': (68, 114, 196),
            'lt2': (237, 125, 49),
            'accent1': (68, 114, 196),
            'accent2': (237, 125, 49),
            'accent3': (165, 165, 165),
            'accent4': (255, 192, 0),
            'accent5': (91, 155, 213),
            'accent6': (112, 173, 71),
            'hlink': (5, 99, 193),
            'folHlink': (149, 79, 114)
        }
    
    return theme_colors

def get_resolved_rgb(shape, color_attr: str, theme_colors: Dict) -> Optional[Tuple[int, int, int]]:
    """
    Resolve a color attribute (fill or font) to RGB, handling scheme colors
    
    Args:
        shape: Shape object
        color_attr: 'fill' or 'font'
        theme_colors: Theme color mapping from parse_theme_colors()
    
    Returns:
        RGB tuple or None if cannot resolve
    """
    try:
        if color_attr == 'fill':
            if not shape.fill or shape.fill.type != 1:  # 1 = SOLID
                return None
            color_obj = shape.fill.fore_color
        elif color_attr == 'font':
            if not hasattr(shape, 'text_frame') or not shape.text_frame.text:
                return None
            if not shape.text_frame.paragraphs:
                return None
            para = shape.text_frame.paragraphs[0]
            if not para.runs:
                return None
            color_obj = para.runs[0].font.color
        else:
            return None
        
        # Try direct RGB first
        if hasattr(color_obj, 'rgb'):
            try:
                rgb = color_obj.rgb
                return (rgb[0], rgb[1], rgb[2])
            except:
                pass
        
        # Try theme color resolution
        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
            theme_idx = color_obj.theme_color
            
            # Map MSO_THEME_COLOR enum to scheme names
            theme_map = {
                0: 'dk1',      # DARK_1
                1: 'lt1',      # LIGHT_1
                2: 'dk2',      # DARK_2
                3: 'lt2',      # LIGHT_2
                4: 'accent1',  # ACCENT_1
                5: 'accent2',  # ACCENT_2
                6: 'accent3',  # ACCENT_3
                7: 'accent4',  # ACCENT_4
                8: 'accent5',  # ACCENT_5
                9: 'accent6',  # ACCENT_6
                10: 'hlink',   # HYPERLINK
                11: 'folHlink' # FOLLOWED_HYPERLINK
            }
            
            scheme_name = theme_map.get(theme_idx)
            if scheme_name and scheme_name in theme_colors:
                return theme_colors[scheme_name]
        
        return None
        
    except Exception as e:
        return None

def calculate_luminance(rgb: Tuple[int, int, int]) -> float:
    """Calculate relative luminance using WCAG 2.1 formula"""
    def adjust(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    
    r, g, b = rgb
    return 0.2126 * adjust(r) + 0.7152 * adjust(g) + 0.0722 * adjust(b)

def calculate_contrast_ratio(fg_rgb: Tuple[int, int, int], bg_rgb: Tuple[int, int, int]) -> float:
    """Calculate contrast ratio between foreground and background"""
    l1 = calculate_luminance(fg_rgb)
    l2 = calculate_luminance(bg_rgb)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def is_light_background(rgb: Tuple[int, int, int]) -> bool:
    """Check if background is light (luminance > 0.9)"""
    return calculate_luminance(rgb) > 0.9

def extract_subtitle_baseline(template_path: str) -> Dict:
    """Extract subtitle shape dimensions from template for comparison"""
    baseline = {}
    
    try:
        template_prs = Presentation(template_path)
        
        for layout_idx, layout in enumerate(template_prs.slide_layouts):
            layout_name = layout.name
            
            # Find subtitle shapes in each layout
            for shape_idx, shape in enumerate(layout.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    text = shape.text_frame.text.strip().lower()
                    
                    # Identify subtitle by position or placeholder type
                    if shape.top < Inches(2.5) and shape.top > Inches(1.0):
                        if shape_idx not in [0]:  # Not title
                            baseline[f"layout_{layout_idx}_{shape_idx}"] = {
                                'layout_name': layout_name,
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'text_sample': text[:50]
                            }
        
        print(f"[INFO] Extracted {len(baseline)} subtitle baselines from template")
        
    except Exception as e:
        print(f"[WARNING] Failed to extract template baseline: {e}")
    
    return baseline

def review_pptx(target_path: str, template_path: Optional[str] = None) -> Dict:
    """
    Comprehensive PPTX review with theme color resolution
    """
    print(f"[STEP 1/13] Opening PPTX: {target_path}")
    
    try:
        prs = Presentation(target_path)
    except Exception as e:
        return {
            'error': f'Failed to open PPTX: {e}',
            'status': 'failed'
        }
    
    print(f"[STEP 1/13] ✓ Opened: {len(prs.slides)} slides, {prs.slide_width/EMU_PER_INCH:.3f}\" × {prs.slide_height/EMU_PER_INCH:.3f}\"")
    
    # Parse theme colors for scheme color resolution
    print(f"[STEP 2/13] Extracting theme colors from theme1.xml...")
    theme_colors = parse_theme_colors(prs)
    print(f"[STEP 2/13] ✓ Theme colors: {json.dumps(theme_colors, indent=2)}")
    
    # Extract template baseline if provided
    subtitle_baseline = {}
    if template_path and Path(template_path).exists():
        print(f"[STEP 3/13] Extracting subtitle baseline from template...")
        subtitle_baseline = extract_subtitle_baseline(template_path)
        print(f"[STEP 3/13] ✓ Subtitle baseline extracted: {len(subtitle_baseline)} shapes")
    else:
        print(f"[STEP 3/13] ⚠ Template baseline not available")
    
    # Initialize results
    results = {
        'file': target_path,
        'slide_count': len(prs.slides),
        'dimensions': f"{prs.slide_width/EMU_PER_INCH:.3f}\" × {prs.slide_height/EMU_PER_INCH:.3f}\"",
        'theme_colors': theme_colors,
        'issues': [],
        'slides': []
    }
    
    # Step 4: Analyze each slide
    print(f"[STEP 4/13] Analyzing slide structures...")
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            'index': slide_idx,
            'layout_name': slide.slide_layout.name,
            'shape_count': len(slide.shapes),
            'shapes': [],
            'violations': []
        }
        
        # Track subtitle violations
        subtitle_violations = []
        
        # Track white-on-light violations
        white_on_light_violations = []
        
        # Track center alignment
        content_shapes = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_data = {
                'index': shape_idx,
                'name': shape.name,
                'type': str(shape.shape_type),
                'left_in': shape.left / EMU_PER_INCH,
                'top_in': shape.top / EMU_PER_INCH,
                'width_in': shape.width / EMU_PER_INCH,
                'height_in': shape.height / EMU_PER_INCH
            }
            
            # Check if subtitle (top < 2.5", top > 1.0", not first shape)
            is_subtitle = (shape.top < Inches(2.5) and shape.top > Inches(1.0) and shape_idx > 0)
            
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                text = shape.text_frame.text
                shape_data['text'] = text[:100]
                shape_data['text_length'] = len(text)
                
                # Count lines
                line_count = text.count('\n') + 1
                shape_data['line_count'] = line_count
                
                # Check subtitle violations
                if is_subtitle:
                    if line_count > 2:
                        subtitle_violations.append(f"Slide {slide_idx}, Shape {shape_idx}: {line_count} lines (max 2)")
                    
                    # Check word breaks
                    lines = text.split('\n')
                    for line_no, line in enumerate(lines):
                        if line and not line[0].isspace() and line_no > 0:
                            prev_line = lines[line_no - 1]
                            if prev_line and not prev_line[-1].isspace():
                                # Check if word is broken
                                if prev_line[-3:].isalpha() and line[:3].isalpha():
                                    subtitle_violations.append(f"Slide {slide_idx}, Shape {shape_idx}: Word break at line {line_no}")
                
                # Extract font properties
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            shape_data['font_size_pt'] = run.font.size.pt
                        if run.font.name:
                            shape_data['font_name'] = run.font.name
                        if run.font.bold is not None:
                            shape_data['font_bold'] = run.font.bold
                        
                        # Get font color using theme resolution
                        font_rgb = get_resolved_rgb(shape, 'font', theme_colors)
                        if font_rgb:
                            shape_data['font_color_rgb'] = font_rgb
                            
                            # Check for white text
                            if font_rgb == WHITE or calculate_luminance(font_rgb) > 0.95:
                                # Get background color
                                bg_rgb = get_resolved_rgb(shape, 'fill', theme_colors)
                                if bg_rgb and is_light_background(bg_rgb):
                                    contrast = calculate_contrast_ratio(font_rgb, bg_rgb)
                                    white_on_light_violations.append({
                                        'slide': slide_idx,
                                        'shape': shape_idx,
                                        'font_rgb': font_rgb,
                                        'bg_rgb': bg_rgb,
                                        'contrast': f"{contrast:.2f}:1",
                                        'severity': 'CRITICAL'
                                    })
            
            # Get fill color using theme resolution
            fill_rgb = get_resolved_rgb(shape, 'fill', theme_colors)
            if fill_rgb:
                shape_data['fill_color_rgb'] = fill_rgb
            
            # Check for Oval placeholders
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if hasattr(shape, 'auto_shape_type'):
                        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                        if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                            shape_data['is_oval_placeholder'] = True
                            slide_info['violations'].append(f"CRITICAL: Oval placeholder at shape {shape_idx}")
                except:
                    pass
            
            # Track body content for center alignment check
            if shape.top >= BODY_START_Y and shape.top < BODY_LIMIT_Y:
                content_shapes.append(shape)
            
            slide_info['shapes'].append(shape_data)
        
        # Check center alignment of body content
        if content_shapes:
            leftmost = min(s.left for s in content_shapes)
            rightmost = max(s.left + s.width for s in content_shapes)
            content_width = rightmost - leftmost
            
            center_pos = leftmost + content_width / 2
            slide_center = prs.slide_width / 2
            center_offset_in = abs(center_pos - slide_center) / EMU_PER_INCH
            
            slide_info['center_offset_in'] = center_offset_in
            
            if center_offset_in > 0.1:
                slide_info['violations'].append(f"MAJOR: Center offset {center_offset_in:.2f}\" (tolerance ±0.1\")")
        
        # Add subtitle violations
        if subtitle_violations:
            slide_info['subtitle_violations'] = subtitle_violations
        
        # Add white-on-light violations
        if white_on_light_violations:
            slide_info['white_on_light_violations'] = white_on_light_violations
            results['issues'].extend([
                f"CRITICAL: Slide {v['slide']}, Shape {v['shape']}: White text (RGB{v['font_rgb']}) on light background (RGB{v['bg_rgb']}), contrast {v['contrast']}"
                for v in white_on_light_violations
            ])
        
        results['slides'].append(slide_info)
    
    print(f"[STEP 4/13] ✓ Analyzed {len(prs.slides)} slides")
    
    # Compile all violations
    print(f"[STEP 5-12/13] Compiling violations...")
    
    critical_count = 0
    major_count = 0
    minor_count = 0
    
    for slide_info in results['slides']:
        for violation in slide_info.get('violations', []):
            if 'CRITICAL' in violation:
                critical_count += 1
            elif 'MAJOR' in violation:
                major_count += 1
            else:
                minor_count += 1
    
    # Add white-on-light critical violations
    critical_count += sum(len(s.get('white_on_light_violations', [])) for s in results['slides'])
    
    results['violation_summary'] = {
        'critical': critical_count,
        'major': major_count,
        'minor': minor_count
    }
    
    # Calculate quality score
    if critical_count > 0:
        quality_score = 30
    elif major_count > 0:
        quality_score = max(30, 70 - (major_count * 10))
    elif minor_count > 0:
        quality_score = max(70, 90 - (minor_count * 5))
    else:
        quality_score = 100
    
    results['quality_score'] = quality_score
    
    print(f"[STEP 13/13] Quality Score: {quality_score}/100")
    print(f"[STEP 13/13] Violations: {critical_count} CRITICAL, {major_count} MAJOR, {minor_count} MINOR")
    
    return results

if __name__ == '__main__':
    target_path = 'results/pptx/AWS_MSK_Expert_Intro.pptx'
    template_path = 'templates/pptx_template.pptx'
    
    results = review_pptx(target_path, template_path)
    
    # Write results to file
    output_path = 'results/pptx_review_comprehensive.json'
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    
    print(f"\n[COMPLETE] Results written to: {output_path}")
    print(f"Quality Score: {results.get('quality_score', 0)}/100")

