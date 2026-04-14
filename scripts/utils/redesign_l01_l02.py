from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import shutil

INPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
TEMP  = "results/pptx/AWS_MSK_Expert_Intro_tmp.pptx"

PRIMARY    = RGBColor(0, 67, 218)
SUB_ORANGE = RGBColor(238, 129, 80)
SUB_GREEN  = RGBColor(76, 184, 143)
DARK_GRAY  = RGBColor(33, 33, 33)
GRAY_MID   = RGBColor(150, 150, 150)
WHITE      = RGBColor(255, 255, 255)
F8F9FA     = RGBColor(248, 249, 250)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def set_anchor(shape, anchor):
    bodyPr = shape._element.find(f".//{{{A_NS}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)

def add_rr(slide, left, top, w, h, fill, line_color=None, line_pt=1.0, no_line=False):
    s = slide.shapes.add_shape(5, Emu(left), Emu(top), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if no_line:
        s.line.fill.background()
    elif line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    return s

def add_arrow(slide, left, top, w, h, fill):
    s = slide.shapes.add_shape(13, Emu(left), Emu(top), Emu(w), Emu(h))  # RIGHT_ARROW=13
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def add_tb(slide, left, top, w, h, text, font_pt, bold, color,
           align=PP_ALIGN.LEFT, anchor="t", font_name="Freesentation"):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    set_anchor(tb, anchor)
    bodyPr = tb._element.find(f".//{{{A_NS}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("lIns", "91440"); bodyPr.set("rIns", "91440")
        bodyPr.set("tIns", "45720"); bodyPr.set("bIns", "45720")
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def replace_text_keep_format(shape, new_text):
    """Replace text in existing shape while preserving run format."""
    tf = shape.text_frame
    txBody = tf._txBody
    # Clear existing paragraphs
    for p in txBody.findall(f"{{{A_NS}}}p"):
        txBody.remove(p)
    lines = new_text.split("\n")
    for line in lines:
        p_elem = etree.SubElement(txBody, f"{{{A_NS}}}p")
        r_elem = etree.SubElement(p_elem, f"{{{A_NS}}}r")
        rPr = etree.SubElement(r_elem, f"{{{A_NS}}}rPr")
        rPr.set("lang", "ko-KR")
        t_elem = etree.SubElement(r_elem, f"{{{A_NS}}}t")
        t_elem.text = line

prs = Presentation(INPUT)

# ═══════════════════════════════════════════════
# SLIDE idx 6 (L02 Three Cards) — REVERT + color
# ═══════════════════════════════════════════════
slide6 = prs.slides[6]
print("=== idx6: remove wrong shapes ===")
remove_names = {
    "Card_Left_BG","Card_Left_Title","Card_Left_Content",
    "Rectangle 18","VPC_Label","Rectangle 20","TextBox 21",
    "Rounded Rectangle 22","Broker_AZa_Label","Rounded Rectangle 24","TextBox 25",
    "Rounded Rectangle 26","TextBox 27","TextBox 28","TextBox 29",
}
for shape in list(slide6.shapes):
    if shape.name in remove_names:
        shape._element.getparent().remove(shape._element)
        print(f"  Removed: {shape.name}")

print("=== idx6: add 3 cards ===")
CARDS = [
    (457200,  "핵심 구성요소",
     "• Broker: 메시지 저장/전달 노드\n• Topic: 메시지 카테고리 단위\n• Partition: 병렬 처리 단위\n• Replication Factor: 데이터 복제 수\n• ZooKeeper / KRaft: 메타데이터 관리"),
    (4292600, "네트워크 구성",
     "• VPC 내 Private Subnet 배치\n• Security Group 기반 접근 제어\n• ENI: 브로커별 네트워크 인터페이스\n• PrivateLink: Cross-VPC 연결\n• Multi-AZ 자동 분산 배치"),
    (8128000, "스토리지",
     "• EBS gp3: 기본 브로커 스토리지\n• Provisioned IOPS: 고성능 워크로드\n• Tiered Storage: S3 기반 장기 보관\n• 자동 확장: 디스크 용량 자동 증설\n• 보존 기간별 계층화 전략"),
]
for left, title, content in CARDS:
    add_rr(slide6, left, 4023360, 3606800, 2286000, F8F9FA, no_line=True)
    add_tb(slide6, left+91440, 4160520, 3424920, 304800, title, 14, True, PRIMARY)
    add_tb(slide6, left+91440, 4526280, 3424920, 1645920, content, 13, False, DARK_GRAY)
print("  3 cards added")

print("=== idx6: diversify flow colors ===")
COLOR_MAP = {
    "Rounded Rectangle 6":  (SUB_ORANGE, WHITE),
    "Right Arrow 7":        (SUB_ORANGE, None),
    "Right Arrow 13":       (SUB_GREEN,  None),
    "Rounded Rectangle 14": (SUB_GREEN,  WHITE),
}
for shape in slide6.shapes:
    if shape.name in COLOR_MAP:
        fill_c, text_c = COLOR_MAP[shape.name]
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_c
        if text_c and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = text_c
        print(f"  Recolored: {shape.name} → {fill_c}")

# ═══════════════════════════════════════════════
# SLIDE idx 5 (L01 Bento Grid) — left card + flow diagram
# ═══════════════════════════════════════════════
slide5 = prs.slides[5]
print("\n=== idx5: update left card content ===")

# Update TextBox 5 (card title)
for shape in slide5.shapes:
    if shape.name == "TextBox 5" and shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = "MSK 아키텍처 핵심 개요"
                break
            break
        print(f"  TextBox 5 updated")

# Update TextBox 6 (card content) — preserve existing format by replacing runs
for shape in slide5.shapes:
    if shape.name == "TextBox 6" and shape.has_text_frame:
        tf = shape.text_frame
        txBody = tf._txBody
        # Remove all existing paragraphs
        for p in list(txBody.findall(f"{{{A_NS}}}p")):
            txBody.remove(p)
        # Re-add with new content
        new_content = [
            ("클러스터 유형", True),
            ("• Provisioned: 직접 용량 지정, 고성능", False),
            ("• Serverless: 자동 확장, 사용량 과금", False),
            ("• Express: 단일 AZ, 개발/테스트용", False),
            ("", False),
            ("보안 체계", True),
            ("• 인증: IAM / SASL-SCRAM / mTLS", False),
            ("• 암호화: TLS in-transit + EBS at-rest", False),
            ("• 접근: Security Group + VPC Endpoint", False),
        ]
        for line, is_bold in new_content:
            p_el = etree.SubElement(txBody, f"{{{A_NS}}}p")
            r_el = etree.SubElement(p_el, f"{{{A_NS}}}r")
            rPr = etree.SubElement(r_el, f"{{{A_NS}}}rPr")
            rPr.set("lang", "ko-KR")
            sz = "1300"
            rPr.set("sz", sz)
            rPr.set("b", "1" if is_bold else "0")
            fill_el = etree.SubElement(rPr, f"{{{A_NS}}}solidFill")
            clr_el = etree.SubElement(fill_el, f"{{{A_NS}}}srgbClr")
            clr_el.set("val", "212121")
            lat_el = etree.SubElement(rPr, f"{{{A_NS}}}latin")
            lat_el.set("typeface", "Freesentation")
            t_el = etree.SubElement(r_el, f"{{{A_NS}}}t")
            t_el.text = line
        print(f"  TextBox 6 updated")

print("=== idx5: add flow diagram (resize RR15 to fit content) ===")
# Remove existing RR15 and rebuild with proper size for the flow diagram
for shape in list(slide5.shapes):
    if shape.name == "Rounded Rectangle 15":
        shape._element.getparent().remove(shape._element)
        print(f"  Removed old RR15")

# Flow diagram layout — sized to content
# Producer → MSK Cluster → Consumer with labels
# Boxes: w=1097280 (1.20"), h=457200 (0.50")
# Arrows: w=365760 (0.40"), h=228600 (0.25")
# Total width: 1.20 + 0.40 + 1.20 + 0.40 + 1.20 = 4.40" = 4023360 EMU
# Plus padding 0.30" each side → container w = 5.00" = 4572000 EMU
# Container h: label(0.25") + box(0.50") + sublabel(0.25") + padding = 1.30" = 1188720 EMU

FLOW_LEFT   = 823320    # 0.90" (same left as original RR15)
FLOW_TOP    = 4796280   # position naturally below left card content
FLOW_W      = 4846080   # 5.30" — fits all 3 boxes + arrows + padding
FLOW_H      = 1371600   # 1.50"
PADDING     = 274320    # 0.30" internal padding

# Container rounded rect
add_rr(slide5, FLOW_LEFT, FLOW_TOP, FLOW_W, FLOW_H, F8F9FA, PRIMARY, 1.0)

# Calculate inner positions
inner_left = FLOW_LEFT + PADDING
inner_top  = FLOW_TOP + PADDING
BOX_W = 1097280   # 1.20"
BOX_H = 457200    # 0.50"
ARR_W = 274320    # 0.30"
ARR_H = 228600    # 0.25"

positions = [
    ("Producer",  SUB_ORANGE, inner_left),
    ("MSK",       PRIMARY,    inner_left + BOX_W + ARR_W),
    ("Consumer",  SUB_GREEN,  inner_left + BOX_W + ARR_W + BOX_W + ARR_W),
]
arrow_positions = [
    inner_left + BOX_W,
    inner_left + BOX_W + ARR_W + BOX_W,
]

for label, color, left in positions:
    add_rr(slide5, left, inner_top, BOX_W, BOX_H, color, no_line=True)
    add_tb(slide5, left, inner_top, BOX_W, BOX_H,
           label, 11, True, WHITE, PP_ALIGN.CENTER, "ctr")

for arr_left in arrow_positions:
    add_arrow(slide5, arr_left, inner_top + (BOX_H - ARR_H)//2, ARR_W, ARR_H, GRAY_MID)

# Sub-labels
sub_top = inner_top + BOX_H + 91440
sub_labels = [
    ("App / SDK",   inner_left,                    SUB_ORANGE),
    ("Topic/Part.", inner_left + BOX_W + ARR_W,    PRIMARY),
    ("Consumer Grp",inner_left + 2*(BOX_W+ARR_W),  SUB_GREEN),
]
for sub_text, sub_left, sub_color in sub_labels:
    add_tb(slide5, sub_left, sub_top, BOX_W, 228600,
           sub_text, 9, False, sub_color, PP_ALIGN.CENTER, "t")

print(f"  Flow diagram added: container top={FLOW_TOP/914400:.2f}\" bot={(FLOW_TOP+FLOW_H)/914400:.2f}\"")

prs.save(TEMP)
shutil.move(TEMP, INPUT)
print("\nSaved.")

# Verify
prs2 = Presentation(INPUT)
print("\n[verify idx6]")
for s in prs2.slides[6].shapes:
    if s.name in {"Text Placeholder 1","TextBox 2","TextBox 3"}: continue
    bot = (s.top+s.height)/914400
    print(f"  {s.name}: top={s.top/914400:.2f}\" bot={bot:.2f}\" {'⚠' if bot>7.01 else ''}")

print("\n[verify idx5]")
for s in prs2.slides[5].shapes:
    if s.name in {"Text Placeholder 1","TextBox 2","TextBox 3"}: continue
    bot = (s.top+s.height)/914400
    print(f"  {s.name}: top={s.top/914400:.2f}\" bot={bot:.2f}\" {'⚠' if bot>7.01 else ''}")
