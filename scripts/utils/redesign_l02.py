from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import shutil

INPUT = "results/pptx/AWS_MSK_Expert_Intro.pptx"
TEMP  = "results/pptx/AWS_MSK_Expert_Intro_tmp.pptx"

PRIMARY   = RGBColor(0, 67, 218)
DARK_GRAY = RGBColor(33, 33, 33)
WHITE     = RGBColor(255, 255, 255)
VPC_FILL  = RGBColor(235, 244, 255)
SUB_FILL  = RGBColor(214, 234, 255)
F8F9FA    = RGBColor(248, 249, 250)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def set_anchor(shape, anchor):
    bodyPr = shape._element.find(f".//{{{A_NS}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)

def add_rr(slide, left, top, w, h, fill, line_color, line_pt=1.0):
    s = slide.shapes.add_shape(5, Emu(left), Emu(top), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    return s

def add_rect(slide, left, top, w, h, fill, line_color, line_pt=1.0):
    s = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    return s

def add_tb(slide, left, top, w, h, text, font_name, font_pt, bold, color,
           align=PP_ALIGN.LEFT, anchor="t"):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    set_anchor(tb, anchor)
    bodyPr = tb._element.find(f".//{{{A_NS}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("lIns", "91440"); bodyPr.set("rIns", "91440")
        bodyPr.set("tIns", "45720"); bodyPr.set("bIns", "45720")

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

prs = Presentation(INPUT)
slide = prs.slides[6]

# 1. Remove old card shapes
remove_names = {
    "Rounded Rectangle 15","TextBox 16","TextBox 17",
    "Rounded Rectangle 18","TextBox 19","TextBox 20",
    "Rounded Rectangle 21","TextBox 22","TextBox 23",
}
for shape in list(slide.shapes):
    if shape.name in remove_names:
        shape._element.getparent().remove(shape._element)
        print(f"  Removed: {shape.name}")

# 2. Left card (wider)
add_rr(slide, 457200, 4023360, 4572000, 2286000, F8F9FA, PRIMARY, 1.5)
add_tb(slide, 548640, 4080960, 4389120, 304800,
       "MSK 인프라 핵심 구성", "Freesentation", 14, True, PRIMARY)
content = ("• 네트워크: VPC Private Subnet, Security Group, Multi-AZ 분산\n"
           "• 스토리지: EBS gp3, Tiered Storage (S3), 자동 확장\n"
           "• 브로커: 메시지 저장/전달, Topic/Partition 단위 분리\n"
           "• 복제: Replication Factor, KRaft 메타데이터 관리")
add_tb(slide, 548640, 4432560, 4389120, 1828800,
       content, "Freesentation", 11, False, DARK_GRAY)
print("  Left card added")

# 3. VPC outer box
add_rect(slide, 5486400, 4023360, 6248160, 2286000, VPC_FILL, PRIMARY, 1.5)
add_tb(slide, 5577840, 4054560, 914400, 228600,
       "VPC", "Freesentation", 10, True, PRIMARY)
print("  VPC outer box added")

# 4. Private Subnet inner box (dashed border via XML)
subnet = add_rect(slide, 5669280, 4298160, 4800600, 1828800, SUB_FILL, PRIMARY, 1.0)
# dashed line
ln = subnet._element.find(f".//{{{A_NS}}}ln")
if ln is not None:
    prstDash = etree.SubElement(ln, f"{{{A_NS}}}prstDash")
    prstDash.set("val", "dash")
add_tb(slide, 5760720, 4329360, 1828800, 228600,
       "Private Subnet", "Freesentation", 9, False, PRIMARY)
print("  Subnet box added")

# 5. Three broker boxes
BROKER_TOP = 4600800
BROKER_H   = 640080
for label, left in [("Broker\nAZ-a", 5760720), ("Broker\nAZ-b", 7223520), ("Broker\nAZ-c", 8686320)]:
    add_rr(slide, left, BROKER_TOP, 1371600, BROKER_H, PRIMARY, WHITE, 0.75)
    add_tb(slide, left, BROKER_TOP, 1371600, BROKER_H,
           label, "Freesentation", 10, True, WHITE, PP_ALIGN.CENTER, "ctr")
print("  3 broker boxes added")

# 6. EBS + S3 labels
add_tb(slide, 5760720, 5349600, 2743200, 228600,
       "↓ EBS gp3  |  Tiered Storage → S3", "Freesentation", 9, False, DARK_GRAY)
print("  Storage labels added")

# 7. Right-of-VPC: PrivateLink label
add_tb(slide, 10606560, 4389120, 1097280, 914400,
       "Private\nLink\n/\nMulti-AZ", "Freesentation", 9, False, PRIMARY, PP_ALIGN.CENTER, "ctr")
print("  PrivateLink label added")

prs.save(TEMP)
shutil.move(TEMP, INPUT)
print("\nSaved.")

# Verify
prs2 = Presentation(INPUT)
slide2 = prs2.slides[6]
print("\n[verify idx6]")
for s in slide2.shapes:
    bot = (s.top + s.height) / 914400
    flag = " ⚠OVERFLOW" if bot > 7.01 else ""
    if s.name not in {"Text Placeholder 1","TextBox 2","TextBox 3"}:
        print(f"  {s.name}: top={s.top/914400:.2f}\" bot={bot:.2f}\"{flag}")
