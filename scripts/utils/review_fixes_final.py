#!/usr/bin/env python3
"""
Reviewer verification script for AWS_MSK_Expert_Intro.pptx
Checks fixes applied to slide idx 8 (L04) and slide idx 9 (L05).
Shape references are by NAME (e.g., "TextBox 7"), not list index.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

EMU_PER_INCH = 914400

def emu_to_inch(emu):
    return emu / EMU_PER_INCH if emu else 0

def get_fill_type(shape):
    """Returns fill info: 'noFill', 'solid:#RRGGBB', or error string"""
    try:
        fill = shape.fill
        fill_type = fill.type
        if fill_type is None:
            return 'noFill'
        fill_type_str = str(fill_type)
        if 'BACKGROUND' in fill_type_str or fill_type_str in ('5', 'None'):
            return 'noFill'
        elif 'SOLID' in fill_type_str or fill_type_str == '1':
            try:
                rgb = fill.fore_color.rgb
                return f'solid:#{str(rgb).upper()}'
            except:
                return 'solid:unknown'
        else:
            return f'type:{fill_type}'
    except AttributeError:
        # Pictures don't have .fill attribute
        return 'N/A(picture)'
    except Exception as e:
        return f'error:{e}'


def get_anchor(shape):
    """Returns text anchor attribute from bodyPr: 't', 'ctr', 'b', or None"""
    try:
        txBody = shape.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody')
        if txBody is None:
            return None
        bodyPr = txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if bodyPr is None:
            return None
        anchor = bodyPr.get('anchor')
        return anchor  # 't', 'ctr', 'b', or None (None means default=ctr for most)
    except:
        return None


def get_shape_position(shape):
    left = emu_to_inch(shape.left)
    top = emu_to_inch(shape.top)
    width = emu_to_inch(shape.width)
    height = emu_to_inch(shape.height)
    return left, top, width, height


def build_shape_name_map(slide):
    """Build dict: shape_name -> shape"""
    return {shape.name: shape for shape in slide.shapes}


def main():
    prs = Presentation(PPTX_PATH)

    checks = {}
    issues = []
    fix_required = []

    # =========================================================
    # SLIDE IDX 8 (L04)
    # =========================================================
    slide8 = prs.slides[8]
    shapes8_by_name = build_shape_name_map(slide8)

    print(f"Slide idx 8 shapes ({len(slide8.shapes)} total):")
    for i, shape in enumerate(slide8.shapes):
        left, top, width, height = get_shape_position(shape)
        anchor = get_anchor(shape)
        fill = get_fill_type(shape)
        print(f"  [{i}] name='{shape.name}' type={shape.shape_type} pos=({left:.2f}\", {top:.2f}\") size=({width:.2f}\"x{height:.2f}\") anchor={anchor} fill={fill}")

    print()

    # CHECK 1: TextBox 7, 11, 15, 19 anchor == 't' (by name)
    tb_names_s8 = ['TextBox 7', 'TextBox 11', 'TextBox 15', 'TextBox 19']
    anchor_checks_s8 = {}
    for name in tb_names_s8:
        if name in shapes8_by_name:
            shape = shapes8_by_name[name]
            anchor = get_anchor(shape)
            passed = anchor == 't'
            anchor_checks_s8[name] = {
                'anchor': anchor,
                'pass': passed
            }
            if not passed:
                issues.append(f"Slide8 '{name}': anchor='{anchor}' (expected 't')")
                fix_required.append(f"slide8_{name.replace(' ','_')}_anchor")
        else:
            anchor_checks_s8[name] = {'error': 'shape not found', 'pass': False}
            issues.append(f"Slide8: shape '{name}' not found")

    checks['slide8_textbox_anchors'] = anchor_checks_s8
    checks['slide8_textbox_anchors_pass'] = all(v.get('pass', False) for v in anchor_checks_s8.values())

    # CHECK 2: PICTURE shape near bottom-right (~10.5", 6.2"), size ~0.6"×0.6"
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    icon_found = False
    icon_details = []
    for shape in slide8.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = get_shape_position(shape)
            icon_details.append({
                'name': shape.name,
                'left': round(left, 3),
                'top': round(top, 3),
                'width': round(width, 3),
                'height': round(height, 3)
            })
            # Near bottom-right: left ~10.5" (±1"), top ~6.2" (±0.7")
            if 9.5 <= left <= 11.5 and 5.5 <= top <= 6.9:
                icon_found = True

    checks['slide8_icon_picture'] = {
        'pictures_found': icon_details,
        'bottom_right_icon_found': icon_found,
        'pass': icon_found
    }
    if not icon_found:
        issues.append("Slide8: No PICTURE shape found near bottom-right (~10.5\", 6.2\")")
        fix_required.append("slide8_add_icon_bottom_right")

    # =========================================================
    # SLIDE IDX 9 (L05)
    # =========================================================
    slide9 = prs.slides[9]
    shapes9_by_name = build_shape_name_map(slide9)

    print(f"\nSlide idx 9 shapes ({len(slide9.shapes)} total):")
    for i, shape in enumerate(slide9.shapes):
        left, top, width, height = get_shape_position(shape)
        anchor = get_anchor(shape)
        fill = get_fill_type(shape)
        print(f"  [{i}] name='{shape.name}' type={shape.shape_type} pos=({left:.2f}\", {top:.2f}\") size=({width:.2f}\"x{height:.2f}\") anchor={anchor} fill={fill}")

    print()

    # CHECK 3: TextBox 11 fill is noFill (by name)
    tb11_name = 'TextBox 11'
    if tb11_name in shapes9_by_name:
        shape = shapes9_by_name[tb11_name]
        fill = get_fill_type(shape)
        passed = fill == 'noFill'
        checks['slide9_tb11_fill'] = {
            'name': tb11_name,
            'fill': fill,
            'pass': passed
        }
        if not passed:
            issues.append(f"Slide9 '{tb11_name}': fill='{fill}' (expected noFill)")
            fix_required.append("slide9_TextBox11_remove_fill")
    else:
        checks['slide9_tb11_fill'] = {'error': f"'{tb11_name}' not found", 'pass': False}
        issues.append(f"Slide9: shape '{tb11_name}' not found")

    # CHECK 4: TextBox 21 fill is noFill (by name)
    tb21_name = 'TextBox 21'
    if tb21_name in shapes9_by_name:
        shape = shapes9_by_name[tb21_name]
        fill = get_fill_type(shape)
        passed = fill == 'noFill'
        checks['slide9_tb21_fill'] = {
            'name': tb21_name,
            'fill': fill,
            'pass': passed
        }
        if not passed:
            issues.append(f"Slide9 '{tb21_name}': fill='{fill}' (expected noFill)")
            fix_required.append("slide9_TextBox21_remove_fill")
    else:
        checks['slide9_tb21_fill'] = {'error': f"'{tb21_name}' not found on slide9", 'pass': False}
        issues.append(f"Slide9: shape '{tb21_name}' not found (slide has {len(slide9.shapes)} shapes)")

    # CHECK 5: TextBox 7, 11, 15, 19 anchor == 't' (by name)
    tb_names_s9 = ['TextBox 7', 'TextBox 11', 'TextBox 15', 'TextBox 19']
    anchor_checks_s9 = {}
    for name in tb_names_s9:
        if name in shapes9_by_name:
            shape = shapes9_by_name[name]
            anchor = get_anchor(shape)
            passed = anchor == 't'
            anchor_checks_s9[name] = {
                'anchor': anchor,
                'pass': passed
            }
            if not passed:
                issues.append(f"Slide9 '{name}': anchor='{anchor}' (expected 't')")
                fix_required.append(f"slide9_{name.replace(' ','_')}_anchor")
        else:
            anchor_checks_s9[name] = {'error': 'shape not found', 'pass': False}
            issues.append(f"Slide9: shape '{name}' not found")

    checks['slide9_textbox_anchors'] = anchor_checks_s9
    checks['slide9_textbox_anchors_pass'] = all(v.get('pass', False) for v in anchor_checks_s9.values())

    # CHECK 6: Phase header fills unchanged
    # RR4=#001B5E, RR8=#0043DA, RR12=#EE8150, RR16=#4CB88F
    phase_header_expected = {
        'Rounded Rectangle 4': '#001B5E',
        'Rounded Rectangle 8': '#0043DA',
        'Rounded Rectangle 12': '#EE8150',
        'Rounded Rectangle 16': '#4CB88F'
    }
    phase_checks = {}
    for shape_name, expected_hex in phase_header_expected.items():
        if shape_name in shapes9_by_name:
            shape = shapes9_by_name[shape_name]
            fill = get_fill_type(shape)
            expected_fill = f'solid:{expected_hex}'
            passed = fill.upper() == expected_fill.upper() if fill else False
            phase_checks[shape_name] = {
                'fill': fill,
                'expected': expected_fill,
                'pass': passed
            }
            if not passed:
                issues.append(f"Slide9 '{shape_name}': fill='{fill}' (expected '{expected_fill}')")
                fix_required.append(f"slide9_{shape_name.replace(' ','_')}_phase_fill_changed")
        else:
            phase_checks[shape_name] = {'error': 'shape not found', 'pass': False}
            issues.append(f"Slide9: phase header shape '{shape_name}' not found")

    checks['slide9_phase_header_fills'] = phase_checks
    checks['slide9_phase_header_fills_pass'] = all(v.get('pass', False) for v in phase_checks.values())

    # Also check TextBox 5,9,13,17 fills (phase header text overlays)
    phase_text_expected = {
        'TextBox 5': '#001B5E',
        'TextBox 9': '#0043DA',
        'TextBox 13': '#EE8150',
        'TextBox 17': '#4CB88F'
    }
    phase_text_checks = {}
    for shape_name, expected_hex in phase_text_expected.items():
        if shape_name in shapes9_by_name:
            shape = shapes9_by_name[shape_name]
            fill = get_fill_type(shape)
            expected_fill = f'solid:{expected_hex}'
            passed = fill.upper() == expected_fill.upper() if fill else False
            phase_text_checks[shape_name] = {
                'fill': fill,
                'expected': expected_fill,
                'pass': passed
            }
            if not passed:
                issues.append(f"Slide9 '{shape_name}' (phase text): fill='{fill}' (expected '{expected_fill}')")
                fix_required.append(f"slide9_{shape_name.replace(' ','_')}_phase_text_fill_changed")
        else:
            phase_text_checks[shape_name] = {'error': 'shape not found', 'pass': False}
            issues.append(f"Slide9: phase text shape '{shape_name}' not found")

    checks['slide9_phase_textbox_fills'] = phase_text_checks
    checks['slide9_phase_textbox_fills_pass'] = all(v.get('pass', False) for v in phase_text_checks.values())

    # CHECK 7: Content card fills unchanged (RR6,RR10,RR14,RR18 = #F8F9FA)
    content_card_expected = {
        'Rounded Rectangle 6': '#F8F9FA',
        'Rounded Rectangle 10': '#F8F9FA',
        'Rounded Rectangle 14': '#F8F9FA',
        'Rounded Rectangle 18': '#F8F9FA'
    }
    content_card_checks = {}
    for shape_name, expected_hex in content_card_expected.items():
        if shape_name in shapes9_by_name:
            shape = shapes9_by_name[shape_name]
            fill = get_fill_type(shape)
            expected_fill = f'solid:{expected_hex}'
            passed = fill.upper() == expected_fill.upper() if fill else False
            content_card_checks[shape_name] = {
                'fill': fill,
                'expected': expected_fill,
                'pass': passed
            }
            if not passed:
                issues.append(f"Slide9 '{shape_name}': fill='{fill}' (expected '{expected_fill}')")
                fix_required.append(f"slide9_{shape_name.replace(' ','_')}_card_fill_changed")
        else:
            content_card_checks[shape_name] = {'error': 'shape not found', 'pass': False}
            issues.append(f"Slide9: content card shape '{shape_name}' not found")

    checks['slide9_content_card_fills'] = content_card_checks
    checks['slide9_content_card_fills_pass'] = all(v.get('pass', False) for v in content_card_checks.values())

    # =========================================================
    # NOTE: anchor == None vs 't'
    # In PPTX, bodyPr anchor attribute absence means default center.
    # 't' means explicit top. We flag None as FAIL because the Sprint_Contract
    # requires explicit 't'.
    # =========================================================

    # =========================================================
    # VERDICT
    # =========================================================
    critical_checks = [
        ('slide8_textbox_anchors_pass', checks.get('slide8_textbox_anchors_pass', False)),
        ('slide8_icon_picture_pass', checks.get('slide8_icon_picture', {}).get('pass', False)),
        ('slide9_tb11_fill_pass', checks.get('slide9_tb11_fill', {}).get('pass', False)),
        ('slide9_tb21_fill_pass', checks.get('slide9_tb21_fill', {}).get('pass', False)),
        ('slide9_textbox_anchors_pass', checks.get('slide9_textbox_anchors_pass', False)),
        ('slide9_phase_header_fills_pass', checks.get('slide9_phase_header_fills_pass', False)),
        ('slide9_content_card_fills_pass', checks.get('slide9_content_card_fills_pass', False)),
    ]

    passed_count = sum(1 for _, v in critical_checks if v)
    total_count = len(critical_checks)
    confidence = round(passed_count / total_count, 2)
    verdict = "PASS" if passed_count == total_count else "FAIL"

    result = {
        "verdict": verdict,
        "confidence": confidence,
        "checks": checks,
        "issues": issues,
        "fix_required": fix_required,
        "summary": {
            "passed": passed_count,
            "total": total_count,
            "critical_check_results": {k: v for k, v in critical_checks}
        }
    }

    print("\n" + "="*60)
    print(f"VERDICT: {verdict}")
    print(f"CONFIDENCE: {confidence}")
    print(f"PASSED: {passed_count} / {total_count}")
    print("="*60)

    for name, val in critical_checks:
        status = "PASS" if val else "FAIL"
        print(f"  [{status}] {name}")

    if issues:
        print("\nISSUES:")
        for issue in issues:
            print(f"  - {issue}")

    print("\nFULL JSON:")
    print(json.dumps(result, indent=2, ensure_ascii=False))

    return result


if __name__ == "__main__":
    main()
