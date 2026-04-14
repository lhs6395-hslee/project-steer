#!/usr/bin/env python3
"""
Reviewer: Verify slides idx 8 and idx 9 in AWS_MSK_Expert_Intro.pptx
Checks flow diagram zone, shape positions, fill contamination, etc.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

# Constants
EMU_PER_INCH = 914400
TOLERANCE = 5000  # EMU

PRIMARY_COLOR = RGBColor(0x00, 0x43, 0xDA)
FILL_COLOR_HEX = "F8F9FA"

# Expected values
FLOW_ZONE_TOP = 2332260
FLOW_ZONE_HEIGHT = 822960
FLOW_ZONE_LEFT = 457200
FLOW_ZONE_WIDTH = 11277600

IDX8_FLOW_TEXT = "1. 설계    →    2. 생성    →    3. 연결    →    4. 운영"
IDX9_FLOW_TEXT = "1. 평가    →    2. 설계    →    3. 구축    →    4. 운영"

# ~3.90" in EMU
TOP_3_90 = int(3.90 * EMU_PER_INCH)
# ~4.80" in EMU
TOP_4_80 = int(4.80 * EMU_PER_INCH)
# ~4.60" in EMU
TOP_4_60 = int(4.60 * EMU_PER_INCH)
# 7.0" in EMU
BOTTOM_7_0 = int(7.0 * EMU_PER_INCH)
# ~0.61" in EMU
TOP_0_61 = int(0.61 * EMU_PER_INCH)
# ~1.75" in EMU
TOP_1_75 = int(1.75 * EMU_PER_INCH)
# ~2.15" in EMU
TOP_2_15 = int(2.15 * EMU_PER_INCH)

POSITION_TOLERANCE = int(0.10 * EMU_PER_INCH)  # ±0.10" tolerance for shifted shapes
EXEMPT_TOLERANCE = int(0.05 * EMU_PER_INCH)    # ±0.05" tolerance for exempt shapes


def emu_to_inches(emu):
    return emu / EMU_PER_INCH


def rgb_from_theme_or_solid(fill):
    """Try to get RGB from solid fill."""
    try:
        if fill.type is not None and fill.type.name == 'SOLID':
            fg = fill.fore_color
            if fg.type is not None and fg.type.name == 'RGB':
                return fg.rgb
    except Exception:
        pass
    return None


def get_shape_text(shape):
    try:
        if shape.has_text_frame:
            return shape.text_frame.text
    except Exception:
        pass
    return ""


def check_color_match(rgb, expected_hex):
    if rgb is None:
        return False
    return str(rgb).upper() == expected_hex.upper()


def check_flow_zone(slide, slide_idx, flow_text, issues, checks_detail):
    """Check 1: Flow diagram zone exists with correct properties."""
    prefix = f"slide_idx{slide_idx}"
    flow_zone_found = False
    flow_zone_issues = []

    for shape in slide.shapes:
        # Look for rounded rectangle with flow text
        shape_text = get_shape_text(shape)
        if flow_text not in shape_text and flow_text.strip() not in shape_text.strip():
            continue

        flow_zone_found = True
        top = shape.top
        height = shape.height
        left = shape.left
        width = shape.width

        # Check position
        if abs(top - FLOW_ZONE_TOP) > TOLERANCE:
            flow_zone_issues.append(
                f"{prefix} flow zone top={top} (expected {FLOW_ZONE_TOP} ±{TOLERANCE}, got diff={abs(top-FLOW_ZONE_TOP)})"
            )
        if abs(height - FLOW_ZONE_HEIGHT) > TOLERANCE:
            flow_zone_issues.append(
                f"{prefix} flow zone height={height} (expected {FLOW_ZONE_HEIGHT} ±{TOLERANCE})"
            )
        if abs(left - FLOW_ZONE_LEFT) > TOLERANCE:
            flow_zone_issues.append(
                f"{prefix} flow zone left={left} (expected {FLOW_ZONE_LEFT})"
            )
        if abs(width - FLOW_ZONE_WIDTH) > TOLERANCE:
            flow_zone_issues.append(
                f"{prefix} flow zone width={width} (expected {FLOW_ZONE_WIDTH})"
            )

        # Check fill color
        try:
            fill_rgb = rgb_from_theme_or_solid(shape.fill)
            if fill_rgb is None or str(fill_rgb).upper() != FILL_COLOR_HEX.upper():
                flow_zone_issues.append(
                    f"{prefix} flow zone fill={fill_rgb} (expected #{FILL_COLOR_HEX})"
                )
        except Exception as e:
            flow_zone_issues.append(f"{prefix} flow zone fill check error: {e}")

        # Check border (line color)
        try:
            line = shape.line
            if line.color and line.color.type is not None:
                line_rgb = line.color.rgb
                if str(line_rgb).upper() != str(PRIMARY_COLOR).upper():
                    flow_zone_issues.append(
                        f"{prefix} flow zone border={line_rgb} (expected {PRIMARY_COLOR})"
                    )
        except Exception as e:
            flow_zone_issues.append(f"{prefix} flow zone border check error: {e}")

        # Check text formatting
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt
                        if abs(pt - 14) > 0.5:
                            flow_zone_issues.append(
                                f"{prefix} flow zone font size={pt}pt (expected 14pt)"
                            )
                    if run.font.bold is False:
                        flow_zone_issues.append(f"{prefix} flow zone font not bold")
                    try:
                        color = run.font.color.rgb
                        if str(color).upper() != str(PRIMARY_COLOR).upper():
                            flow_zone_issues.append(
                                f"{prefix} flow zone text color={color} (expected {PRIMARY_COLOR})"
                            )
                    except Exception:
                        pass
        except Exception as e:
            flow_zone_issues.append(f"{prefix} flow zone text format check error: {e}")

        break  # Found the flow zone shape

    if not flow_zone_found:
        flow_zone_issues.append(f"{prefix} flow zone NOT FOUND (text: '{flow_text[:30]}...')")

    checks_detail[f"{prefix}_flow_zone"] = {
        "found": flow_zone_found,
        "issues": flow_zone_issues
    }
    issues.extend(flow_zone_issues)
    return len(flow_zone_issues) == 0


def check_no_picture_shapes(slide, slide_idx, issues, checks_detail):
    """Check 3: No PICTURE shapes remain."""
    prefix = f"slide_idx{slide_idx}"
    picture_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append(f"shape_idx={shape.shape_id} name='{shape.name}'")

    checks_detail[f"{prefix}_no_pictures"] = {
        "picture_count": len(picture_shapes),
        "pictures": picture_shapes
    }
    if picture_shapes:
        issues.append(f"{prefix} has PICTURE shapes: {picture_shapes}")
        return False
    return True


def check_exempt_shapes(slide, slide_idx, issues, checks_detail):
    """Check 4: Exempt shapes unchanged."""
    prefix = f"slide_idx{slide_idx}"
    exempt_issues = []

    # Expected exempt shapes by name patterns or idx
    # TextBox 2, 3: top ≈ 0.61"
    # TextBox 20: top ≈ 1.75"
    # TextBox 21: top ≈ 2.15"

    # Map shape index (0-based in slide.shapes) to expected top
    exempt_by_idx = {
        2: TOP_0_61,
        3: TOP_0_61,
        20: TOP_1_75,
        21: TOP_2_15,
    }

    shapes_list = list(slide.shapes)
    for shape_idx, expected_top in exempt_by_idx.items():
        if shape_idx < len(shapes_list):
            shape = shapes_list[shape_idx]
            actual_top = shape.top
            if abs(actual_top - expected_top) > EXEMPT_TOLERANCE:
                exempt_issues.append(
                    f"{prefix} exempt shape[{shape_idx}] '{shape.name}' top={actual_top} ({emu_to_inches(actual_top):.3f}\") "
                    f"expected {expected_top} ({emu_to_inches(expected_top):.3f}\")"
                )
        else:
            exempt_issues.append(f"{prefix} exempt shape[{shape_idx}] not found (total shapes: {len(shapes_list)})")

    checks_detail[f"{prefix}_exempt_shapes"] = {
        "issues": exempt_issues
    }
    issues.extend(exempt_issues)
    return len(exempt_issues) == 0


def check_content_no_bottom_overflow(slide, slide_idx, issues, checks_detail):
    """Check 5: No content textbox bottom > 7.0"."""
    prefix = f"slide_idx{slide_idx}"
    overflow_issues = []
    shapes_list = list(slide.shapes)

    # Check all shapes (except exempt 2,3,20,21 and flow zone)
    exempt_indices = {2, 3, 20, 21}

    for i, shape in enumerate(shapes_list):
        if i in exempt_indices:
            continue
        bottom = shape.top + shape.height
        if bottom > BOTTOM_7_0:
            overflow_issues.append(
                f"{prefix} shape[{i}] '{shape.name}' bottom={bottom} ({emu_to_inches(bottom):.3f}\") > 7.0\""
            )

    checks_detail[f"{prefix}_no_bottom_overflow"] = {
        "overflow_count": len(overflow_issues),
        "issues": overflow_issues
    }
    issues.extend(overflow_issues)
    return len(overflow_issues) == 0


def check_fill_contamination(slide, slide_idx, content_tb_indices, issues, checks_detail):
    """Check 6: No fill contamination on content textboxes."""
    prefix = f"slide_idx{slide_idx}"
    contamination_issues = []
    shapes_list = list(slide.shapes)

    for i in content_tb_indices:
        if i >= len(shapes_list):
            contamination_issues.append(f"{prefix} TextBox[{i}] not found")
            continue
        shape = shapes_list[i]
        try:
            fill = shape.fill
            if fill.type is not None and fill.type.name not in ('BACKGROUND', 'NO_FILL'):
                contamination_issues.append(
                    f"{prefix} TextBox[{i}] '{shape.name}' has fill type={fill.type.name} (expected no fill)"
                )
        except Exception as e:
            contamination_issues.append(f"{prefix} TextBox[{i}] fill check error: {e}")

    checks_detail[f"{prefix}_no_fill_contamination"] = {
        "checked_indices": content_tb_indices,
        "issues": contamination_issues
    }
    issues.extend(contamination_issues)
    return len(contamination_issues) == 0


def check_shapes_shifted_down(slide, slide_idx, shape_spec, issues, checks_detail):
    """Check 2: Content shapes shifted down."""
    prefix = f"slide_idx{slide_idx}"
    shift_issues = []
    shapes_list = list(slide.shapes)

    for spec in shape_spec:
        idx = spec["idx"]
        expected_top = spec["expected_top"]
        label = spec["label"]
        if idx >= len(shapes_list):
            shift_issues.append(f"{prefix} {label}[{idx}] not found")
            continue
        shape = shapes_list[idx]
        actual_top = shape.top
        if abs(actual_top - expected_top) > POSITION_TOLERANCE:
            shift_issues.append(
                f"{prefix} {label}[{idx}] '{shape.name}' top={actual_top} ({emu_to_inches(actual_top):.3f}\") "
                f"expected ~{emu_to_inches(expected_top):.2f}\" ({expected_top})"
            )

    checks_detail[f"{prefix}_shapes_shifted"] = {
        "issues": shift_issues
    }
    issues.extend(shift_issues)
    return len(shift_issues) == 0


def print_slide_shapes(slide, slide_idx):
    """Debug: print all shapes with indices."""
    print(f"\n=== Slide idx {slide_idx} shapes ===")
    for i, shape in enumerate(slide.shapes):
        top_in = emu_to_inches(shape.top)
        bottom_in = emu_to_inches(shape.top + shape.height)
        text_preview = get_shape_text(shape)[:40].replace('\n', '|') if get_shape_text(shape) else ""
        print(f"  [{i:2d}] id={shape.shape_id:3d} type={str(shape.shape_type)} name='{shape.name}' "
              f"top={top_in:.3f}\" bottom={bottom_in:.3f}\" text='{text_preview}'")


def main():
    prs = Presentation(PPTX_PATH)
    slides = prs.slides

    print(f"Total slides: {len(slides)}")
    if len(slides) < 10:
        print(f"ERROR: Need at least 10 slides, got {len(slides)}")
        sys.exit(1)

    slide8 = slides[8]
    slide9 = slides[9]

    # Debug: print all shapes
    print_slide_shapes(slide8, 8)
    print_slide_shapes(slide9, 9)

    issues = []
    checks_detail = {}

    # === Slide idx 8 ===
    # Check 1: Flow zone
    check_flow_zone(slide8, 8, IDX8_FLOW_TEXT, issues, checks_detail)

    # Check 3: No pictures
    check_no_picture_shapes(slide8, 8, issues, checks_detail)

    # Check 4: Exempt shapes
    check_exempt_shapes(slide8, 8, issues, checks_detail)

    # Check 2: Content shapes shifted
    # idx8: Chevron 4,8,12,16 and TextBox 5,9,13,17 → top ≈ 3.90"
    # idx8: Rounded Rectangle 6,10,14,18 → top ≈ 4.80"
    shape_spec_8 = []
    for idx in [4, 8, 12, 16]:
        shape_spec_8.append({"idx": idx, "expected_top": TOP_3_90, "label": "Chevron"})
    for idx in [5, 9, 13, 17]:
        shape_spec_8.append({"idx": idx, "expected_top": TOP_3_90, "label": "TextBox"})
    for idx in [6, 10, 14, 18]:
        shape_spec_8.append({"idx": idx, "expected_top": TOP_4_80, "label": "RndRect"})
    check_shapes_shifted_down(slide8, 8, shape_spec_8, issues, checks_detail)

    # Check 5: No bottom overflow
    check_content_no_bottom_overflow(slide8, 8, issues, checks_detail)

    # Check 6: Fill contamination on TextBox 7,11,15,19
    check_fill_contamination(slide8, 8, [7, 11, 15, 19], issues, checks_detail)

    # === Slide idx 9 ===
    # Check 1: Flow zone
    check_flow_zone(slide9, 9, IDX9_FLOW_TEXT, issues, checks_detail)

    # Check 3: No pictures
    check_no_picture_shapes(slide9, 9, issues, checks_detail)

    # Check 4: Exempt shapes
    check_exempt_shapes(slide9, 9, issues, checks_detail)

    # Check 2: Content shapes shifted
    # idx9: Rounded Rectangle 4,8,12,16 and TextBox 5,9,13,17 → top ≈ 3.90"
    # idx9: Rounded Rectangle 6,10,14,18 → top ≈ 4.60"
    TOP_4_60_VAL = int(4.60 * EMU_PER_INCH)
    shape_spec_9 = []
    for idx in [4, 8, 12, 16]:
        shape_spec_9.append({"idx": idx, "expected_top": TOP_3_90, "label": "RndRect"})
    for idx in [5, 9, 13, 17]:
        shape_spec_9.append({"idx": idx, "expected_top": TOP_3_90, "label": "TextBox"})
    for idx in [6, 10, 14, 18]:
        shape_spec_9.append({"idx": idx, "expected_top": TOP_4_60_VAL, "label": "RndRect"})
    check_shapes_shifted_down(slide9, 9, shape_spec_9, issues, checks_detail)

    # Check 5: No bottom overflow
    check_content_no_bottom_overflow(slide9, 9, issues, checks_detail)

    # Check 6: Fill contamination on TextBox 7,11,15,19
    check_fill_contamination(slide9, 9, [7, 11, 15, 19], issues, checks_detail)

    # Build verdict
    verdict = "PASS" if len(issues) == 0 else "FAIL"
    confidence = 1.0 if verdict == "PASS" else max(0.0, 1.0 - len(issues) * 0.05)

    fix_required = list(set(
        issue.split(" ")[0] + " " + issue.split(" ")[1] if len(issue.split(" ")) > 1 else issue
        for issue in issues
    ))

    result = {
        "verdict": verdict,
        "confidence": round(confidence, 2),
        "checks": checks_detail,
        "issues": issues,
        "fix_required": fix_required[:20]  # Cap at 20
    }

    print("\n\n=== VERDICT ===")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return result


if __name__ == "__main__":
    main()
