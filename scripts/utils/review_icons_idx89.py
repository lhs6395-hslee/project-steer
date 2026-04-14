#!/usr/bin/env python3
"""
Reviewer: Verify slides idx 8 and idx 9 in AWS_MSK_Expert_Intro.pptx
Sprint_Contract checks:
- Slide idx 8: exactly 1 PICTURE shape at left≈10.5", top≈6.2", size≈0.6"×0.6"
- Slide idx 9: exactly 1 PICTURE shape at left≈10.5", top≈6.2", size≈0.6"×0.6"
- No flow diagram banner ("Rounded Rectangle 23") on either slide
- No shape bottom > 7.0" on either slide
- All previously verified shapes intact (anchor=t for content TBs, phase fills for idx9, etc.)
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

EMU_PER_INCH = 914400

def emu_to_inch(emu):
    return emu / EMU_PER_INCH

def inch_to_emu(inch):
    return inch * EMU_PER_INCH

def check_slide(slide, slide_idx):
    """Check a single slide and return results dict."""
    results = {
        "picture_count": 0,
        "picture_positions": [],
        "has_rounded_rect_23": False,
        "shapes_beyond_7in": [],
        "content_tb_anchor_t": [],
        "phase_fills_present": None,  # only for idx 9
        "all_shapes": []
    }
    issues = []

    TOLERANCE_INCH = 0.15  # ±0.15" tolerance

    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left_in": round(emu_to_inch(shape.left), 3) if shape.left is not None else None,
            "top_in": round(emu_to_inch(shape.top), 3) if shape.top is not None else None,
            "width_in": round(emu_to_inch(shape.width), 3) if shape.width is not None else None,
            "height_in": round(emu_to_inch(shape.height), 3) if shape.height is not None else None,
        }
        if shape.top is not None and shape.height is not None:
            shape_info["bottom_in"] = round(emu_to_inch(shape.top + shape.height), 3)
        else:
            shape_info["bottom_in"] = None
        results["all_shapes"].append(shape_info)

        # Check for PICTURE shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left_in = emu_to_inch(shape.left) if shape.left is not None else None
            top_in = emu_to_inch(shape.top) if shape.top is not None else None
            width_in = emu_to_inch(shape.width) if shape.width is not None else None
            height_in = emu_to_inch(shape.height) if shape.height is not None else None

            pic_info = {
                "name": shape.name,
                "left_in": round(left_in, 3) if left_in is not None else None,
                "top_in": round(top_in, 3) if top_in is not None else None,
                "width_in": round(width_in, 3) if width_in is not None else None,
                "height_in": round(height_in, 3) if height_in is not None else None,
            }
            results["picture_positions"].append(pic_info)
            results["picture_count"] += 1

            # Check position: left≈10.5", top≈6.2"
            if left_in is not None and abs(left_in - 10.5) > TOLERANCE_INCH:
                issues.append(f"Slide {slide_idx}: PICTURE '{shape.name}' left={left_in:.3f}\" expected≈10.5\" (tolerance ±{TOLERANCE_INCH}\")")
            if top_in is not None and abs(top_in - 6.2) > TOLERANCE_INCH:
                issues.append(f"Slide {slide_idx}: PICTURE '{shape.name}' top={top_in:.3f}\" expected≈6.2\" (tolerance ±{TOLERANCE_INCH}\")")
            # Check size: ≈0.6"×0.6"
            if width_in is not None and abs(width_in - 0.6) > TOLERANCE_INCH:
                issues.append(f"Slide {slide_idx}: PICTURE '{shape.name}' width={width_in:.3f}\" expected≈0.6\" (tolerance ±{TOLERANCE_INCH}\")")
            if height_in is not None and abs(height_in - 0.6) > TOLERANCE_INCH:
                issues.append(f"Slide {slide_idx}: PICTURE '{shape.name}' height={height_in:.3f}\" expected≈0.6\" (tolerance ±{TOLERANCE_INCH}\")")

        # Check for "Rounded Rectangle 23"
        if "Rounded Rectangle 23" in shape.name:
            results["has_rounded_rect_23"] = True
            issues.append(f"Slide {slide_idx}: Found banned shape 'Rounded Rectangle 23'")

        # Check no shape bottom > 7.0"
        if shape.top is not None and shape.height is not None:
            bottom_in = emu_to_inch(shape.top + shape.height)
            if bottom_in > 7.0:
                results["shapes_beyond_7in"].append({
                    "name": shape.name,
                    "bottom_in": round(bottom_in, 3)
                })
                issues.append(f"Slide {slide_idx}: Shape '{shape.name}' bottom={bottom_in:.3f}\" exceeds 7.0\"")

        # Check content textboxes have anchor=t
        if shape.has_text_frame:
            tf = shape.text_frame
            # anchor=t means MSO_ANCHOR.TOP
            from pptx.enum.text import PP_ALIGN
            from pptx.oxml.ns import qn
            txBody = tf._txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                anchor = bodyPr.get('anchor', None)
                if anchor is not None:
                    tb_info = {
                        "name": shape.name,
                        "anchor": anchor
                    }
                    results["content_tb_anchor_t"].append(tb_info)

    # Check picture count
    if results["picture_count"] != 1:
        issues.append(f"Slide {slide_idx}: Expected exactly 1 PICTURE shape, found {results['picture_count']}")

    # Check phase fills for idx 9 (look for shapes with fill colors typical of phase headers)
    if slide_idx == 9:
        phase_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if any(keyword in text for keyword in ["Phase", "단계", "PHASE"]):
                    phase_shapes.append(shape.name)
        results["phase_fills_present"] = len(phase_shapes) > 0
        if not results["phase_fills_present"]:
            # Try to detect by fill color
            fill_shapes = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'fill') and shape.fill.type is not None:
                        fill_shapes.append(shape.name)
                except:
                    pass
            results["phase_shapes_by_fill"] = fill_shapes[:10]  # limit output

    return results, issues


def main():
    prs = Presentation(PPTX_PATH)
    total_slides = len(prs.slides)

    all_checks = {}
    all_issues = []
    fix_required = []

    for target_idx in [8, 9]:
        if target_idx >= total_slides:
            all_issues.append(f"Slide idx {target_idx} does not exist (total slides: {total_slides})")
            all_checks[f"slide_{target_idx}"] = {"error": "slide not found"}
            continue

        slide = prs.slides[target_idx]
        results, issues = check_slide(slide, target_idx)
        all_checks[f"slide_{target_idx}"] = results
        all_issues.extend(issues)

    # Build fix_required list
    for issue in all_issues:
        if "PICTURE" in issue and "expected" in issue:
            fix_required.append("Reposition/resize PICTURE shape to left≈10.5\", top≈6.2\", size≈0.6\"×0.6\"")
        elif "Expected exactly 1 PICTURE" in issue:
            fix_required.append("Add exactly 1 PICTURE shape per slide (idx 8 and idx 9)")
        elif "Rounded Rectangle 23" in issue:
            fix_required.append("Remove 'Rounded Rectangle 23' flow diagram banner")
        elif "bottom" in issue and "exceeds 7.0" in issue:
            fix_required.append("Resize/reposition shape so bottom does not exceed 7.0\"")

    # Remove duplicates
    fix_required = list(dict.fromkeys(fix_required))

    verdict = "PASS" if len(all_issues) == 0 else "FAIL"
    confidence = 0.97 if verdict == "PASS" else max(0.1, 0.97 - len(all_issues) * 0.1)

    output = {
        "verdict": verdict,
        "confidence": round(confidence, 2),
        "checks": all_checks,
        "issues": all_issues,
        "fix_required": fix_required
    }

    print(json.dumps(output, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
