"""
Reviewer: Independent verification of slides idx8 and idx9
Checks Sprint_Contract execution results without access to executor reasoning.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"
EMU_TOLERANCE = 5000

def emu_to_inches(emu):
    return round(emu / 914400, 4)

def check_color(run, expected_hex):
    """Check if a text run's color matches expected hex."""
    try:
        color = run.font.color.rgb
        return str(color).upper() == expected_hex.upper()
    except Exception:
        return False

def get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None

def get_all_text_colors(shape):
    """Return all unique RGB colors from text runs in a shape."""
    colors = set()
    if not shape.has_text_frame:
        return colors
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color.rgb
                colors.add(str(c).upper())
            except Exception:
                colors.add("INHERITED")
    return colors

def verify_slide(slide, slide_label, checks, issues, fix_required):
    """Generic slide verifier."""
    shape_names = {s.name: s for s in slide.shapes}
    results = {}

    # --- Check: No Rounded Rectangle 23 ---
    rr23_exists = "Rounded Rectangle 23" in shape_names
    results["no_rounded_rectangle_23"] = not rr23_exists
    if rr23_exists:
        issues.append(f"{slide_label}: 'Rounded Rectangle 23' still exists")
        fix_required.append(f"{slide_label}: Remove Rounded Rectangle 23")

    # --- Check: No PICTURE shapes ---
    pic_shapes = [s.name for s in slide.shapes if s.shape_type == 13]
    results["no_picture_shapes"] = len(pic_shapes) == 0
    if pic_shapes:
        issues.append(f"{slide_label}: PICTURE shapes found: {pic_shapes}")
        fix_required.append(f"{slide_label}: Remove picture shapes: {pic_shapes}")

    return results

def verify_idx8(slide, checks, issues, fix_required):
    label = "idx8(L04)"
    shape_names = {s.name: s for s in slide.shapes}
    results = verify_slide(slide, label, checks, issues, fix_required)

    # Chevron/TextBox tops ≈ 2.80" = 2560320 EMU
    TOP_EXPECTED_HEADER = 2560320
    for name in ["Chevron 4", "Chevron 8", "Chevron 12", "Chevron 16",
                 "TextBox 5", "TextBox 9", "TextBox 13", "TextBox 17"]:
        if name in shape_names:
            top = shape_names[name].top
            ok = abs(top - TOP_EXPECTED_HEADER) <= EMU_TOLERANCE
            results[f"{name}_top"] = {"value": top, "expected": TOP_EXPECTED_HEADER, "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} top={top} expected≈{TOP_EXPECTED_HEADER} (diff={abs(top-TOP_EXPECTED_HEADER)})")
                fix_required.append(f"{label}: Fix {name} top to {TOP_EXPECTED_HEADER}")
        else:
            results[f"{name}_top"] = {"value": None, "expected": TOP_EXPECTED_HEADER, "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # Rounded Rectangle 6/10/14/18 tops ≈ 3.70" = 3383280 EMU
    TOP_EXPECTED_CARD = 3383280
    for name in ["Rounded Rectangle 6", "Rounded Rectangle 10",
                 "Rounded Rectangle 14", "Rounded Rectangle 18"]:
        if name in shape_names:
            top = shape_names[name].top
            ok = abs(top - TOP_EXPECTED_CARD) <= EMU_TOLERANCE
            results[f"{name}_top"] = {"value": top, "expected": TOP_EXPECTED_CARD, "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} top={top} expected≈{TOP_EXPECTED_CARD} (diff={abs(top-TOP_EXPECTED_CARD)})")
                fix_required.append(f"{label}: Fix {name} top to {TOP_EXPECTED_CARD}")
        else:
            results[f"{name}_top"] = {"value": None, "expected": TOP_EXPECTED_CARD, "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # TextBox 7/11/15/19 text color = #212121
    for name in ["TextBox 7", "TextBox 11", "TextBox 15", "TextBox 19"]:
        if name in shape_names:
            colors = get_all_text_colors(shape_names[name])
            # Allow empty (no runs) or all 212121
            bad_colors = colors - {"212121", "INHERITED"}
            ok = len(bad_colors) == 0
            results[f"{name}_color"] = {"colors": list(colors), "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} has non-DARK_GRAY colors: {bad_colors}")
                fix_required.append(f"{label}: Set {name} text color to #212121")
        else:
            results[f"{name}_color"] = {"colors": [], "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # No content bottom > 7.0" = 6401520 EMU
    MAX_BOTTOM = 6401520
    for s in slide.shapes:
        bottom = s.top + s.height
        if bottom > MAX_BOTTOM:
            results[f"overflow_{s.name}"] = {"bottom": bottom, "pass": False}
            issues.append(f"{label}: '{s.name}' bottom={bottom} > {MAX_BOTTOM} (7.0\")")
            fix_required.append(f"{label}: Fix overflow for '{s.name}'")

    checks[label] = results
    return results

def verify_idx9(slide, checks, issues, fix_required):
    label = "idx9(L05)"
    shape_names = {s.name: s for s in slide.shapes}
    results = verify_slide(slide, label, checks, issues, fix_required)

    # RR4/8/12/16 and TB5/9/13/17 tops ≈ 2.80" = 2560320 EMU
    TOP_EXPECTED_HEADER = 2560320
    for name in ["Rounded Rectangle 4", "Rounded Rectangle 8",
                 "Rounded Rectangle 12", "Rounded Rectangle 16",
                 "TextBox 5", "TextBox 9", "TextBox 13", "TextBox 17"]:
        if name in shape_names:
            top = shape_names[name].top
            ok = abs(top - TOP_EXPECTED_HEADER) <= EMU_TOLERANCE
            results[f"{name}_top"] = {"value": top, "expected": TOP_EXPECTED_HEADER, "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} top={top} expected≈{TOP_EXPECTED_HEADER} (diff={abs(top-TOP_EXPECTED_HEADER)})")
                fix_required.append(f"{label}: Fix {name} top to {TOP_EXPECTED_HEADER}")
        else:
            results[f"{name}_top"] = {"value": None, "expected": TOP_EXPECTED_HEADER, "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # RR6/10/14/18 tops ≈ 3.50" = 3200400 EMU
    TOP_EXPECTED_CARD = 3200400
    for name in ["Rounded Rectangle 6", "Rounded Rectangle 10",
                 "Rounded Rectangle 14", "Rounded Rectangle 18"]:
        if name in shape_names:
            top = shape_names[name].top
            ok = abs(top - TOP_EXPECTED_CARD) <= EMU_TOLERANCE
            results[f"{name}_top"] = {"value": top, "expected": TOP_EXPECTED_CARD, "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} top={top} expected≈{TOP_EXPECTED_CARD} (diff={abs(top-TOP_EXPECTED_CARD)})")
                fix_required.append(f"{label}: Fix {name} top to {TOP_EXPECTED_CARD}")
        else:
            results[f"{name}_top"] = {"value": None, "expected": TOP_EXPECTED_CARD, "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # CRITICAL: TextBox 11 text color = #212121 (NOT #FFFFFF)
    if "TextBox 11" in shape_names:
        colors = get_all_text_colors(shape_names["TextBox 11"])
        has_white = "FFFFFF" in colors
        all_dark = len(colors - {"212121", "INHERITED"}) == 0
        ok = not has_white and (all_dark or len(colors) == 0)
        results["TextBox_11_color_CRITICAL"] = {
            "colors": list(colors),
            "has_white": has_white,
            "pass": ok
        }
        if not ok:
            issues.append(f"{label}: CRITICAL - TextBox 11 has colors {colors}, expected #212121 only")
            fix_required.append(f"{label}: CRITICAL - Set TextBox 11 text color to #212121")
    else:
        results["TextBox_11_color_CRITICAL"] = {"colors": [], "pass": False}
        issues.append(f"{label}: CRITICAL - TextBox 11 not found")

    # TextBox 7/15/19 text color = #212121
    for name in ["TextBox 7", "TextBox 15", "TextBox 19"]:
        if name in shape_names:
            colors = get_all_text_colors(shape_names[name])
            bad_colors = colors - {"212121", "INHERITED"}
            ok = len(bad_colors) == 0
            results[f"{name}_color"] = {"colors": list(colors), "pass": ok}
            if not ok:
                issues.append(f"{label}: {name} has non-DARK_GRAY colors: {bad_colors}")
                fix_required.append(f"{label}: Set {name} text color to #212121")
        else:
            results[f"{name}_color"] = {"colors": [], "pass": False}
            issues.append(f"{label}: Shape '{name}' not found")

    # Phase header fills: RR4=#001B5E, RR8=#0043DA, RR12=#EE8150, RR16=#4CB88F
    phase_fills = {
        "Rounded Rectangle 4": "001B5E",
        "Rounded Rectangle 8": "0043DA",
        "Rounded Rectangle 12": "EE8150",
        "Rounded Rectangle 16": "4CB88F",
    }
    for name, expected_hex in phase_fills.items():
        if name in shape_names:
            shape = shape_names[name]
            try:
                fill = shape.fill
                fg = fill.fore_color.rgb
                ok = str(fg).upper() == expected_hex.upper()
                results[f"{name}_fill"] = {"value": str(fg).upper(), "expected": expected_hex, "pass": ok}
                if not ok:
                    issues.append(f"{label}: {name} fill={str(fg).upper()} expected={expected_hex}")
                    fix_required.append(f"{label}: Restore {name} fill to #{expected_hex}")
            except Exception as e:
                results[f"{name}_fill"] = {"value": "ERROR", "expected": expected_hex, "pass": False, "error": str(e)}
                issues.append(f"{label}: Could not read fill for {name}: {e}")
        else:
            results[f"{name}_fill"] = {"value": None, "expected": expected_hex, "pass": False}
            issues.append(f"{label}: Shape '{name}' not found for fill check")

    # No content bottom > 7.0" = 6401520 EMU
    MAX_BOTTOM = 6401520
    for s in slide.shapes:
        bottom = s.top + s.height
        if bottom > MAX_BOTTOM:
            results[f"overflow_{s.name}"] = {"bottom": bottom, "pass": False}
            issues.append(f"{label}: '{s.name}' bottom={bottom} > {MAX_BOTTOM} (7.0\")")
            fix_required.append(f"{label}: Fix overflow for '{s.name}'")

    checks[label] = results
    return results

def main():
    prs = Presentation(PPTX_PATH)
    slides = list(prs.slides)

    if len(slides) < 10:
        result = {
            "verdict": "FAIL",
            "confidence": 0.0,
            "checks": {},
            "issues": [f"Presentation has only {len(slides)} slides, expected at least 10"],
            "fix_required": ["Verify correct PPTX file"]
        }
        print(json.dumps(result, indent=2, ensure_ascii=False))
        return

    slide_idx8 = slides[8]
    slide_idx9 = slides[9]

    checks = {}
    issues = []
    fix_required = []

    verify_idx8(slide_idx8, checks, issues, fix_required)
    verify_idx9(slide_idx9, checks, issues, fix_required)

    # Count total checks
    total_checks = 0
    passed_checks = 0
    for slide_label, slide_checks in checks.items():
        for check_name, check_val in slide_checks.items():
            total_checks += 1
            if isinstance(check_val, dict):
                if check_val.get("pass", False):
                    passed_checks += 1
            elif isinstance(check_val, bool):
                if check_val:
                    passed_checks += 1

    confidence = round(passed_checks / total_checks, 3) if total_checks > 0 else 0.0
    verdict = "PASS" if len(issues) == 0 else "FAIL"

    result = {
        "verdict": verdict,
        "confidence": confidence,
        "passed_checks": passed_checks,
        "total_checks": total_checks,
        "checks": checks,
        "issues": issues,
        "fix_required": fix_required
    }

    print(json.dumps(result, indent=2, ensure_ascii=False))

if __name__ == "__main__":
    main()
