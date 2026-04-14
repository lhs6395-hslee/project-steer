#!/usr/bin/env python3
"""
Comprehensive quality reviewer for slides idx 5~9 in AWS_MSK_Expert_Intro.pptx
Checks: overlap rules, WHITE text rules, overflow, L01/L02/L04/L05 layout rules

Fix history:
  v2: Fixed XML namespace for spPr (presentationml, not drawingml) and prstGeom search.
      Fixed WHITE text rule: medium-vibrant fills (SUB_ORANGE/SUB_GREEN variants) are
      acceptable backgrounds for white text when TB exactly overlays the colored RR.
      Rule: only flag as CRITICAL if the *smallest* enclosing colored RR is explicitly light.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"
MAX_BOTTOM_EMU = 6401280  # 7.0 inches

# PML namespace (for spPr in slide shapes)
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Known dark fill colors (hex strings, uppercase, no #)
DARK_COLORS = {
    "1E3A5F",  # PRIMARY / DARK_NAVY
    "0D2137",  # deeper navy
    "FF6B35",  # SUB_ORANGE (spec)
    "2ECC71",  # SUB_GREEN (spec)
    "2C3E50",  # dark slate
    "003366",  # dark blue
    "1A2B3C",
    "0F2744",
    "162840",
    "1B3A5C",
    "1F3F6B",
    "243B55",
    "1E3C5F",
    "17304E",
    "152B45",
    "0043DA",  # blue (used in MSK card)
}

# Non-light accent colors that are acceptable backgrounds for WHITE text
# (medium-vibrant: not dark by luminance but clearly colored, not #FFF/F8F9FA)
# Rule: white text on these = PASS (they are SUB_ORANGE/SUB_GREEN variants)
ACCENT_OK_FOR_WHITE = {
    "EE8150",  # SUB_ORANGE variant
    "4CB88F",  # SUB_GREEN variant
    "0043DA",  # blue
    "FF6B35",  # SUB_ORANGE
    "2ECC71",  # SUB_GREEN
    "E6863C",
    "3BB07F",
    "F07840",
    "4ABBA0",
    "E8793A",
    "2D9CDB",
    "EB6B44",
}

# Theme colors considered dark
DARK_THEME_COLORS = {
    MSO_THEME_COLOR.DARK_1,
    MSO_THEME_COLOR.DARK_2,
    MSO_THEME_COLOR.ACCENT_1,
    MSO_THEME_COLOR.ACCENT_2,
    MSO_THEME_COLOR.ACCENT_3,
    MSO_THEME_COLOR.ACCENT_4,
    MSO_THEME_COLOR.ACCENT_5,
    MSO_THEME_COLOR.ACCENT_6,
}

LIGHT_COLORS = {"F8F9FA", "FFFFFF", "FAFAFA", "F0F0F0", "EEEEEE", "E8E8E8", "E6F0FF"}


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def get_fill_info(shape):
    """Return (fill_type, hex_color_or_None, is_dark, is_light, is_transparent)
    is_dark: True if fill is dark enough to support white text
    is_light: True if fill is explicitly light (#FFF, F8F9FA, etc.) — white text would be invisible
    is_transparent: True if shape has no fill (background/none)
    Note: ACCENT_OK_FOR_WHITE colors return is_dark=True so white text is accepted.
    """
    try:
        fill = shape.fill
        ft = fill.type
        if ft is None:
            return ("none", None, False, False, True)
        ft_name = str(ft)

        if "SOLID" in ft_name or ft == 1:  # MSO_FILL_TYPE.SOLID = 1
            try:
                fore = fill.fore_color
                if fore.type is not None:
                    try:
                        rgb = fore.rgb
                        hex_c = rgb_to_hex(rgb)
                        is_dark = (hex_c in DARK_COLORS
                                   or hex_c in ACCENT_OK_FOR_WHITE
                                   or _is_dark_by_luminance(rgb))
                        is_light = (hex_c in LIGHT_COLORS
                                    and hex_c not in ACCENT_OK_FOR_WHITE
                                    and hex_c not in DARK_COLORS)
                        # Also mark as light by luminance only if not in accent list
                        if not is_dark and _is_light_by_luminance(rgb):
                            is_light = True
                        return ("solid", hex_c, is_dark, is_light, False)
                    except Exception:
                        pass
                    try:
                        theme_color = fore.theme_color
                        is_dark = theme_color in DARK_THEME_COLORS
                        return ("solid_theme", str(theme_color), is_dark, False, False)
                    except Exception:
                        pass
            except Exception:
                pass
        if "BACKGROUND" in ft_name or ft == 9:
            return ("background", None, False, False, True)
        if "GRADIENT" in ft_name:
            return ("gradient", None, True, False, False)  # assume dark for gradient
        return (str(ft_name), None, False, False, True)
    except Exception:
        return ("error", None, False, False, True)


def _is_dark_by_luminance(rgb: RGBColor) -> bool:
    r, g, b = rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return lum < 0.35


def _is_light_by_luminance(rgb: RGBColor) -> bool:
    r, g, b = rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return lum > 0.7


def get_text_color_info(run_or_para):
    """Return (hex_color_or_None, is_white, is_explicit_light)"""
    try:
        font = run_or_para.font
        color = font.color
        if color.type is not None:
            try:
                rgb = color.rgb
                hex_c = rgb_to_hex(rgb)
                is_white = hex_c in ("FFFFFF", "FEFEFE", "FDFDFD")
                is_explicit_light = _is_light_by_luminance(rgb)
                return (hex_c, is_white, is_explicit_light)
            except Exception:
                pass
    except Exception:
        pass
    return (None, False, False)


def get_shape_bounds(shape):
    """Return (left, top, right, bottom) in EMU"""
    l = shape.left or 0
    t = shape.top or 0
    w = shape.width or 0
    h = shape.height or 0
    return (l, t, l + w, t + h)


def shapes_overlap(b1, b2):
    """Check if two bounding boxes overlap (not just touch)"""
    l1, t1, r1, bot1 = b1
    l2, t2, r2, bot2 = b2
    return not (r1 <= l2 or r2 <= l1 or bot1 <= t2 or bot2 <= t1)


def is_contained_in(inner, outer, tolerance=18288):  # ~0.02 inch tolerance
    """Check if inner bbox is fully contained in outer bbox"""
    li, ti, ri, bi = inner
    lo, to, ro, bo = outer
    return (li >= lo - tolerance and ti >= to - tolerance and
            ri <= ro + tolerance and bi <= bo + tolerance)


def get_all_shapes_flat(slide):
    """Return flat list of all shapes including those in groups"""
    result = []
    for shape in slide.shapes:
        result.append(shape)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in shape.shapes:
                result.append(sub)
                if sub.shape_type == 6:
                    for subsub in sub.shapes:
                        result.append(subsub)
    return result


def get_picture_shapes(slide):
    """Return list of picture shapes"""
    pics = []
    all_shapes = get_all_shapes_flat(slide)
    for s in all_shapes:
        if s.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            pics.append(s)
    return pics


def has_upper_flow_diagram(slide):
    """
    Check if slide has a flow/diagram structure in the upper portion (top ~40% of slide).
    Heuristics: connectors, or shapes with arrows, or 3+ shapes in a row in upper area.
    """
    slide_height = slide.height if hasattr(slide, 'height') else 6858000  # default
    upper_threshold = slide_height * 0.45

    all_shapes = get_all_shapes_flat(slide)
    # Check for connector shapes
    connectors = [s for s in all_shapes if s.shape_type == 9]  # MSO_SHAPE_TYPE.LINE/CONNECTOR
    upper_connectors = [s for s in connectors if (s.top or 0) < upper_threshold]
    if len(upper_connectors) >= 1:
        return True

    # Check for shapes with arrows (auto shapes with arrow preset geometry)
    # spPr is in PML namespace for slide shapes; prstGeom is in DML namespace
    for s in all_shapes:
        try:
            # Search prstGeom anywhere in the element tree (works regardless of namespace)
            prstGeom = s._element.find(f'.//{{{DML_NS}}}prstGeom')
            if prstGeom is None:
                # Try PML namespace spPr containing DML prstGeom
                sp_pr = s._element.find(f'.//{{{PML_NS}}}spPr')
                if sp_pr is not None:
                    prstGeom = sp_pr.find(f'{{{DML_NS}}}prstGeom')
            if prstGeom is not None:
                prst = prstGeom.get('prst', '')
                if 'arrow' in prst.lower() or 'chevron' in prst.lower():
                    if (s.top or 0) < upper_threshold:
                        return True
        except Exception:
            pass

    # Check for horizontal arrangement of 3+ shapes in upper area
    upper_shapes = [s for s in all_shapes
                    if (s.top or 0) < upper_threshold
                    and s.shape_type not in (13,)  # exclude pictures
                    and s.width and s.width > 200000]

    if len(upper_shapes) >= 3:
        # Check if they are roughly horizontally aligned
        tops = [(s.top or 0) for s in upper_shapes]
        if max(tops) - min(tops) < 500000:  # within ~0.5 inch
            return True

    return False


def check_white_text_rules(slide, slide_idx):
    """
    Check WHITE text rules for a slide.
    Returns list of issues.
    """
    issues = []
    all_shapes = get_all_shapes_flat(slide)

    # Build a map of shapes with their bounds for parent lookup
    shape_bounds = {}
    for s in all_shapes:
        try:
            shape_bounds[id(s)] = get_shape_bounds(s)
        except Exception:
            pass

    # For each text-bearing shape, check white text
    for shape in all_shapes:
        if not shape.has_text_frame:
            continue
        try:
            bounds = get_shape_bounds(shape)
        except Exception:
            continue

        fill_type, fill_color, fill_is_dark, fill_is_light, fill_is_transparent = get_fill_info(shape)

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                hex_c, is_white, is_explicit_light = get_text_color_info(run)
                if not is_white:
                    continue

                # WHITE text found
                if fill_is_transparent or fill_type in ("none", "background"):
                    # Need to check parent/sibling RR for effective background color.
                    # Strategy: find all colored shapes that enclose this TB.
                    # Select the SMALLEST enclosing colored shape as the immediate parent
                    # (the closest visual background). Only that fill determines PASS/FAIL.
                    enclosing = []
                    for other in all_shapes:
                        if id(other) == id(shape):
                            continue
                        try:
                            other_bounds = get_shape_bounds(other)
                            if is_contained_in(bounds, other_bounds):
                                o_fill_type, o_fill_color, o_is_dark, o_is_light, o_is_transparent = get_fill_info(other)
                                if not o_is_transparent:
                                    # Compute area to find smallest (most immediate) parent
                                    ol, ot, or_, ob = other_bounds
                                    area = (or_ - ol) * (ob - ot)
                                    enclosing.append((area, other.name, o_fill_color, o_is_dark, o_is_light))
                        except Exception:
                            pass

                    if enclosing:
                        # Sort by area ascending — smallest area = most immediate parent
                        enclosing.sort(key=lambda x: x[0])
                        immediate_parent = enclosing[0]
                        _, p_name, p_color, p_is_dark, p_is_light = immediate_parent

                        if p_is_light and not p_is_dark:
                            # Immediate parent is light — white text invisible: CRITICAL
                            issues.append({
                                "type": "WHITE_TEXT_NO_DARK_PARENT",
                                "severity": "CRITICAL",
                                "shape": shape.name,
                                "text": run.text[:50],
                                "immediate_parent": p_name,
                                "parent_fill": p_color,
                                "all_enclosing": [(e[1], e[2]) for e in enclosing],
                                "message": (
                                    f"WHITE text on transparent fill; immediate parent '{p_name}' "
                                    f"has light fill #{p_color} — white text not visible"
                                )
                            })
                        # else: immediate parent is dark/accent — PASS
                    else:
                        # No enclosing shape found — floating white text, warn only
                        issues.append({
                            "type": "WHITE_TEXT_NO_PARENT",
                            "severity": "WARNING",
                            "shape": shape.name,
                            "text": run.text[:50],
                            "message": "WHITE text on transparent fill with no enclosing shape found"
                        })
                elif fill_is_light:
                    # Explicit light fill with white text
                    issues.append({
                        "type": "WHITE_TEXT_ON_LIGHT_FILL",
                        "severity": "CRITICAL",
                        "shape": shape.name,
                        "text": run.text[:50],
                        "fill_color": fill_color,
                        "message": f"WHITE text on explicit light fill #{fill_color}"
                    })
                elif fill_is_dark:
                    pass  # PASS: white text on dark fill
                # else: unknown fill, skip

    return issues


def check_overflow(slide):
    """Check if any shape bottom exceeds MAX_BOTTOM_EMU"""
    issues = []
    all_shapes = get_all_shapes_flat(slide)
    for shape in all_shapes:
        try:
            _, _, _, bottom = get_shape_bounds(shape)
            if bottom > MAX_BOTTOM_EMU:
                issues.append({
                    "type": "OVERFLOW",
                    "severity": "FAIL",
                    "shape": shape.name,
                    "bottom_emu": bottom,
                    "max_emu": MAX_BOTTOM_EMU,
                    "overflow_emu": bottom - MAX_BOTTOM_EMU,
                    "message": f"Shape bottom {bottom} EMU > {MAX_BOTTOM_EMU} EMU (7.0\")"
                })
        except Exception:
            pass
    return issues


def check_tb_overlap(slide):
    """
    Check for overlapping TextBoxes from different containers.
    PASS: child TBs overlapping their own parent RR.
    FAIL: TBs from different containers overlapping each other.
    """
    issues = []
    # Get top-level shapes only (not digging into groups for this check)
    top_shapes = list(slide.shapes)

    # Collect text-bearing shapes with bounds
    tb_shapes = []
    for s in top_shapes:
        if s.shape_type == 6:  # group
            for sub in s.shapes:
                if sub.has_text_frame:
                    try:
                        tb_shapes.append((sub, get_shape_bounds(sub), s.name))
                    except Exception:
                        pass
        else:
            if s.has_text_frame:
                try:
                    tb_shapes.append((s, get_shape_bounds(s), None))
                except Exception:
                    pass

    # Check pairs for overlap
    for i in range(len(tb_shapes)):
        for j in range(i + 1, len(tb_shapes)):
            s1, b1, g1 = tb_shapes[i]
            s2, b2, g2 = tb_shapes[j]

            # Skip if same group (parent container)
            if g1 is not None and g1 == g2:
                continue

            if shapes_overlap(b1, b2):
                # Check if one contains the other (z-order child layering — PASS)
                if is_contained_in(b1, b2) or is_contained_in(b2, b1):
                    continue  # legitimate card bg + child TB layering

                issues.append({
                    "type": "TB_OVERLAP",
                    "severity": "FAIL",
                    "shape1": s1.name,
                    "shape2": s2.name,
                    "group1": g1,
                    "group2": g2,
                    "message": f"TextBoxes from different containers overlap: '{s1.name}' (group:{g1}) vs '{s2.name}' (group:{g2})"
                })

    return issues


def check_l01_content_uniqueness(slide):
    """
    L01: left card content must NOT be same as right card content.
    Heuristic: collect text from leftmost card vs other cards, compare.
    """
    issues = []
    all_shapes = get_all_shapes_flat(slide)

    # Group shapes by horizontal region
    slide_width = 9144000  # default 10 inches in EMU
    try:
        slide_width = slide.width
    except Exception:
        pass

    left_texts = []
    right_texts = []

    for s in all_shapes:
        if not s.has_text_frame:
            continue
        try:
            left, top, right, bottom = get_shape_bounds(s)
            center_x = (left + right) / 2
            text = s.text_frame.text.strip()
            if not text:
                continue
            if center_x < slide_width * 0.35:
                left_texts.append(text)
            elif center_x > slide_width * 0.4:
                right_texts.append(text)
        except Exception:
            pass

    # Check for duplicate content
    left_set = set(left_texts)
    right_set = set(right_texts)
    duplicates = left_set & right_set
    if duplicates:
        issues.append({
            "type": "L01_CONTENT_DUPLICATE",
            "severity": "FAIL",
            "duplicates": list(duplicates)[:5],
            "message": f"Left card content duplicated in right cards: {list(duplicates)[:3]}"
        })

    return issues


def check_picture_icon(slide, position="bottom-right"):
    """Check if slide has a PICTURE icon in expected position"""
    pics = get_picture_shapes(slide)
    if not pics:
        return [{
            "type": "MISSING_PICTURE_ICON",
            "severity": "FAIL",
            "message": f"No PICTURE shape found; expected {position} icon"
        }]

    slide_width = 9144000
    slide_height = 6858000
    try:
        slide_width = slide.width
        slide_height = slide.height
    except Exception:
        pass

    # Check if any picture is in bottom-right quadrant
    for pic in pics:
        try:
            l, t, r, b = get_shape_bounds(pic)
            cx = (l + r) / 2
            cy = (t + b) / 2
            if position == "bottom-right":
                if cx > slide_width * 0.55 and cy > slide_height * 0.55:
                    return []  # PASS
        except Exception:
            pass

    return [{
        "type": "PICTURE_ICON_WRONG_POSITION",
        "severity": "FAIL",
        "picture_count": len(pics),
        "message": f"PICTURE shape exists but not in {position} position"
    }]


def review_slide(slide, slide_idx, layout_name):
    """Run all checks for a slide. Returns dict with verdict and issues."""
    issues = []

    # 1. Overflow check
    overflow_issues = check_overflow(slide)
    issues.extend(overflow_issues)

    # 2. WHITE text check
    white_issues = check_white_text_rules(slide, slide_idx)
    issues.extend(white_issues)

    # 3. TB overlap check
    overlap_issues = check_tb_overlap(slide)
    issues.extend(overlap_issues)

    # 4. Layout-specific checks
    if slide_idx == 5:  # L01
        l01_issues = check_l01_content_uniqueness(slide)
        issues.extend(l01_issues)

    if slide_idx == 6:  # L02
        if not has_upper_flow_diagram(slide):
            issues.append({
                "type": "L02_MISSING_FLOW_DIAGRAM",
                "severity": "FAIL",
                "message": "L02 (idx6): upper flow diagram not detected"
            })

    if slide_idx in (8, 9):  # L04, L05
        pic_issues = check_picture_icon(slide, "bottom-right")
        issues.extend(pic_issues)

    # Determine verdict
    critical_or_fail = [i for i in issues if i.get("severity") in ("CRITICAL", "FAIL")]
    verdict = "FAIL" if critical_or_fail else "PASS"

    return {
        "verdict": verdict,
        "layout": layout_name,
        "issue_count": len(issues),
        "critical_count": len(critical_or_fail),
        "issues": issues
    }


def main():
    prs = Presentation(PPTX_PATH)
    slides = prs.slides

    total_slides = len(slides)
    print(f"Total slides: {total_slides}", file=sys.stderr)

    results = {}
    overall_issues = []
    all_fix_required = []

    layout_names = {
        5: "L01",
        6: "L02",
        7: "L03",
        8: "L04",
        9: "L05",
    }

    for idx in range(5, 10):
        if idx >= total_slides:
            results[f"idx{idx}"] = {
                "verdict": "SKIP",
                "issues": [{"type": "SLIDE_NOT_FOUND", "message": f"Slide idx {idx} does not exist (total: {total_slides})"}]
            }
            continue

        slide = slides[idx]
        layout_name = layout_names.get(idx, f"L0{idx-4}")

        # Try to get actual layout name
        try:
            actual_layout = slide.slide_layout.name
            print(f"  Slide idx{idx}: layout='{actual_layout}'", file=sys.stderr)
        except Exception:
            actual_layout = "unknown"

        result = review_slide(slide, idx, actual_layout)
        results[f"idx{idx}"] = {
            "verdict": result["verdict"],
            "layout_detected": actual_layout,
            "layout_expected": layout_name,
            "issue_count": result["issue_count"],
            "critical_count": result["critical_count"],
            "issues": result["issues"]
        }

        if result["issues"]:
            overall_issues.append(f"idx{idx} ({layout_name}): {result['issue_count']} issues")

        for issue in result["issues"]:
            if issue.get("severity") in ("CRITICAL", "FAIL"):
                all_fix_required.append({
                    "slide": f"idx{idx}",
                    "type": issue["type"],
                    "message": issue.get("message", "")
                })

    # Overall verdict
    any_fail = any(r["verdict"] == "FAIL" for r in results.values())
    any_critical = any(
        any(i.get("severity") == "CRITICAL" for i in r["issues"])
        for r in results.values()
    )

    if any_critical:
        overall_verdict = "FAIL"
        confidence = 0.95
    elif any_fail:
        overall_verdict = "FAIL"
        confidence = 0.85
    else:
        overall_verdict = "PASS"
        confidence = 0.90

    output = {
        "verdict": overall_verdict,
        "confidence": confidence,
        "slides": results,
        "overall_issues": overall_issues,
        "fix_required": all_fix_required
    }

    print(json.dumps(output, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
