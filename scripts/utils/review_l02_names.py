#!/usr/bin/env python3
"""
Independent Reviewer: Verify shape renaming and layout checks on slide idx 6.
Checks:
1. OLD names absent
2. NEW names present
3. No overflow (all shapes bottom <= 7.0")
4. Left card content has 4 keywords
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"
SLIDE_IDX = 6

OLD_NAMES = [
    "Rounded Rectangle 15",
    "TextBox 16",
    "TextBox 17",
    "Rounded Rectangle 18",
    "TextBox 19",
    "TextBox 20",
    "Rounded Rectangle 21",
    "TextBox 22",
    "TextBox 23",
]

NEW_NAMES = [
    "Card_Left_BG",
    "Card_Left_Title",
    "Card_Left_Content",
    "VPC_Label",
    "Broker_AZa_Label",
]

KEYWORDS = ["네트워크", "스토리지", "브로커", "복제"]
MAX_BOTTOM_INCHES = 7.0


def emu_to_inches(emu):
    return emu / 914400.0


def run_checks():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[SLIDE_IDX]

    shape_names = [shape.name for shape in slide.shapes]

    checks = {}
    issues = []
    fix_required = []

    # Check 1: OLD names absent
    old_names_found = [n for n in OLD_NAMES if n in shape_names]
    checks["old_names_absent"] = {
        "pass": len(old_names_found) == 0,
        "old_names_still_present": old_names_found,
        "all_shape_names": shape_names,
    }
    if old_names_found:
        issues.append(f"Old names still present: {old_names_found}")
        fix_required.append(f"Rename shapes: {old_names_found}")

    # Check 2: NEW names present
    new_names_missing = [n for n in NEW_NAMES if n not in shape_names]
    checks["new_names_present"] = {
        "pass": len(new_names_missing) == 0,
        "new_names_missing": new_names_missing,
    }
    if new_names_missing:
        issues.append(f"New names missing: {new_names_missing}")
        fix_required.append(f"Add/rename shapes to: {new_names_missing}")

    # Check 3: No overflow (all shapes bottom <= 7.0")
    overflow_shapes = []
    for shape in slide.shapes:
        top_emu = shape.top if shape.top is not None else 0
        height_emu = shape.height if shape.height is not None else 0
        bottom_inches = emu_to_inches(top_emu + height_emu)
        if bottom_inches > MAX_BOTTOM_INCHES:
            overflow_shapes.append({
                "name": shape.name,
                "bottom_inches": round(bottom_inches, 4),
            })

    checks["no_overflow"] = {
        "pass": len(overflow_shapes) == 0,
        "max_allowed_inches": MAX_BOTTOM_INCHES,
        "overflow_shapes": overflow_shapes,
    }
    if overflow_shapes:
        issues.append(f"Shapes exceed 7.0\" bottom boundary: {[s['name'] for s in overflow_shapes]}")
        fix_required.append("Resize or reposition overflowing shapes")

    # Check 4: Left card content has 4 keywords
    card_content_shape = None
    for shape in slide.shapes:
        if shape.name == "Card_Left_Content":
            card_content_shape = shape
            break

    if card_content_shape is None:
        checks["left_card_keywords"] = {
            "pass": False,
            "reason": "Card_Left_Content shape not found",
            "keywords_found": [],
            "keywords_missing": KEYWORDS,
        }
        issues.append("Card_Left_Content shape not found — cannot verify keywords")
        fix_required.append("Ensure Card_Left_Content shape exists with 4 keywords")
    else:
        text = ""
        if card_content_shape.has_text_frame:
            text = card_content_shape.text_frame.text
        keywords_found = [kw for kw in KEYWORDS if kw in text]
        keywords_missing = [kw for kw in KEYWORDS if kw not in text]
        checks["left_card_keywords"] = {
            "pass": len(keywords_missing) == 0,
            "text_snippet": text[:200],
            "keywords_found": keywords_found,
            "keywords_missing": keywords_missing,
        }
        if keywords_missing:
            issues.append(f"Missing keywords in Card_Left_Content: {keywords_missing}")
            fix_required.append(f"Add missing keywords to Card_Left_Content: {keywords_missing}")

    # Determine verdict
    all_pass = all(
        checks[k]["pass"] for k in ["old_names_absent", "new_names_present", "no_overflow", "left_card_keywords"]
    )
    verdict = "PASS" if all_pass else "FAIL"

    # Confidence: based on number of checks passing
    num_checks = 4
    num_passing = sum(
        1 for k in ["old_names_absent", "new_names_present", "no_overflow", "left_card_keywords"]
        if checks[k]["pass"]
    )
    confidence = round(num_passing / num_checks, 2)

    result = {
        "verdict": verdict,
        "confidence": confidence,
        "checks": checks,
        "issues": issues,
        "fix_required": fix_required,
    }

    print(json.dumps(result, ensure_ascii=False, indent=2))
    return result


if __name__ == "__main__":
    result = run_checks()
    sys.exit(0 if result["verdict"] == "PASS" else 1)
