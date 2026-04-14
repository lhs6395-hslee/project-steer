#!/usr/bin/env python3
"""
Reviewer agent: Verify slide idx 6 in AWS_MSK_Expert_Intro.pptx
against Sprint_Contract success criteria.
"""
import json
import sys
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"
SLIDE_IDX = 6

# Constants
EMU_PER_INCH = 914400
INCH_7 = 7 * EMU_PER_INCH  # 6401280 EMU

PRIMARY_HEX = "0043DA"
LIGHT_BLUE_HEX = "EBF4FF"
WHITE_HEX = "FFFFFF"

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# spPr is under p: namespace in shape XML
SP_PR_FILL_PATH = 'p:spPr/a:solidFill/a:srgbClr'
SP_PR_NOFILL_PATH = 'p:spPr/a:noFill'


def hex_to_rgb(h):
    h = h.lstrip('#').upper()
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def color_approx_hex(h1, h2, tol=15):
    """Compare two hex color strings approximately."""
    if h1 is None or h2 is None:
        return False
    try:
        r1, g1, b1 = hex_to_rgb(h1)
        r2, g2, b2 = hex_to_rgb(h2)
        return (abs(r1 - r2) <= tol and abs(g1 - g2) <= tol and abs(b1 - b2) <= tol)
    except Exception:
        return h1.upper() == h2.upper()


def get_fill_hex(shape):
    """
    Extract solid fill hex color from shape XML (spPr/solidFill/srgbClr).
    Returns hex string (e.g. 'EBF4FF') or None.
    """
    xml = shape.element
    # spPr is in 'p:' namespace for shape elements
    for srgb in xml.findall(SP_PR_FILL_PATH, NS):
        val = srgb.get('val')
        if val:
            return val.upper()
    return None


def get_fill_type(shape):
    """Return fill type description: 'solid:XXXXXX', 'none', 'theme:xxx', 'gradient', 'unknown'"""
    xml = shape.element
    # noFill
    if xml.find(SP_PR_NOFILL_PATH, NS) is not None:
        return 'none'
    # solidFill with srgbClr
    for srgb in xml.findall(SP_PR_FILL_PATH, NS):
        val = srgb.get('val', '')
        return f'solid:{val.upper()}'
    # solidFill with schemeClr (theme color)
    for sc in xml.findall('p:spPr/a:solidFill/a:schemeClr', NS):
        val = sc.get('val', '')
        return f'theme:{val}'
    # gradFill
    if xml.find('p:spPr/a:gradFill', NS) is not None:
        return 'gradient'
    # pattFill
    if xml.find('p:spPr/a:pattFill', NS) is not None:
        return 'pattern'
    return 'unknown'


def get_text(shape):
    """Get all text from a shape."""
    try:
        if shape.has_text_frame:
            return shape.text_frame.text
    except Exception:
        pass
    return ""


def get_text_colors(shape):
    """Get all unique text colors from runs in the shape. Returns list of hex strings."""
    colors = set()
    try:
        xml = shape.element
        # Run-level: a:r/a:rPr/a:solidFill/a:srgbClr (under txBody which is in p: namespace)
        for srgb in xml.findall('.//a:r/a:rPr/a:solidFill/a:srgbClr', NS):
            val = srgb.get('val')
            if val:
                colors.add(val.upper())
        # Theme color runs
        for sc in xml.findall('.//a:r/a:rPr/a:solidFill/a:schemeClr', NS):
            val = sc.get('val', '')
            colors.add(f'theme:{val}')
    except Exception:
        pass
    return list(colors)


def get_font_color_default(shape):
    """Get default paragraph/body text color for a shape."""
    try:
        xml = shape.element
        # lstStyle or bodyPr default
        for srgb in xml.findall('.//a:lstStyle//a:solidFill/a:srgbClr', NS):
            val = srgb.get('val')
            if val:
                return val.upper()
    except Exception:
        pass
    return None


def main():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[SLIDE_IDX]
    shapes = list(slide.shapes)

    checks = {}
    issues = []
    fix_required = []

    # Build shape lookup by name
    shape_by_name = {s.name: s for s in shapes}

    print("\n--- Shape fill details (slide 6) ---", file=sys.stderr)
    for s in shapes:
        fill_hex = get_fill_hex(s)
        fill_type = get_fill_type(s)
        text_colors = get_text_colors(s)
        text_preview = get_text(s)[:40].replace('\n', '|')
        print(
            f"  [{s.name}] top={s.top/EMU_PER_INCH:.2f}\" left={s.left/EMU_PER_INCH:.2f}\" "
            f"w={s.width/EMU_PER_INCH:.2f}\" fill_hex=#{fill_hex} fill_type={fill_type} "
            f"text_colors={text_colors} text='{text_preview}'",
            file=sys.stderr
        )

    # ----------------------------------------------------------------
    # Check 1: Old shapes ABSENT
    # ----------------------------------------------------------------
    old_shapes = [
        "Rounded Rectangle 15", "TextBox 16", "TextBox 17",
        "Rounded Rectangle 18", "TextBox 19", "TextBox 20",
        "Rounded Rectangle 21", "TextBox 22", "TextBox 23"
    ]
    present_old = [n for n in old_shapes if n in shape_by_name]
    if present_old:
        checks["old_shapes_absent"] = f"FAIL: Old shapes still present: {present_old}"
        issues.append(f"Old shapes not removed: {present_old}")
        fix_required.append("Remove old shapes from slide 6")
    else:
        checks["old_shapes_absent"] = "PASS: All old shapes absent"

    # ----------------------------------------------------------------
    # Check 2: New left card exists (width >= 4572000, top=4023360, left=457200)
    # ----------------------------------------------------------------
    LEFT_CARD_TOP = 4023360
    LEFT_CARD_LEFT = 457200
    LEFT_CARD_MIN_WIDTH = 4572000
    TOP_TOLERANCE = 100000

    left_card_candidates = []
    for s in shapes:
        if (abs(s.top - LEFT_CARD_TOP) <= TOP_TOLERANCE and
                abs(s.left - LEFT_CARD_LEFT) <= TOP_TOLERANCE and
                s.width >= LEFT_CARD_MIN_WIDTH):
            left_card_candidates.append(s)

    if left_card_candidates:
        lc = left_card_candidates[0]
        checks["left_card_exists"] = (
            f"PASS: Left card found '{lc.name}' "
            f"width={lc.width}, top={lc.top}, left={lc.left}"
        )
    else:
        # Relaxed search: any shape at bottom-left area
        bottom_left = [s for s in shapes
                       if s.top >= 3800000 and s.left <= 600000 and s.width >= 3000000]
        if bottom_left:
            lc = bottom_left[0]
            checks["left_card_exists"] = (
                f"FAIL: Left card position/size mismatch. Found '{lc.name}' "
                f"width={lc.width} (need>={LEFT_CARD_MIN_WIDTH}), "
                f"top={lc.top} (need~{LEFT_CARD_TOP}), left={lc.left} (need~{LEFT_CARD_LEFT})"
            )
        else:
            checks["left_card_exists"] = (
                f"FAIL: No left card found near top={LEFT_CARD_TOP}, left={LEFT_CARD_LEFT}"
            )
        issues.append("Left card shape not found with correct dimensions/position")
        fix_required.append("Create left card with width>=4572000, top=4023360, left=457200")

    # ----------------------------------------------------------------
    # Check 3: Left card title contains "MSK 인프라 핵심 구성"
    # ----------------------------------------------------------------
    TITLE_TEXT = "MSK 인프라 핵심 구성"
    title_found = False
    for s in shapes:
        if s.top >= 3800000 and s.top <= 4700000:
            text = get_text(s)
            if TITLE_TEXT in text:
                title_found = True
                checks["left_card_title"] = f"PASS: Found title '{TITLE_TEXT}' in shape '{s.name}'"
                break

    if not title_found:
        checks["left_card_title"] = f"FAIL: Title '{TITLE_TEXT}' not found in bottom area"
        issues.append(f"Left card title '{TITLE_TEXT}' missing")
        fix_required.append(f"Add title textbox with '{TITLE_TEXT}' in left card area")

    # ----------------------------------------------------------------
    # Check 4: Left card content has all 4 bullet keywords
    # ----------------------------------------------------------------
    KEYWORDS = ["네트워크", "스토리지", "브로커", "복제"]
    found_keywords = set()
    for s in shapes:
        if s.top >= 3800000:
            text = get_text(s)
            for kw in KEYWORDS:
                if kw in text:
                    found_keywords.add(kw)

    missing_kw = [kw for kw in KEYWORDS if kw not in found_keywords]
    if not missing_kw:
        checks["left_card_keywords"] = f"PASS: All 4 keywords found: {KEYWORDS}"
    else:
        checks["left_card_keywords"] = f"FAIL: Missing keywords: {missing_kw}"
        issues.append(f"Missing bullet keywords: {missing_kw}")
        fix_required.append(f"Add missing keywords to left card content: {missing_kw}")

    # ----------------------------------------------------------------
    # Check 5: VPC outer box exists (left >= 5486400, fill ≈ #EBF4FF)
    # ----------------------------------------------------------------
    VPC_MIN_LEFT = 5486400
    vpc_found = False
    vpc_candidates = []
    for s in shapes:
        if s.left >= VPC_MIN_LEFT:
            fill_hex = get_fill_hex(s)
            if fill_hex and color_approx_hex(fill_hex, LIGHT_BLUE_HEX, tol=20):
                vpc_candidates.append(s)

    if vpc_candidates:
        v = vpc_candidates[0]
        checks["vpc_box_exists"] = (
            f"PASS: VPC box found '{v.name}' "
            f"left={v.left}, fill=#{get_fill_hex(v)}"
        )
        vpc_found = True
    else:
        right_shapes = [(s.name, get_fill_hex(s), s.left) for s in shapes if s.left >= VPC_MIN_LEFT]
        checks["vpc_box_exists"] = (
            f"FAIL: No VPC box with #EBF4FF fill found at left>={VPC_MIN_LEFT}. "
            f"Right-side shapes (name, fill_hex, left): {right_shapes}"
        )
        issues.append("VPC outer box with #EBF4FF fill not found at correct position")
        fix_required.append("Add VPC outer box with left>=5486400 and solid fill #EBF4FF")

    # ----------------------------------------------------------------
    # Check 6: Three broker boxes with PRIMARY (#0043DA) fill and white text "Broker"
    # The executor may implement this as: RoundedRectangle (primary fill, no text) +
    # overlapping TextBox (white text "Broker"). Both patterns are acceptable.
    # ----------------------------------------------------------------

    # Pattern A: single shape with primary fill AND broker text
    broker_boxes_single = []
    for s in shapes:
        fill_hex = get_fill_hex(s)
        text = get_text(s)
        is_primary = (fill_hex and color_approx_hex(fill_hex, PRIMARY_HEX, tol=10))
        has_broker_text = ("Broker" in text or "broker" in text.lower())
        if is_primary and has_broker_text:
            broker_boxes_single.append(s)

    # Pattern B: primary-fill shapes (no text) paired with textboxes having "Broker" text + white color
    primary_fill_shapes = [s for s in shapes
                           if get_fill_hex(s) and color_approx_hex(get_fill_hex(s), PRIMARY_HEX, tol=10)]

    # Find textboxes with "Broker" text and white color that overlap a primary shape
    broker_text_boxes = []
    for s in shapes:
        text = get_text(s)
        if "Broker" in text or "broker" in text.lower():
            tc_list = get_text_colors(s)
            has_white = any(color_approx_hex(c, WHITE_HEX, tol=10) for c in tc_list
                            if not c.startswith('theme:'))
            if has_white:
                broker_text_boxes.append(s)

    # Pair: for each primary shape, check if a broker text box overlaps it
    broker_pairs = []
    for pshape in primary_fill_shapes:
        for tbox in broker_text_boxes:
            # Check overlap: text box center within primary shape bounds
            px1, py1 = pshape.left, pshape.top
            px2, py2 = pshape.left + pshape.width, pshape.top + pshape.height
            tx_center = tbox.left + tbox.width // 2
            ty_center = tbox.top + tbox.height // 2
            if px1 <= tx_center <= px2 and py1 <= ty_center <= py2:
                broker_pairs.append((pshape.name, tbox.name))

    total_broker_units = len(broker_boxes_single) + len(set(p[0] for p in broker_pairs))

    if total_broker_units >= 3:
        if broker_boxes_single:
            checks["broker_boxes"] = (
                f"PASS: 3+ broker units found (single-shape pattern): "
                f"{[s.name for s in broker_boxes_single[:3]]}"
            )
        else:
            checks["broker_boxes"] = (
                f"PASS: 3 broker units found (split shape+textbox pattern): "
                f"pairs={broker_pairs}"
            )
    else:
        primary_details = [(s.name, get_fill_hex(s), get_text(s)[:20]) for s in primary_fill_shapes]
        broker_text_details = [(s.name, get_text_colors(s), get_text(s)[:20]) for s in broker_text_boxes]
        checks["broker_boxes"] = (
            f"FAIL: Only {total_broker_units} broker units found (need 3). "
            f"Primary shapes: {primary_details}. "
            f"Broker text boxes: {broker_text_details}. "
            f"Matched pairs: {broker_pairs}"
        )
        issues.append(f"Only {total_broker_units} broker units found, need 3")
        fix_required.append("Add 3 broker boxes with PRIMARY (#0043DA) solid fill and white text containing 'Broker'")

    # ----------------------------------------------------------------
    # Check 7: No shape bottom > 7.0" (6401280 EMU)
    # ----------------------------------------------------------------
    overflow_shapes = []
    for s in shapes:
        bottom = s.top + s.height
        if bottom > INCH_7:
            overflow_shapes.append((s.name, bottom))

    if not overflow_shapes:
        checks["no_overflow"] = "PASS: No shape extends beyond 7.0\""
    else:
        checks["no_overflow"] = f"FAIL: Shapes exceeding 7.0\": {overflow_shapes}"
        issues.append(f"Shapes overflow 7.0\": {overflow_shapes}")
        fix_required.append("Resize/reposition shapes that extend beyond 7.0\"")

    # ----------------------------------------------------------------
    # Check 8: Upper diagram shapes UNCHANGED (top < 4.10" = 3748440 EMU)
    # ----------------------------------------------------------------
    INCH_4_10 = int(4.10 * EMU_PER_INCH)

    upper_check_results = {}

    # Rounded Rectangle 6 and 14: top ≈ 2.80", fill = PRIMARY
    for shape_name in ["Rounded Rectangle 6", "Rounded Rectangle 14"]:
        if shape_name in shape_by_name:
            s = shape_by_name[shape_name]
            top_inches = s.top / EMU_PER_INCH
            fill_hex = get_fill_hex(s)
            fill_type = get_fill_type(s)
            fill_is_primary = fill_hex and color_approx_hex(fill_hex, PRIMARY_HEX, tol=10)
            top_ok = s.top < INCH_4_10

            if top_ok and fill_is_primary:
                upper_check_results[shape_name] = (
                    f"PASS: top={top_inches:.2f}\", fill=#{fill_hex} (PRIMARY)"
                )
            elif not top_ok:
                upper_check_results[shape_name] = (
                    f"FAIL: top={top_inches:.2f}\" (should be <4.10\"), fill=#{fill_hex}"
                )
            elif not fill_is_primary:
                upper_check_results[shape_name] = (
                    f"FAIL: top={top_inches:.2f}\" ok, but fill=#{fill_hex} (expected PRIMARY #{PRIMARY_HEX})"
                )
            else:
                upper_check_results[shape_name] = (
                    f"FAIL: top={top_inches:.2f}\", fill=#{fill_hex}"
                )
        else:
            upper_check_results[shape_name] = f"FAIL: Shape '{shape_name}' not found"

    # Rounded Rectangle 8: top ≈ 2.80" (any fill)
    for shape_name in ["Rounded Rectangle 8"]:
        if shape_name in shape_by_name:
            s = shape_by_name[shape_name]
            top_inches = s.top / EMU_PER_INCH
            top_ok = s.top < INCH_4_10
            if top_ok:
                upper_check_results[shape_name] = f"PASS: top={top_inches:.2f}\""
            else:
                upper_check_results[shape_name] = (
                    f"FAIL: top={top_inches:.2f}\" (should be <4.10\")"
                )
        else:
            upper_check_results[shape_name] = f"FAIL: Shape 'Rounded Rectangle 8' not found"

    # Determine overall: FAIL only if shape is missing or top is wrong
    # WARN is acceptable (fill cannot be determined without theme resolution)
    hard_failures = {k: v for k, v in upper_check_results.items()
                     if v.startswith("FAIL")}
    if not hard_failures:
        checks["upper_shapes_unchanged"] = f"PASS: Upper shapes intact: {upper_check_results}"
    else:
        checks["upper_shapes_unchanged"] = f"FAIL: Upper shape issues: {upper_check_results}"
        issues.append(f"Upper diagram shapes have issues: {hard_failures}")
        fix_required.append("Restore upper diagram shapes to original positions/fills")

    # ----------------------------------------------------------------
    # Check 9: Exempt shapes UNCHANGED
    # ----------------------------------------------------------------
    EXEMPT = {
        "TextBox 2": {"top_inches": 0.61, "tol": 0.1},
        "TextBox 3": {"top_inches": 0.61, "tol": 0.1},
        "TextBox 4": {"top_inches": 1.75, "tol": 0.1},
        "TextBox 5": {"top_inches": 2.15, "tol": 0.1},
    }
    exempt_results = {}
    for name, spec in EXEMPT.items():
        if name in shape_by_name:
            s = shape_by_name[name]
            actual_inches = s.top / EMU_PER_INCH
            expected = spec["top_inches"]
            tol = spec["tol"]
            if abs(actual_inches - expected) <= tol:
                exempt_results[name] = f"PASS: top={actual_inches:.2f}\""
            else:
                exempt_results[name] = (
                    f"FAIL: top={actual_inches:.2f}\" (expected ~{expected}\")"
                )
        else:
            exempt_results[name] = f"FAIL: Shape not found"

    exempt_pass = all("PASS" in v for v in exempt_results.values())
    if exempt_pass:
        checks["exempt_shapes_unchanged"] = f"PASS: Exempt shapes intact: {exempt_results}"
    else:
        checks["exempt_shapes_unchanged"] = f"FAIL: Exempt shape issues: {exempt_results}"
        failed_exempt = {k: v for k, v in exempt_results.items() if "FAIL" in v}
        issues.append(f"Exempt shapes changed: {failed_exempt}")
        fix_required.append("Restore exempt shapes (TextBox 2-5) to original positions")

    # ----------------------------------------------------------------
    # Check 10: No content textbox has WHITE (#FFFFFF) text on white/light background
    # ----------------------------------------------------------------
    # A textbox with white text and noFill is acceptable ONLY if it sits on top of
    # a dark (primary) shape (split shape+textbox pattern).
    # It's a violation if the textbox with white text has no dark shape underneath.
    white_on_white_violations = []

    # Build set of shape names that are "covered" by a primary-fill shape
    covered_by_primary = set()
    for pshape in primary_fill_shapes:
        px1, py1 = pshape.left, pshape.top
        px2, py2 = pshape.left + pshape.width, pshape.top + pshape.height
        for s in shapes:
            if s.name == pshape.name:
                continue
            tx_center = s.left + s.width // 2
            ty_center = s.top + s.height // 2
            if px1 <= tx_center <= px2 and py1 <= ty_center <= py2:
                covered_by_primary.add(s.name)

    for s in shapes:
        tc_list = get_text_colors(s)
        has_explicit_white_text = any(
            color_approx_hex(c, WHITE_HEX, tol=10)
            for c in tc_list if not c.startswith('theme:')
        )
        if has_explicit_white_text:
            fill_hex = get_fill_hex(s)
            fill_type = get_fill_type(s)
            text_preview = get_text(s)[:50]
            if not text_preview.strip():
                continue

            # Check if shape itself has a dark fill
            if fill_hex:
                r, g, b = hex_to_rgb(fill_hex)
                avg = (r + g + b) / 3
                if avg <= 150:
                    # Dark fill — white text is fine
                    continue
                else:
                    white_on_white_violations.append(
                        f"'{s.name}': white text on light fill #{fill_hex}, text='{text_preview}'"
                    )
            elif fill_type == 'none':
                # noFill textbox — check if it's covered by a dark shape
                if s.name in covered_by_primary:
                    # Sitting on a primary-fill shape — acceptable
                    continue
                else:
                    white_on_white_violations.append(
                        f"'{s.name}': white text, noFill, not covered by dark shape, "
                        f"text='{text_preview}'"
                    )
            elif fill_type == 'unknown':
                # No explicit fill defined — likely transparent
                if s.name in covered_by_primary:
                    continue
                else:
                    white_on_white_violations.append(
                        f"'{s.name}': white text, unknown fill (transparent?), not covered by dark shape, "
                        f"text='{text_preview}'"
                    )
            # theme fills: can't determine luminosity without resolving theme — skip

    if not white_on_white_violations:
        checks["no_white_on_white"] = "PASS: No white text on light/transparent background"
    else:
        checks["no_white_on_white"] = (
            f"FAIL: White text on white/light background detected: {white_on_white_violations}"
        )
        issues.extend([f"White-on-white: {v}" for v in white_on_white_violations])
        fix_required.append("Fix text colors: shapes with white text must have dark fill or sit on dark background")

    # ----------------------------------------------------------------
    # Summary
    # ----------------------------------------------------------------
    all_pass = all(v.startswith("PASS") for v in checks.values())
    fail_count = sum(1 for v in checks.values() if v.startswith("FAIL"))
    pass_count = sum(1 for v in checks.values() if v.startswith("PASS"))
    total = len(checks)
    confidence = pass_count / total if total > 0 else 0.0

    verdict = {
        "verdict": "PASS" if all_pass else "FAIL",
        "confidence": round(confidence, 2),
        "checks": checks,
        "issues": issues,
        "fix_required": fix_required,
    }

    print(json.dumps(verdict, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
