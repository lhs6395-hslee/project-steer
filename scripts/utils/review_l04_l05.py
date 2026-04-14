#!/usr/bin/env python3
"""
Reviewer script for Sprint_Contract verification.
Checks slides idx 8 (L04) and idx 9 (L05) in AWS_MSK_Expert_Intro.pptx.
"""

import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

# Template palette (R, G, B)
TEMPLATE_PALETTE = {
    "PRIMARY":    (0, 67, 218),
    "SUB_BLUE":   (42, 111, 222),
    "SUB_ORANGE": (238, 129, 80),
    "SUB_GREEN":  (76, 184, 143),
    "DARK_NAVY":  (0, 27, 94),
}
PALETTE_RGBS = set(TEMPLATE_PALETTE.values())

EXEMPT_SHAPES = {"Text Placeholder 1", "TextBox 2", "TextBox 3"}

# Expected body title textbox geometry (EMU)
BODY_TITLE_GEOMETRY = {
    "left":   457200,
    "top":    1600200,
    "width":  11277600,
    "height": 320040,
}

def get_text_with_newlines(tf):
    """Returns text with paragraph breaks as \\n."""
    return "\n".join(para.text for para in tf.paragraphs)

def get_anchor(shape):
    """Returns the txBody anchor attribute value."""
    if not shape.has_text_frame:
        return None
    txBody = shape.text_frame._txBody
    anchor = txBody.get("anchor")
    return anchor

def get_solid_fill_rgb(shape):
    """Returns (R, G, B) if shape has solid fill, else None."""
    sp = shape._element
    spPr = sp.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr")
    if spPr is None:
        return None
    solidFill = spPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if solidFill is None:
        return None
    srgbClr = solidFill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    if srgbClr is not None:
        val = srgbClr.get("val")
        if val:
            r = int(val[0:2], 16)
            g = int(val[2:4], 16)
            b = int(val[4:6], 16)
            return (r, g, b)
    # Check schemeClr (theme color) — report as None since we can't resolve easily
    return None

def rgb_close(rgb1, rgb2, tol=5):
    return all(abs(a - b) <= tol for a, b in zip(rgb1, rgb2))

def is_palette_color(rgb):
    for pc in PALETTE_RGBS:
        if rgb_close(rgb, pc):
            return True
    return False

def palette_name(rgb):
    for name, pc in TEMPLATE_PALETTE.items():
        if rgb_close(rgb, pc):
            return name
    return f"RGB{rgb}"

def main():
    prs = Presentation(PPTX_PATH)
    slides = prs.slides
    total_slides = len(slides)
    print(f"Total slides: {total_slides}")

    results = {}
    issues = []
    fix_required = []

    # ---- SLIDE IDX 8 (L04) ----
    slide8 = slides[8]
    print(f"\n=== Slide idx 8: {slide8.name if hasattr(slide8, 'name') else 'N/A'} ===")

    textbox2_8 = None
    for shape in slide8.shapes:
        print(f"  Shape: name={shape.name!r}, has_text={shape.has_text_frame}")
        if shape.name == "TextBox 2":
            textbox2_8 = shape

    # Check 1: idx8 TextBox 2 line break
    if textbox2_8 and textbox2_8.has_text_frame:
        text = get_text_with_newlines(textbox2_8.text_frame)
        print(f"  TextBox 2 text: {text!r}")
        paragraphs = [p.text for p in textbox2_8.text_frame.paragraphs]
        print(f"  Paragraphs: {paragraphs}")
        expected = "L04. Process\nArrow"
        if text == expected:
            results["idx8_textbox2_linebreak"] = f"PASS: {text!r}"
        else:
            results["idx8_textbox2_linebreak"] = f"FAIL: actual={text!r}, expected={expected!r}"
            issues.append(f"idx8 TextBox 2 text mismatch: got {text!r}")
            fix_required.append("Slide idx 8 TextBox 2: set text to 'L04. Process\\nArrow' (2 paragraphs)")
    else:
        results["idx8_textbox2_linebreak"] = "FAIL: TextBox 2 not found or has no text frame"
        issues.append("idx8 TextBox 2 not found")
        fix_required.append("Slide idx 8: find and fix TextBox 2")

    # Check 3: idx8 TextBox 2 width unchanged
    if textbox2_8:
        w = textbox2_8.width
        h = textbox2_8.height
        print(f"  TextBox 2 width={w}, height={h}")
        # original expected width from sprint contract
        expected_w = 11277600
        if w == expected_w:
            results["idx8_textbox2_size_unchanged"] = f"PASS: width={w}"
        else:
            results["idx8_textbox2_size_unchanged"] = f"FAIL: width={w} (expected {expected_w})"
            issues.append(f"idx8 TextBox 2 width changed: {w} != {expected_w}")
            fix_required.append(f"Slide idx 8 TextBox 2: restore width to {expected_w}")
    else:
        results["idx8_textbox2_size_unchanged"] = "FAIL: TextBox 2 not found"

    # Check 5: idx8 anchor='ctr' for non-exempt shapes
    anchor_issues_8 = []
    anchor_checked_8 = 0
    for shape in slide8.shapes:
        if shape.name in EXEMPT_SHAPES:
            continue
        if not shape.has_text_frame:
            continue
        anchor = get_anchor(shape)
        anchor_checked_8 += 1
        print(f"  Shape {shape.name!r}: anchor={anchor}")
        if anchor != "ctr":
            anchor_issues_8.append(f"{shape.name!r}(anchor={anchor})")

    if not anchor_issues_8:
        results["idx8_anchor_ctr"] = f"PASS: {anchor_checked_8} shapes checked"
    else:
        results["idx8_anchor_ctr"] = f"FAIL: {anchor_checked_8} shapes checked, issues: {anchor_issues_8}"
        issues.append(f"idx8 anchor issues: {anchor_issues_8}")
        fix_required.append(f"Slide idx 8: set anchor='ctr' on: {anchor_issues_8}")

    # ---- SLIDE IDX 9 (L05) ----
    slide9 = slides[9]
    print(f"\n=== Slide idx 9: {slide9.name if hasattr(slide9, 'name') else 'N/A'} ===")

    textbox2_9 = None
    for shape in slide9.shapes:
        print(f"  Shape: name={shape.name!r}, has_text={shape.has_text_frame}")
        if shape.name == "TextBox 2":
            textbox2_9 = shape

    # Check 2: idx9 TextBox 2 line break
    if textbox2_9 and textbox2_9.has_text_frame:
        text = get_text_with_newlines(textbox2_9.text_frame)
        print(f"  TextBox 2 text: {text!r}")
        paragraphs = [p.text for p in textbox2_9.text_frame.paragraphs]
        print(f"  Paragraphs: {paragraphs}")
        expected = "L05. Phased\nColumns"
        if text == expected:
            results["idx9_textbox2_linebreak"] = f"PASS: {text!r}"
        else:
            results["idx9_textbox2_linebreak"] = f"FAIL: actual={text!r}, expected={expected!r}"
            issues.append(f"idx9 TextBox 2 text mismatch: got {text!r}")
            fix_required.append("Slide idx 9 TextBox 2: set text to 'L05. Phased\\nColumns' (2 paragraphs)")
    else:
        results["idx9_textbox2_linebreak"] = "FAIL: TextBox 2 not found or has no text frame"
        issues.append("idx9 TextBox 2 not found")
        fix_required.append("Slide idx 9: find and fix TextBox 2")

    # Check 4: idx9 TextBox 2 width unchanged
    if textbox2_9:
        w = textbox2_9.width
        h = textbox2_9.height
        print(f"  TextBox 2 width={w}, height={h}")
        expected_w = 11277600
        if w == expected_w:
            results["idx9_textbox2_size_unchanged"] = f"PASS: width={w}"
        else:
            results["idx9_textbox2_size_unchanged"] = f"FAIL: width={w} (expected {expected_w})"
            issues.append(f"idx9 TextBox 2 width changed: {w} != {expected_w}")
            fix_required.append(f"Slide idx 9 TextBox 2: restore width to {expected_w}")
    else:
        results["idx9_textbox2_size_unchanged"] = "FAIL: TextBox 2 not found"

    # Check 6: idx9 anchor='ctr' for non-exempt shapes
    anchor_issues_9 = []
    anchor_checked_9 = 0
    for shape in slide9.shapes:
        if shape.name in EXEMPT_SHAPES:
            continue
        if not shape.has_text_frame:
            continue
        anchor = get_anchor(shape)
        anchor_checked_9 += 1
        print(f"  Shape {shape.name!r}: anchor={anchor}")
        if anchor != "ctr":
            anchor_issues_9.append(f"{shape.name!r}(anchor={anchor})")

    if not anchor_issues_9:
        results["idx9_anchor_ctr"] = f"PASS: {anchor_checked_9} shapes checked"
    else:
        results["idx9_anchor_ctr"] = f"FAIL: {anchor_checked_9} shapes checked, issues: {anchor_issues_9}"
        issues.append(f"idx9 anchor issues: {anchor_issues_9}")
        fix_required.append(f"Slide idx 9: set anchor='ctr' on: {anchor_issues_9}")

    # Check 7: idx9 phase header colors varied (Rounded Rectangle 4, 8, 12, 16)
    phase_header_names = {"Rounded Rectangle 4", "Rounded Rectangle 8", "Rounded Rectangle 12", "Rounded Rectangle 16"}
    phase_colors = {}
    for shape in slide9.shapes:
        if shape.name in phase_header_names:
            rgb = get_solid_fill_rgb(shape)
            phase_colors[shape.name] = rgb
            print(f"  Phase header {shape.name!r}: fill={rgb}")

    found_names = set(phase_colors.keys())
    missing_headers = phase_header_names - found_names
    if missing_headers:
        results["idx9_phase_colors_varied"] = f"FAIL: shapes not found: {missing_headers}"
        issues.append(f"idx9 phase header shapes not found: {missing_headers}")
        fix_required.append(f"Slide idx 9: verify phase header shape names: {missing_headers}")
    else:
        color_values = list(phase_colors.values())
        unique_colors = set(str(c) for c in color_values if c is not None)
        all_palette = all(is_palette_color(c) for c in color_values if c is not None)
        are_distinct = len(unique_colors) == len(phase_header_names)
        none_colors = [name for name, c in phase_colors.items() if c is None]

        color_report = {name: (palette_name(c) if c else "schemeClr/None") for name, c in phase_colors.items()}
        print(f"  Phase color report: {color_report}")

        if none_colors:
            results["idx9_phase_colors_varied"] = f"PARTIAL: some colors use schemeClr (unresolved): {color_report}"
            issues.append(f"idx9 phase headers use schemeClr colors (cannot verify as sRGB): {none_colors}")
            fix_required.append("Slide idx 9: ensure phase headers use explicit sRGB fill from template palette, not theme/scheme colors")
        elif are_distinct and all_palette:
            results["idx9_phase_colors_varied"] = f"PASS: distinct palette colors: {color_report}"
        elif not are_distinct:
            results["idx9_phase_colors_varied"] = f"FAIL: colors not distinct: {color_report}"
            issues.append(f"idx9 phase header colors not distinct: {color_report}")
            fix_required.append("Slide idx 9: assign distinct template palette colors to each phase header")
        elif not all_palette:
            results["idx9_phase_colors_varied"] = f"FAIL: non-palette colors used: {color_report}"
            issues.append(f"idx9 phase header uses non-palette color: {color_report}")
            fix_required.append("Slide idx 9: use only template palette colors for phase headers")

    # Check 8: idx9 content textboxes should not have unexpected solid fills
    content_fill_issues = []
    for shape in slide9.shapes:
        if shape.name in EXEMPT_SHAPES:
            continue
        if shape.name in phase_header_names:
            continue
        if not shape.has_text_frame:
            continue
        rgb = get_solid_fill_rgb(shape)
        if rgb is not None:
            content_fill_issues.append(f"{shape.name!r}: fill={rgb}")
            print(f"  WARNING: content shape {shape.name!r} has solid fill {rgb}")

    if not content_fill_issues:
        results["idx9_content_fill_clean"] = "PASS: no unexpected solid fills on content textboxes"
    else:
        results["idx9_content_fill_clean"] = f"FAIL: unexpected solid fills: {content_fill_issues}"
        issues.append(f"idx9 content shapes have unexpected fills: {content_fill_issues}")
        fix_required.append(f"Slide idx 9: remove solid fills from content textboxes: {content_fill_issues}")

    # ---- VERDICT ----
    all_pass = all("PASS" in v or "PARTIAL" in v for v in results.values())
    has_fail = any("FAIL" in v for v in results.values())
    verdict = "FAIL" if has_fail else "PASS"
    confidence = 0.95 if not has_fail else (0.5 if not issues else 0.85)

    import json
    output = {
        "verdict": verdict,
        "confidence": confidence,
        "checks": results,
        "issues": issues,
        "fix_required": fix_required
    }
    print("\n\n=== VERDICT JSON ===")
    print(json.dumps(output, indent=2, ensure_ascii=False))
    return output

if __name__ == "__main__":
    main()
