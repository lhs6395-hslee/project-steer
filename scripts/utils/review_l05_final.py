#!/usr/bin/env python3
"""
Reviewer: Verify slide idx 9 phase header colors in AWS_MSK_Expert_Intro.pptx

Architecture on slide idx 9:
- Phase header background: Rounded Rectangle N (RR4, RR8, RR12, RR16) — solid fill
- Phase header text:       TextBox N+1 (TB5, TB9, TB13, TB17) — overlaid at same position

Checks:
1. All 4 RR fills are from template palette
2. All 4 fills are distinct
3. Text color in each TextBox is WCAG AA compliant for its background
"""

import json
import sys
from pptx import Presentation
from pptx.oxml.ns import qn

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

# Template palette (r, g, b)
PALETTE = {
    "DARK_NAVY":  (0,   27,  94),
    "PRIMARY":    (0,   67,  218),
    "SUB_BLUE":   (42,  111, 222),
    "SUB_ORANGE": (238, 129, 80),
    "SUB_GREEN":  (76,  184, 143),
}

# WCAG required text color per background name
WCAG_REQUIRED = {
    "DARK_NAVY":  "WHITE",
    "PRIMARY":    "WHITE",
    "SUB_BLUE":   "WHITE",
    "SUB_ORANGE": "BLACK",
    "SUB_GREEN":  "BLACK",
}

WHITE_RGB = (255, 255, 255)
BLACK_RGB = (0, 0, 0)


def rgb_to_hex(r, g, b):
    return f"#{r:02X}{g:02X}{b:02X}"


def relative_luminance(r, g, b):
    """WCAG 2.1 relative luminance"""
    def ch(c):
        s = c / 255.0
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def contrast_ratio(fg, bg):
    l1 = relative_luminance(*fg)
    l2 = relative_luminance(*bg)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def match_palette(rgb_tuple):
    for name, pal_rgb in PALETTE.items():
        if rgb_tuple == pal_rgb:
            return name
    return None


def get_rr_fill(shape):
    """Get solid fill RGB from a Rounded Rectangle via fore_color.rgb
    RGBColor in python-pptx is a tuple (r, g, b) subclass."""
    rgb_obj = shape.fill.fore_color.rgb
    # RGBColor is a tuple (r, g, b)
    return tuple(rgb_obj)


def get_textbox_text_color(shape):
    """Get text color RGB from first run's solidFill srgbClr in TextBox"""
    sp = shape.element
    for r_el in sp.iter(qn('a:r')):
        rpr = r_el.find(qn('a:rPr'))
        if rpr is not None:
            sf = rpr.find('.//' + qn('a:srgbClr'))
            if sf is not None:
                val = sf.get('val')
                if val and len(val) == 6:
                    return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
    return None


def main():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[9]

    # Map shape names to roles
    rr_names = {
        "Rounded Rectangle 4":  "RR4",
        "Rounded Rectangle 8":  "RR8",
        "Rounded Rectangle 12": "RR12",
        "Rounded Rectangle 16": "RR16",
    }
    tb_names = {
        "TextBox 5":  "Phase1",
        "TextBox 9":  "Phase2",
        "TextBox 13": "Phase3",
        "TextBox 17": "Phase4",
    }
    # Phase order
    phase_order = [("RR4", "Phase1"), ("RR8", "Phase2"), ("RR12", "Phase3"), ("RR16", "Phase4")]

    rr_shapes = {}
    tb_shapes = {}
    for shape in slide.shapes:
        if shape.name in rr_names:
            rr_shapes[rr_names[shape.name]] = shape
        if shape.name in tb_names:
            tb_shapes[tb_names[shape.name]] = shape

    issues = []
    phase_colors = {}
    wcag_check = {}
    fill_data = {}  # rr_key -> (fill_rgb, pal_name)

    # --- Check 1 & 2: RR fills from palette, all distinct ---
    for rr_key, phase_label in phase_order:
        if rr_key not in rr_shapes:
            issues.append(f"{rr_key}: shape NOT FOUND on slide idx 9")
            phase_colors[rr_key] = "NOT FOUND"
            continue

        shape = rr_shapes[rr_key]
        fill_rgb = get_rr_fill(shape)
        fill_hex = rgb_to_hex(*fill_rgb)
        pal_name = match_palette(fill_rgb)
        pal_match = "YES" if pal_name else "NO"

        if not pal_name:
            issues.append(f"{rr_key}: fill {fill_hex} {fill_rgb} is NOT in template palette")

        phase_colors[rr_key] = f"color {fill_hex} {fill_rgb} — palette match: {pal_match}" + (f" [{pal_name}]" if pal_name else "")
        fill_data[rr_key] = (fill_rgb, pal_name)

    # Distinct check
    fill_values = [rgb for rgb, _ in fill_data.values()]
    if len(fill_values) != len(set(fill_values)):
        seen = {}
        for rr_key, (rgb, _) in fill_data.items():
            seen.setdefault(rgb, []).append(rr_key)
        dupes = {str(rgb): keys for rgb, keys in seen.items() if len(keys) > 1}
        issues.append(f"Duplicate fill colors: {dupes}")

    # --- Check 3: WCAG text color compliance ---
    for rr_key, phase_label in phase_order:
        if rr_key not in fill_data:
            wcag_check[f"{phase_label}_text"] = "SKIPPED (RR not found)"
            continue
        if phase_label not in tb_shapes:
            issues.append(f"{phase_label}: TextBox NOT FOUND")
            wcag_check[f"{phase_label}_text"] = "NOT FOUND"
            continue

        fill_rgb, pal_name = fill_data[rr_key]
        fill_hex = rgb_to_hex(*fill_rgb)
        tb = tb_shapes[phase_label]
        text_rgb = get_textbox_text_color(tb)

        if text_rgb is None:
            issues.append(f"{phase_label}: could not read text color from TextBox")
            wcag_check[f"{phase_label}_text"] = "UNREADABLE"
            continue

        text_hex = rgb_to_hex(*text_rgb)
        ratio = contrast_ratio(text_rgb, fill_rgb)
        wcag_aa_pass = ratio >= 4.5

        required = WCAG_REQUIRED.get(pal_name, "UNKNOWN") if pal_name else "UNKNOWN"

        if required == "WHITE":
            color_correct = (text_rgb[0] >= 200 and text_rgb[1] >= 200 and text_rgb[2] >= 200)
        elif required == "BLACK":
            color_correct = (text_rgb[0] < 50 and text_rgb[1] < 50 and text_rgb[2] < 50)
        else:
            color_correct = True  # unknown bg, only check ratio

        wcag_status = "PASS" if (wcag_aa_pass and color_correct) else "FAIL"

        if not wcag_aa_pass:
            issues.append(f"{phase_label}: text {text_hex} on {fill_hex} — contrast {ratio:.2f}:1 < 4.5 WCAG AA FAIL")
        if not color_correct and pal_name:
            issues.append(f"{phase_label}: bg is {pal_name}, requires {required} text, but got {text_hex}")

        wcag_check[f"{phase_label}_text"] = (
            f"{text_hex} on {fill_hex} — contrast {ratio:.2f}:1 — {wcag_status}"
            + (f" (required: {required})" if pal_name else "")
        )

    verdict = "PASS" if len(issues) == 0 else "FAIL"
    confidence = 1.0 if verdict == "PASS" else max(0.1, round(1.0 - 0.2 * len(issues), 2))

    result = {
        "verdict": verdict,
        "confidence": confidence,
        "phase_colors": phase_colors,
        "wcag_check": wcag_check,
        "issues": issues,
    }

    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0 if verdict == "PASS" else 1


if __name__ == "__main__":
    sys.exit(main())
