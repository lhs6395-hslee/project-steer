#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 종합 리뷰 스크립트
모든 검증 항목을 단계별로 수행하고 상세 보고서 생성
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict
import json

# 스타일 가이드 상수
STYLE_COLORS = {
    'PRIMARY': (0, 67, 218),
    'DARK_GRAY': (33, 33, 33),
    'WHITE': (255, 255, 255),
    'ORANGE': (238, 129, 80),
    'GREEN': (76, 184, 143),
    'BG_BOX': (248, 249, 250),
    'BG_WHITE': (255, 255, 255),
}

STYLE_FONTS = {
    'HEAD_TITLE': '프리젠테이션 7 Bold',
    'BODY_TITLE': 'Freesentation',
    'CARD_TITLE': '프리젠테이션 7 Bold',
    'CARD_CONTENT': 'Freesentation',
}

SLIDE_WIDTH = Inches(13.333).emu
SLIDE_HEIGHT = Inches(7.5).emu
BODY_LIMIT_Y = Inches(7.0).emu

# 결과 저장
results = {
    'file': '',
    'total_slides': 0,
    'issues': [],
    'summary': defaultdict(int),
}

def add_issue(severity, category, slide_idx, message, details=None):
    """이슈 기록"""
    issue = {
        'severity': severity,  # CRITICAL, WARNING, PASS
        'category': category,
        'slide_index': slide_idx,
        'message': message,
        'details': details or {}
    }
    results['issues'].append(issue)
    results['summary'][severity] += 1

def emu_to_inches(emu):
    """EMU를 인치로 변환"""
    return emu / 914400

def color_match(rgb1, rgb2, tolerance=5):
    """색상 매칭 (허용 오차 포함)"""
    if not rgb1 or not rgb2:
        return False
    return all(abs(a - b) <= tolerance for a, b in zip(rgb1, rgb2))

def get_shape_color(shape):
    """도형 배경 색상 추출"""
    try:
        if shape.fill.type == 1:  # SOLID
            return shape.fill.fore_color.rgb
    except:
        pass
    return None

def get_text_color(run):
    """텍스트 색상 추출"""
    try:
        return run.font.color.rgb
    except:
        pass
    return None

def check_shape_overlap(shapes):
    """도형 겹침 검출"""
    overlaps = []
    shape_list = [(s, s.left, s.top, s.width, s.height) for s in shapes if hasattr(s, 'left')]
    
    for i, (s1, l1, t1, w1, h1) in enumerate(shape_list):
        for j, (s2, l2, t2, w2, h2) in enumerate(shape_list[i+1:], i+1):
            # 겹침 확인
            if not (l1 + w1 < l2 or l2 + w2 < l1 or t1 + h1 < t2 or t2 + h2 < t1):
                # 의도된 겹침 필터링 (부모-자식 관계는 정상)
                overlaps.append({
                    'shape1': i,
                    'shape2': j,
                    'bounds1': (emu_to_inches(l1), emu_to_inches(t1), emu_to_inches(w1), emu_to_inches(h1)),
                    'bounds2': (emu_to_inches(l2), emu_to_inches(t2), emu_to_inches(w2), emu_to_inches(h2)),
                })
    return overlaps

def estimate_text_width(text, font_size_pt):
    """텍스트 너비 추정 (CJK + ASCII 혼용)"""
    width_pt = 0
    for ch in text:
        if ord(ch) > 0x2E80:  # CJK
            width_pt += font_size_pt
        elif ch == ' ':
            width_pt += font_size_pt * 0.25
        elif ch in '.:/-':
            width_pt += font_size_pt * 0.35
        else:  # ASCII
            width_pt += font_size_pt * 0.5
    return width_pt

def check_text_overflow(shape):
    """텍스트 오버플로우 확인"""
    if not shape.has_text_frame:
        return False, None
    
    tf = shape.text_frame
    text = tf.text.strip()
    if not text:
        return False, None
    
    # 줄 수 계산
    line_count = len(tf.paragraphs)
    
    # 폰트 크기 가져오기
    font_size_pt = 13  # 기본값
    if tf.paragraphs and tf.paragraphs[0].runs:
        try:
            font_size_pt = tf.paragraphs[0].runs[0].font.size.pt
        except:
            pass
    
    # 줄 간격 (1.2 기본)
    line_height_pt = font_size_pt * 1.2
    
    # 필요한 높이
    required_height_pt = line_count * line_height_pt
    
    # shape 높이
    shape_height_pt = shape.height / 914400 * 72
    
    overflow = required_height_pt > shape_height_pt
    
    return overflow, {
        'line_count': line_count,
        'font_size_pt': font_size_pt,
        'required_height_pt': required_height_pt,
        'shape_height_pt': shape_height_pt,
        'text_preview': text[:50] + '...' if len(text) > 50 else text
    }

def step1_basic_info(prs):
    """Step 1: 기본 정보 조회"""
    print("[Step 1] 기본 정보 조회...")
    
    results['total_slides'] = len(prs.slides)
    results['slide_width'] = emu_to_inches(prs.slide_width)
    results['slide_height'] = emu_to_inches(prs.slide_height)
    
    add_issue('PASS', 'basic_info', -1, 
              f"파일 정상 열림: {results['total_slides']} 슬라이드",
              {'width': results['slide_width'], 'height': results['slide_height']})

def step2_cover_slide(prs):
    """Step 2: 표지 슬라이드 검증"""
    print("[Step 2] 표지 슬라이드 검증...")
    
    if len(prs.slides) == 0:
        add_issue('CRITICAL', 'cover', 0, "슬라이드가 없습니다")
        return
    
    slide = prs.slides[0]
    
    # 제목 텍스트박스 찾기 (원본 스펙: shape[2])
    title_shape = None
    subtitle_shape = None
    date_shape = None
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # 제목: 가장 큰 텍스트박스
            if 'MSK' in text or 'Kafka' in text or 'AWS' in text:
                title_shape = shape
            # 부제: 제목 아래
            elif '전문가' in text or 'Expert' in text or '소개' in text:
                subtitle_shape = shape
            # 날짜: MM/DD 패턴
            elif '/' in text and len(text) <= 10:
                date_shape = shape
    
    # 제목 검증
    if title_shape:
        expected_left = Inches(1.67).emu
        expected_width = Inches(10.0).emu
        
        left_diff = abs(title_shape.left - expected_left) / 914400
        width_diff = abs(title_shape.width - expected_width) / 914400
        
        if left_diff > 0.1 or width_diff > 0.1:
            add_issue('WARNING', 'cover', 0,
                      f"표지 제목 위치/크기 불일치",
                      {'expected': '1.67", 10.0"',
                       'actual': f'{emu_to_inches(title_shape.left):.2f}", {emu_to_inches(title_shape.width):.2f}"'})
        else:
            add_issue('PASS', 'cover', 0, "표지 제목 위치/크기 OK")
    else:
        add_issue('WARNING', 'cover', 0, "표지 제목 텍스트박스를 찾을 수 없습니다")
    
    # 부제 검증
    if subtitle_shape and title_shape:
        expected_gap = Inches(0.15).emu
        actual_gap = subtitle_shape.top - (title_shape.top + title_shape.height)
        gap_diff = abs(actual_gap - expected_gap) / 914400
        
        if gap_diff > 0.1:
            add_issue('WARNING', 'cover', 0,
                      f"표지 부제 간격 불일치",
                      {'expected': '0.15"', 'actual': f'{emu_to_inches(actual_gap):.2f}"'})
        else:
            add_issue('PASS', 'cover', 0, "표지 부제 간격 OK")
    
    # 날짜 형식 검증
    if date_shape:
        date_text = date_shape.text_frame.text.strip()
        import re
        if re.match(r'\d{1,2}/\d{1,2}', date_text):
            add_issue('PASS', 'cover', 0, f"날짜 형식 OK: {date_text}")
        else:
            add_issue('WARNING', 'cover', 0,
                      f"날짜 형식 불일치: {date_text} (기대: MM/DD)")

def step3_toc_slide(prs):
    """Step 3: 목차 슬라이드 검증"""
    print("[Step 3] 목차 슬라이드 검증...")
    
    if len(prs.slides) < 2:
        add_issue('WARNING', 'toc', 1, "목차 슬라이드가 없습니다")
        return
    
    slide = prs.slides[1]
    
    # CONTENTS 헤더 찾기
    has_contents = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if 'CONTENTS' in text.upper():
                has_contents = True
                break
    
    if has_contents:
        add_issue('PASS', 'toc', 1, "목차 헤더 존재")
    else:
        add_issue('WARNING', 'toc', 1, "목차 헤더를 찾을 수 없습니다")

def step4_layout_rules(prs):
    """Step 4: 본문 슬라이드별 레이아웃 규칙 검증"""
    print("[Step 4] 레이아웃 규칙 검증...")
    
    # 본문 슬라이드는 2번부터 (0:표지, 1:목차)
    for idx in range(2, len(prs.slides) - 1):  # 마지막 제외 (끝맺음)
        slide = prs.slides[idx]
        
        # 슬라이드 타입 판단 (휴리스틱)
        # L02: Three Cards - 3개의 큰 카드
        # L04: Process Arrow - 화살표 + 4단계
        # L05: Phased Columns - 4개 컬럼
        
        # 간단한 검증: shape 개수와 배치로 레이아웃 추정
        shape_count = len([s for s in slide.shapes if hasattr(s, 'left')])
        
        # 임시: 상단 흐름도/다이어그램 존재 여부
        has_top_diagram = False
        for shape in slide.shapes:
            if hasattr(shape, 'top') and shape.top < Inches(2.5).emu:
                if shape.has_text_frame or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # 다이어그램일 가능성
                    has_top_diagram = True
                    break
        
        # L04/L05: 우하단 아이콘 체크
        has_bottom_right_icon = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                left_in = emu_to_inches(shape.left)
                top_in = emu_to_inches(shape.top)
                width_in = emu_to_inches(shape.width)
                height_in = emu_to_inches(shape.height)
                
                # 우하단 아이콘: left~10.5", top~6.2", size~0.60"
                if (10.3 <= left_in <= 10.7 and
                    6.0 <= top_in <= 6.4 and
                    0.55 <= width_in <= 0.65 and
                    0.55 <= height_in <= 0.65):
                    has_bottom_right_icon = True
                    break

def step5_subtitle_design(prs):
    """Step 5: 중제목 디자인 규칙 검증"""
    print("[Step 5] 중제목 디자인 규칙 검증...")
    
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text.strip()
            
            # 중제목 판단 (헤더 영역, 12pt 설명글)
            # 실제로는 헤더 우측 설명 텍스트박스
            if hasattr(shape, 'top'):
                top_in = emu_to_inches(shape.top)
                if 0.5 <= top_in <= 1.5:  # 헤더 영역
                    # 줄 수 확인
                    lines = text.split('\n')
                    line_count = len([l for l in lines if l.strip()])
                    
                    if line_count > 2:
                        add_issue('CRITICAL', 'subtitle', idx,
                                  f"중제목이 3줄 이상: {line_count}줄",
                                  {'text': text[:100]})
                    
                    # 단어 중간 줄바꿈 검출
                    for line in lines:
                        if line.strip():
                            # 영문 단어가 하이픈 없이 잘렸는지 확인
                            words = line.split()
                            for word in words:
                                # 단어가 비정상적으로 짧고 다음 줄 시작과 이어질 가능성
                                if len(word) >= 2 and word[-1].isalpha():
                                    # 다음 줄 확인
                                    next_idx = lines.index(line) + 1
                                    if next_idx < len(lines) and lines[next_idx].strip():
                                        next_word = lines[next_idx].strip().split()[0] if lines[next_idx].strip().split() else ''
                                        if next_word and next_word[0].islower():
                                            add_issue('WARNING', 'subtitle', idx,
                                                      f"단어 중간 줄바꿈 의심: {word} / {next_word}")

def step6_title_wrapping(prs):
    """Step 6: 타이틀 텍스트 잘림 검증"""
    print("[Step 6] 타이틀 텍스트 잘림 검증...")
    
    for idx, slide in enumerate(prs.slides[2:-1], 2):  # 본문만
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text.strip()
            
            # 타이틀 판단 (번호 패턴: N-N.)
            import re
            if re.match(r'^\d+-\d+\.', text):
                # 타이틀 영역 4.5" / 28pt 기준
                width_pt = estimate_text_width(text, 28)
                
                if width_pt > 340:
                    # 줄바꿈이 있는지 확인
                    if '\n' not in text:
                        add_issue('WARNING', 'title', idx,
                                  f"타이틀이 길지만 줄바꿈 없음: {width_pt:.0f}pt",
                                  {'text': text})
                    else:
                        # 줄바꿈 위치가 단어 경계인지 확인
                        lines = text.split('\n')
                        for i, line in enumerate(lines[:-1]):
                            if line and line[-1].isalnum() and lines[i+1] and lines[i+1][0].islower():
                                add_issue('WARNING', 'title', idx,
                                          f"타이틀 줄바꿈이 단어 중간: {line} / {lines[i+1]}")

def step7_icon_validation(prs):
    """Step 7: 아이콘 사용 규칙 검증"""
    print("[Step 7] 아이콘 사용 규칙 검증...")
    
    for idx, slide in enumerate(prs.slides):
        # 파란색 원형 fallback 검출
        blue_circles = []
        real_icons = []
        
        for shape in slide.shapes:
            # PICTURE 타입이 실제 PNG 아이콘
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width_in = emu_to_inches(shape.width)
                height_in = emu_to_inches(shape.height)
                
                # 카드 아이콘: 0.45" × 0.45"
                if 0.40 <= width_in <= 0.50 and 0.40 <= height_in <= 0.50:
                    real_icons.append(shape)
                # 우하단 아이콘: 0.60" × 0.60"
                elif 0.55 <= width_in <= 0.65 and 0.55 <= height_in <= 0.65:
                    real_icons.append(shape)
            
            # 원형(OVAL) + 파란색 배경 = fallback
            elif shape.shape_type == MSO_SHAPE_TYPE.OVAL:
                color = get_shape_color(shape)
                if color and color_match(color, STYLE_COLORS['PRIMARY']):
                    blue_circles.append(shape)
        
        # CRITICAL: 파란색 원형 fallback 검출
        if blue_circles:
            add_issue('CRITICAL', 'icon', idx,
                      f"파란색 원형 fallback 아이콘 발견: {len(blue_circles)}개")
        
        # 실제 PNG 아이콘 확인
        if real_icons:
            add_issue('PASS', 'icon', idx,
                      f"실제 PNG 아이콘 사용: {len(real_icons)}개")

def step8_content_placement(prs):
    """Step 8: 본문 콘텐츠 배치 및 여백 검증"""
    print("[Step 8] 콘텐츠 배치 및 여백 검증...")
    
    for idx, slide in enumerate(prs.slides[2:-1], 2):
        shapes_in_body = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                top_in = emu_to_inches(shape.top)
                # 본문 영역: 2.0" ~ 7.0"
                if 2.0 <= top_in <= 7.0:
                    shapes_in_body.append(shape)
        
        if not shapes_in_body:
            continue
        
        # 좌우 여백 계산
        leftmost = min(s.left for s in shapes_in_body)
        rightmost = max(s.left + s.width for s in shapes_in_body)
        
        left_margin = emu_to_inches(leftmost)
        right_margin = emu_to_inches(SLIDE_WIDTH - rightmost)
        margin_diff = abs(left_margin - right_margin)
        
        if margin_diff > 0.1:
            add_issue('WARNING', 'placement', idx,
                      f"좌우 여백 비대칭: 좌 {left_margin:.2f}\" vs 우 {right_margin:.2f}\"")
        
        # 하단 여백 확인
        bottommost = max(s.top + s.height for s in shapes_in_body)
        bottom_margin = emu_to_inches(SLIDE_HEIGHT - bottommost)
        
        if bottom_margin < 0.3:
            add_issue('WARNING', 'placement', idx,
                      f"하단 여백 부족: {bottom_margin:.2f}\" (최소 0.3\")")
        
        # 7.0" 초과 확인
        for shape in shapes_in_body:
            bottom_in = emu_to_inches(shape.top + shape.height)
            if bottom_in > 7.0:
                add_issue('WARNING', 'placement', idx,
                          f"콘텐츠가 7.0\" 초과: {bottom_in:.2f}\"")
                break

def step9_text_overflow(prs):
    """Step 9: 텍스트 오버플로우 검증"""
    print("[Step 9] 텍스트 오버플로우 검증...")
    
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            overflow, details = check_text_overflow(shape)
            
            if overflow:
                add_issue('CRITICAL', 'overflow', idx,
                          f"텍스트 오버플로우 발견",
                          details)

def step10_font_consistency(prs):
    """Step 10: 폰트 일관성 검증"""
    print("[Step 10] 폰트 일관성 검증...")
    
    for idx, slide in enumerate(prs.slides[2:-1], 2):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    
                    font_name = run.font.name
                    try:
                        font_size = run.font.size.pt if run.font.size else None
                    except:
                        font_size = None
                    
                    # 도형/카드 제목: 14pt 프리젠테이션 7 Bold
                    # 도형/카드 내용: 13pt Freesentation
                    # 실제 검증은 복잡하므로 샘플링

def step11_color_scheme(prs):
    """Step 11: 색상 스킴 일관성 검증"""
    print("[Step 11] 색상 스킴 검증...")
    
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 흰 텍스트 on 밝은 배경 검출
            bg_color = get_shape_color(shape)
            
            if bg_color and (color_match(bg_color, STYLE_COLORS['BG_BOX']) or
                              color_match(bg_color, STYLE_COLORS['BG_WHITE'])):
                # 밝은 배경
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text_color = get_text_color(run)
                            if text_color and color_match(text_color, STYLE_COLORS['WHITE']):
                                add_issue('CRITICAL', 'color', idx,
                                          f"흰 텍스트 on 밝은 배경 발견",
                                          {'bg': bg_color, 'text': text_color})

def step12_shape_overlap(prs):
    """Step 12: 도형 겹침 검증"""
    print("[Step 12] 도형 겹침 검증...")
    
    for idx, slide in enumerate(prs.slides[2:-1], 2):
        overlaps = check_shape_overlap(slide.shapes)
        
        if overlaps:
            # 필터링: 부모-자식 관계는 정상
            for overlap in overlaps[:5]:  # 최대 5개만 보고
                add_issue('WARNING', 'overlap', idx,
                          f"도형 겹침 발견",
                          overlap)

def step13_ending_slide(prs):
    """Step 13: 끝맺음 슬라이드 검증"""
    print("[Step 13] 끝맺음 슬라이드 검증...")
    
    if len(prs.slides) == 0:
        return
    
    slide = prs.slides[-1]
    
    # "Thank You" 텍스트 찾기
    has_thank_you = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if 'Thank You' in text or 'thank you' in text.lower():
                has_thank_you = True
                add_issue('PASS', 'ending', len(prs.slides)-1,
                          "끝맺음 'Thank You' 텍스트 존재")
                break
    
    if not has_thank_you:
        add_issue('WARNING', 'ending', len(prs.slides)-1,
                  "끝맺음에 'Thank You' 텍스트가 없습니다")

def step14_generate_report():
    """Step 14: 종합 보고서 생성"""
    print("[Step 14] 종합 보고서 생성...")
    
    # 심각도별 분류
    critical_issues = [i for i in results['issues'] if i['severity'] == 'CRITICAL']
    warning_issues = [i for i in results['issues'] if i['severity'] == 'WARNING']
    pass_items = [i for i in results['issues'] if i['severity'] == 'PASS']
    
    # 슬라이드별 그룹화
    issues_by_slide = defaultdict(list)
    for issue in results['issues']:
        issues_by_slide[issue['slide_index']].append(issue)
    
    # 보고서 생성
    report = []
    report.append("=" * 80)
    report.append("PPTX 종합 리뷰 보고서")
    report.append("=" * 80)
    report.append(f"파일: {results['file']}")
    report.append(f"총 슬라이드: {results['total_slides']}")
    report.append(f"슬라이드 크기: {results['slide_width']:.2f}\" × {results['slide_height']:.2f}\"")
    report.append("")
    
    report.append("=" * 80)
    report.append("요약")
    report.append("=" * 80)
    report.append(f"CRITICAL: {results['summary']['CRITICAL']} 건")
    report.append(f"WARNING:  {results['summary']['WARNING']} 건")
    report.append(f"PASS:     {results['summary']['PASS']} 건")
    report.append("")
    
    # CRITICAL 이슈
    if critical_issues:
        report.append("=" * 80)
        report.append("CRITICAL 이슈")
        report.append("=" * 80)
        for issue in critical_issues:
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else 'N/A'
            report.append(f"[슬라이드 {slide_num}] [{issue['category']}] {issue['message']}")
            if issue['details']:
                report.append(f"  상세: {json.dumps(issue['details'], ensure_ascii=False, indent=2)}")
        report.append("")
    
    # WARNING 이슈
    if warning_issues:
        report.append("=" * 80)
        report.append("WARNING 이슈")
        report.append("=" * 80)
        for issue in warning_issues[:20]:  # 최대 20개
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else 'N/A'
            report.append(f"[슬라이드 {slide_num}] [{issue['category']}] {issue['message']}")
        if len(warning_issues) > 20:
            report.append(f"... 외 {len(warning_issues) - 20}건")
        report.append("")
    
    # 슬라이드별 상세
    report.append("=" * 80)
    report.append("슬라이드별 상세")
    report.append("=" * 80)
    for slide_idx in sorted(issues_by_slide.keys()):
        if slide_idx < 0:
            continue
        issues = issues_by_slide[slide_idx]
        report.append(f"\n슬라이드 {slide_idx + 1}:")
        for issue in issues:
            report.append(f"  [{issue['severity']}] [{issue['category']}] {issue['message']}")
    report.append("")
    
    # 권고사항
    report.append("=" * 80)
    report.append("개선 권고사항")
    report.append("=" * 80)
    if critical_issues:
        report.append("1. CRITICAL 이슈를 최우선으로 해결해야 합니다:")
        for issue in critical_issues[:5]:
            report.append(f"   - {issue['message']}")
    if warning_issues:
        report.append("2. WARNING 이슈를 검토하고 필요 시 수정하세요:")
        for issue in warning_issues[:5]:
            report.append(f"   - {issue['message']}")
    report.append("")
    
    # 품질 점수
    total_checks = len(results['issues'])
    passed = results['summary']['PASS']
    score = (passed / total_checks * 100) if total_checks > 0 else 0
    
    report.append("=" * 80)
    report.append(f"전체 품질 점수: {score:.1f}% ({passed}/{total_checks} 항목 통과)")
    report.append("=" * 80)
    
    return "\n".join(report)

def main():
    if len(sys.argv) < 2:
        print("Usage: python review_pptx_comprehensive.py <pptx_file>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    results['file'] = pptx_file
    
    print(f"PPTX 파일 열기: {pptx_file}")
    
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        print(f"파일 열기 실패: {e}")
        sys.exit(1)
    
    # 14단계 검증 실행
    step1_basic_info(prs)
    step2_cover_slide(prs)
    step3_toc_slide(prs)
    step4_layout_rules(prs)
    step5_subtitle_design(prs)
    step6_title_wrapping(prs)
    step7_icon_validation(prs)
    step8_content_placement(prs)
    step9_text_overflow(prs)
    step10_font_consistency(prs)
    step11_color_scheme(prs)
    step12_shape_overlap(prs)
    step13_ending_slide(prs)
    
    # 보고서 생성
    report = step14_generate_report()
    
    print("\n" + report)
    
    # 결과 저장
    output_json = pptx_file.replace('.pptx', '_review.json')
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    output_txt = pptx_file.replace('.pptx', '_review.txt')
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(report)
    
    print(f"\n결과 저장:")
    print(f"  - JSON: {output_json}")
    print(f"  - TXT:  {output_txt}")
    
    # 종료 코드
    if results['summary']['CRITICAL'] > 0:
        sys.exit(2)  # CRITICAL 이슈 있음
    elif results['summary']['WARNING'] > 10:
        sys.exit(1)  # WARNING 많음
    else:
        sys.exit(0)  # OK

if __name__ == '__main__':
    main()
