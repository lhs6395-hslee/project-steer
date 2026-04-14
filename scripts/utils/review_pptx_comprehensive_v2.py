#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 종합 리뷰 스크립트 v2
모든 검증 항목을 단계별로 수행하고 상세 보고서 생성
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict
import json

# 스타일 가이드 상수
STYLE_COLORS = {
    'PRIMARY': (0, 67, 218),
    'DARK_GRAY': (33, 33, 33),
    'WHITE': (255, 255, 255),
    'ORANGE': (238, 129, 80),
    'GREEN': (76, 184, 143),
    'BG_BOX': (248, 249, 250),
    'BG_WHITE': (255, 255, 255),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BODY_LIMIT_Y = Inches(7.0)

# 결과 저장
results = {
    'file': '',
    'total_slides': 0,
    'issues': [],
    'summary': defaultdict(int),
}

def add_issue(severity, category, slide_idx, message, details=None):
    """이슈 기록"""
    issue = {
        'severity': severity,
        'category': category,
        'slide_index': slide_idx,
        'message': message,
        'details': details or {}
    }
    results['issues'].append(issue)
    results['summary'][severity] += 1

def emu_to_inches(emu):
    """EMU를 인치로 변환"""
    return emu / 914400

def color_match(rgb1, rgb2, tolerance=5):
    """색상 매칭"""
    if not rgb1 or not rgb2:
        return False
    try:
        r1, g1, b1 = rgb1[0], rgb1[1], rgb1[2]
        r2, g2, b2 = rgb2[0], rgb2[1], rgb2[2]
        return abs(r1-r2) <= tolerance and abs(g1-g2) <= tolerance and abs(b1-b2) <= tolerance
    except:
        return False

def get_shape_color(shape):
    """도형 배경 색상 추출"""
    try:
        if shape.fill.type == 1:  # SOLID
            rgb = shape.fill.fore_color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except:
        pass
    return None

def get_text_color(run):
    """텍스트 색상 추출"""
    try:
        rgb = run.font.color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except:
        pass
    return None

def main():
    if len(sys.argv) < 2:
        print("Usage: python review_pptx_comprehensive_v2.py <pptx_file>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    results['file'] = pptx_file
    
    print(f"PPTX 파일 열기: {pptx_file}")
    
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        print(f"파일 열기 실패: {e}")
        sys.exit(1)
    
    results['total_slides'] = len(prs.slides)
    results['slide_width'] = emu_to_inches(prs.slide_width)
    results['slide_height'] = emu_to_inches(prs.slide_height)
    
    add_issue('PASS', 'basic_info', -1, 
              f"파일 정상 열림: {results['total_slides']} 슬라이드")
    
    print(f"[Step 1] 기본 정보: {results['total_slides']} 슬라이드")
    
    # Step 2: 표지
    print("[Step 2] 표지 슬라이드 검증...")
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        found_title = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if 'MSK' in text or 'AWS' in text or 'Kafka' in text:
                    found_title = True
                    add_issue('PASS', 'cover', 0, f"표지 제목 발견: {text[:50]}")
                    break
        if not found_title:
            add_issue('WARNING', 'cover', 0, "표지 제목을 찾을 수 없습니다")
    
    # Step 3: 목차
    print("[Step 3] 목차 슬라이드 검증...")
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        found_contents = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if 'CONTENTS' in text.upper() or '목차' in text:
                    found_contents = True
                    add_issue('PASS', 'toc', 1, "목차 헤더 발견")
                    break
        if not found_contents:
            add_issue('WARNING', 'toc', 1, "목차 헤더를 찾을 수 없습니다")
    
    # Step 7: 아이콘 검증
    print("[Step 7] 아이콘 검증...")
    for idx, slide in enumerate(prs.slides):
        blue_circles = []
        real_icons = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width_in = emu_to_inches(shape.width)
                height_in = emu_to_inches(shape.height)
                if (0.40 <= width_in <= 0.70 and 0.40 <= height_in <= 0.70):
                    real_icons.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.OVAL:
                color = get_shape_color(shape)
                if color and color_match(color, STYLE_COLORS['PRIMARY']):
                    blue_circles.append(shape)
        
        if blue_circles:
            add_issue('CRITICAL', 'icon', idx,
                      f"파란색 원형 fallback 아이콘 발견: {len(blue_circles)}개")
        elif real_icons:
            add_issue('PASS', 'icon', idx,
                      f"실제 PNG 아이콘 사용: {len(real_icons)}개")
    
    # Step 8: 배치 및 여백
    print("[Step 8] 콘텐츠 배치 검증...")
    for idx, slide in enumerate(prs.slides):
        if idx in [0, 1]:  # 표지, 목차 제외
            continue
        if idx == len(prs.slides) - 1:  # 끝맺음 제외
            continue
        
        shapes_in_body = []
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                top_in = emu_to_inches(shape.top)
                if 2.0 <= top_in <= 7.5:
                    shapes_in_body.append(shape)
        
        if shapes_in_body:
            # 하단 여백 확인
            bottommost = max(s.top + s.height for s in shapes_in_body)
            bottom_in = emu_to_inches(bottommost)
            
            if bottom_in > 7.0:
                add_issue('WARNING', 'placement', idx,
                          f"콘텐츠가 7.0\" 초과: {bottom_in:.2f}\"")
    
    # Step 9: 텍스트 오버플로우
    print("[Step 9] 텍스트 오버플로우 검증...")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            text = tf.text.strip()
            if not text or len(text) < 10:
                continue
            
            # 간단한 오버플로우 추정
            line_count = len(tf.paragraphs)
            if line_count > 10:  # 너무 많은 줄
                add_issue('WARNING', 'overflow', idx,
                          f"텍스트 줄 수 과다: {line_count}줄",
                          {'text_preview': text[:100]})
    
    # Step 11: 색상 검증
    print("[Step 11] 색상 스킴 검증...")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            bg_color = get_shape_color(shape)
            
            if bg_color and (color_match(bg_color, STYLE_COLORS['BG_BOX']) or
                              color_match(bg_color, STYLE_COLORS['BG_WHITE'])):
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text_color = get_text_color(run)
                            if text_color and color_match(text_color, STYLE_COLORS['WHITE']):
                                add_issue('CRITICAL', 'color', idx,
                                          "흰 텍스트 on 밝은 배경 발견")
                                break
    
    # Step 13: 끝맺음
    print("[Step 13] 끝맺음 검증...")
    if len(prs.slides) > 0:
        slide = prs.slides[-1]
        found_thank_you = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if 'Thank You' in text or '감사합니다' in text:
                    found_thank_you = True
                    add_issue('PASS', 'ending', len(prs.slides)-1,
                              "'Thank You' 텍스트 발견")
                    break
        if not found_thank_you:
            add_issue('WARNING', 'ending', len(prs.slides)-1,
                      "끝맺음 텍스트를 찾을 수 없습니다")
    
    # 보고서 생성
    print("[Step 14] 보고서 생성...")
    
    critical_issues = [i for i in results['issues'] if i['severity'] == 'CRITICAL']
    warning_issues = [i for i in results['issues'] if i['severity'] == 'WARNING']
    pass_items = [i for i in results['issues'] if i['severity'] == 'PASS']
    
    report = []
    report.append("=" * 80)
    report.append("PPTX 종합 리뷰 보고서")
    report.append("=" * 80)
    report.append(f"파일: {results['file']}")
    report.append(f"총 슬라이드: {results['total_slides']}")
    report.append(f"슬라이드 크기: {results['slide_width']:.2f}\" × {results['slide_height']:.2f}\"")
    report.append("")
    
    report.append("=" * 80)
    report.append("요약")
    report.append("=" * 80)
    report.append(f"CRITICAL: {results['summary']['CRITICAL']} 건")
    report.append(f"WARNING:  {results['summary']['WARNING']} 건")
    report.append(f"PASS:     {results['summary']['PASS']} 건")
    report.append("")
    
    if critical_issues:
        report.append("=" * 80)
        report.append("CRITICAL 이슈")
        report.append("=" * 80)
        for issue in critical_issues:
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else 'N/A'
            report.append(f"[슬라이드 {slide_num}] [{issue['category']}] {issue['message']}")
            if issue['details']:
                report.append(f"  상세: {json.dumps(issue['details'], ensure_ascii=False)}")
        report.append("")
    
    if warning_issues:
        report.append("=" * 80)
        report.append("WARNING 이슈 (최대 20개)")
        report.append("=" * 80)
        for issue in warning_issues[:20]:
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else 'N/A'
            report.append(f"[슬라이드 {slide_num}] [{issue['category']}] {issue['message']}")
        if len(warning_issues) > 20:
            report.append(f"... 외 {len(warning_issues) - 20}건")
        report.append("")
    
    total_checks = len(results['issues'])
    passed = results['summary']['PASS']
    score = (passed / total_checks * 100) if total_checks > 0 else 0
    
    report.append("=" * 80)
    report.append(f"전체 품질 점수: {score:.1f}% ({passed}/{total_checks} 항목 통과)")
    report.append("=" * 80)
    report.append("")
    
    report.append("개선 권고사항:")
    if critical_issues:
        report.append("1. CRITICAL 이슈를 최우선으로 해결:")
        for issue in critical_issues[:5]:
            report.append(f"   - {issue['message']}")
    if warning_issues:
        report.append("2. WARNING 이슈를 검토:")
        for issue in warning_issues[:5]:
            report.append(f"   - {issue['message']}")
    
    report_text = "\n".join(report)
    print("\n" + report_text)
    
    # 결과 저장
    output_json = pptx_file.replace('.pptx', '_review.json')
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    output_txt = pptx_file.replace('.pptx', '_review.txt')
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(report_text)
    
    print(f"\n결과 저장:")
    print(f"  - JSON: {output_json}")
    print(f"  - TXT:  {output_txt}")
    
    if results['summary']['CRITICAL'] > 0:
        sys.exit(2)
    elif results['summary']['WARNING'] > 10:
        sys.exit(1)
    else:
        sys.exit(0)

if __name__ == '__main__':
    main()
