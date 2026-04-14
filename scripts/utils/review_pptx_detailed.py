#!/usr/bin/env python3
"""
PPTX Detailed Review Utility
Performs comprehensive quality checks on PPTX presentations
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import json

# Constants from style guide
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.500)
BODY_START_Y = Inches(2.0)
BODY_LIMIT_Y = Inches(7.0)
MIN_BOTTOM_MARGIN = Inches(0.3)

# Color constants (RGB)
PRIMARY = RGBColor(0, 112, 192)
DARK_GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = [RGBColor(248, 249, 250), RGBColor(255, 255, 255)]

def inches_to_emu(inches):
    """Convert inches to EMU"""
    return int(inches * 914400)

def emu_to_inches(emu):
    """Convert EMU to inches"""
    return emu / 914400

def get_shape_bounds(shape):
    """Get shape boundaries in inches"""
    return {
        'left': emu_to_inches(shape.left),
        'top': emu_to_inches(shape.top),
        'width': emu_to_inches(shape.width),
        'height': emu_to_inches(shape.height),
        'right': emu_to_inches(shape.left + shape.width),
        'bottom': emu_to_inches(shape.top + shape.height)
    }

def check_shape_overlap(shapes):
    """Check for unintended shape overlaps"""
    overlaps = []
    for i, shape1 in enumerate(shapes):
        if not hasattr(shape1, 'left'):
            continue
        bounds1 = get_shape_bounds(shape1)
        
        for j, shape2 in enumerate(shapes[i+1:], start=i+1):
            if not hasattr(shape2, 'left'):
                continue
            bounds2 = get_shape_bounds(shape2)
            
            # Check if rectangles overlap
            if not (bounds1['right'] <= bounds2['left'] or
                   bounds1['left'] >= bounds2['right'] or
                   bounds1['bottom'] <= bounds2['top'] or
                   bounds1['top'] >= bounds2['bottom']):
                overlaps.append({
                    'shape1': {'name': shape1.name, 'index': i, 'bounds': bounds1},
                    'shape2': {'name': shape2.name, 'index': j, 'bounds': bounds2}
                })
    
    return overlaps

def check_text_on_light_background(shape):
    """Check for white text on light background"""
    issues = []
    
    if not shape.has_text_frame:
        return issues
    
    # Check shape fill color
    bg_color = None
    if shape.fill.type == 1:  # SOLID
        try:
            bg_color = shape.fill.fore_color.rgb
        except:
            pass
    
    # Check text color
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                text_color = run.font.color.rgb
                if text_color == WHITE or (text_color.r > 250 and text_color.g > 250 and text_color.b > 250):
                    if bg_color in LIGHT_BG or (bg_color and bg_color.r > 240 and bg_color.g > 240 and bg_color.b > 240):
                        issues.append({
                            'shape': shape.name,
                            'text_color': f"RGB({text_color.r},{text_color.g},{text_color.b})",
                            'bg_color': f"RGB({bg_color.r},{bg_color.g},{bg_color.b})" if bg_color else "Unknown",
                            'text_preview': run.text[:50]
                        })
            except:
                pass
    
    return issues

def check_subtitle_constraints(slide, slide_idx):
    """Check subtitle (중제목) design constraints"""
    issues = []
    
    # Find subtitle shapes (typically TextBox with specific positioning)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text = shape.text_frame.text.strip()
        
        # Heuristic: subtitles are usually in top area (0.5" - 1.5")
        top_inches = emu_to_inches(shape.top)
        if 0.4 < top_inches < 1.6 and len(text) > 0:
            # Check line count
            lines = text.split('\n')
            line_count = len([l for l in lines if l.strip()])
            
            if line_count > 2:
                issues.append({
                    'type': 'subtitle_too_many_lines',
                    'shape': shape.name,
                    'lines': line_count,
                    'text': text[:100]
                })
            
            # Check for word breaks (단어 중간 줄바꿈)
            for line in lines:
                if line.endswith('/') or (len(lines) > 1 and any(line.startswith(c) for c in ['row', 'umns', '처'])):
                    issues.append({
                        'type': 'subtitle_word_break',
                        'shape': shape.name,
                        'text': text
                    })
    
    return issues

def check_text_overflow(shape):
    """Check if text overflows shape boundaries"""
    if not shape.has_text_frame:
        return None
    
    tf = shape.text_frame
    text = tf.text.strip()
    
    if not text:
        return None
    
    # Estimate text height (rough approximation)
    total_lines = 0
    max_font_size = 0
    
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                max_font_size = max(max_font_size, run.font.size.pt)
        
        # Count lines based on text length and shape width
        text_in_para = paragraph.text
        if text_in_para.strip():
            total_lines += max(1, len(text_in_para) // 50)  # Rough estimate
    
    if max_font_size == 0:
        max_font_size = 13  # Default
    
    line_height = max_font_size * 1.2  # Typical line height
    estimated_height = total_lines * line_height
    shape_height = emu_to_inches(shape.height) * 72  # Convert to points
    
    if estimated_height > shape_height * 1.1:  # 10% tolerance
        return {
            'shape': shape.name,
            'estimated_height_pt': estimated_height,
            'shape_height_pt': shape_height,
            'overflow_ratio': estimated_height / shape_height if shape_height > 0 else 0,
            'text_preview': text[:100]
        }
    
    return None

def check_content_alignment(shapes):
    """Check if content is horizontally centered"""
    issues = []
    
    # Group shapes by vertical position (row)
    rows = {}
    for shape in shapes:
        if not hasattr(shape, 'top'):
            continue
        
        top_inches = emu_to_inches(shape.top)
        row_key = round(top_inches * 2) / 2  # Group by 0.5" increments
        
        if row_key not in rows:
            rows[row_key] = []
        rows[row_key].append(shape)
    
    # Check each row for centering
    slide_center = emu_to_inches(inches_to_emu(SLIDE_WIDTH / 2))
    
    for row_key, row_shapes in rows.items():
        if len(row_shapes) < 2:  # Need at least 2 shapes to check alignment
            continue
        
        # Calculate bounding box of all shapes in row
        min_left = min(emu_to_inches(s.left) for s in row_shapes if hasattr(s, 'left'))
        max_right = max(emu_to_inches(s.left + s.width) for s in row_shapes if hasattr(s, 'left'))
        
        row_center = (min_left + max_right) / 2
        left_margin = min_left
        right_margin = emu_to_inches(inches_to_emu(SLIDE_WIDTH)) - max_right
        
        # Check if centered (±0.1")
        if abs(row_center - slide_center) > 0.1:
            issues.append({
                'row_top': f"{row_key:.2f}\"",
                'row_center': f"{row_center:.2f}\"",
                'slide_center': f"{slide_center:.2f}\"",
                'offset': f"{abs(row_center - slide_center):.2f}\"",
                'left_margin': f"{left_margin:.2f}\"",
                'right_margin': f"{right_margin:.2f}\""
            })
    
    return issues

def check_font_consistency(shape):
    """Check font consistency against style guide"""
    issues = []
    
    if not shape.has_text_frame:
        return issues
    
    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        for run_idx, run in enumerate(paragraph.runs):
            font_name = run.font.name
            font_size = run.font.size.pt if run.font.size else None
            
            # Check for expected fonts
            if font_name and font_name not in ['Freesentation', '프리젠테이션 7 Bold', 'Calibri', 'Arial']:
                issues.append({
                    'shape': shape.name,
                    'para': para_idx,
                    'run': run_idx,
                    'unexpected_font': font_name,
                    'text': run.text[:30]
                })
            
            # Check size expectations (rough heuristic)
            if font_size:
                if font_size < 8 or font_size > 72:
                    issues.append({
                        'shape': shape.name,
                        'para': para_idx,
                        'run': run_idx,
                        'unusual_size': font_size,
                        'text': run.text[:30]
                    })
    
    return issues

def review_presentation(pptx_path):
    """Main review function"""
    prs = Presentation(pptx_path)
    
    results = {
        'file': str(pptx_path),
        'total_slides': len(prs.slides),
        'slide_size': {
            'width': f"{emu_to_inches(prs.slide_width):.3f}\"",
            'height': f"{emu_to_inches(prs.slide_height):.3f}\""
        },
        'issues': {
            'critical': [],
            'warning': [],
            'info': []
        },
        'slide_details': []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            'index': slide_idx,
            'layout': slide.slide_layout.name,
            'shape_count': len(slide.shapes),
            'checks': {
                'overlaps': [],
                'light_bg_white_text': [],
                'subtitle_issues': [],
                'text_overflow': [],
                'alignment_issues': [],
                'font_issues': []
            }
        }
        
        # Run all checks
        slide_data['checks']['overlaps'] = check_shape_overlap(slide.shapes)
        slide_data['checks']['subtitle_issues'] = check_subtitle_constraints(slide, slide_idx)
        slide_data['checks']['alignment_issues'] = check_content_alignment(slide.shapes)
        
        for shape in slide.shapes:
            # Light background + white text
            light_bg_issues = check_text_on_light_background(shape)
            slide_data['checks']['light_bg_white_text'].extend(light_bg_issues)
            
            # Text overflow
            overflow = check_text_overflow(shape)
            if overflow:
                slide_data['checks']['text_overflow'].append(overflow)
            
            # Font consistency
            font_issues = check_font_consistency(shape)
            slide_data['checks']['font_issues'].extend(font_issues)
            
            # Check if content is within body area
            if hasattr(shape, 'top') and hasattr(shape, 'height'):
                bottom = emu_to_inches(shape.top + shape.height)
                if bottom > emu_to_inches(inches_to_emu(BODY_LIMIT_Y)):
                    results['issues']['warning'].append({
                        'slide': slide_idx,
                        'type': 'content_below_limit',
                        'shape': shape.name,
                        'bottom': f"{bottom:.2f}\"",
                        'limit': "7.0\""
                    })
        
        # Categorize issues
        if slide_data['checks']['light_bg_white_text']:
            for issue in slide_data['checks']['light_bg_white_text']:
                results['issues']['critical'].append({
                    'slide': slide_idx,
                    'type': 'white_text_light_bg',
                    **issue
                })
        
        if slide_data['checks']['subtitle_issues']:
            for issue in slide_data['checks']['subtitle_issues']:
                results['issues']['critical'].append({
                    'slide': slide_idx,
                    **issue
                })
        
        if slide_data['checks']['text_overflow']:
            for issue in slide_data['checks']['text_overflow']:
                results['issues']['warning'].append({
                    'slide': slide_idx,
                    'type': 'text_overflow',
                    **issue
                })
        
        if slide_data['checks']['overlaps']:
            for issue in slide_data['checks']['overlaps']:
                results['issues']['warning'].append({
                    'slide': slide_idx,
                    'type': 'shape_overlap',
                    **issue
                })
        
        if slide_data['checks']['alignment_issues']:
            for issue in slide_data['checks']['alignment_issues']:
                results['issues']['info'].append({
                    'slide': slide_idx,
                    'type': 'alignment',
                    **issue
                })
        
        results['slide_details'].append(slide_data)
    
    # Calculate quality score
    critical_count = len(results['issues']['critical'])
    warning_count = len(results['issues']['warning'])
    info_count = len(results['issues']['info'])
    
    # Scoring: 100 - (critical * 10) - (warning * 3) - (info * 1)
    quality_score = max(0, 100 - (critical_count * 10) - (warning_count * 3) - (info_count * 1))
    results['quality_score'] = quality_score
    
    return results

if __name__ == '__main__':
    if len(sys.argv) != 2:
        print("Usage: python review_pptx_detailed.py <pptx_file>")
        sys.exit(1)
    
    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    results = review_presentation(pptx_path)
    print(json.dumps(results, indent=2, ensure_ascii=False))
