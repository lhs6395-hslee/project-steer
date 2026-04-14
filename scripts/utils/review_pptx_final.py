#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 종합 리뷰 스크립트 - 최종 버전
Sprint Contract의 모든 검증 항목 수행
"""

import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from collections import defaultdict
import json
import re

# 스타일 가이드 상수
COLORS = {
    'PRIMARY': (0, 67, 218),
    'DARK_GRAY': (33, 33, 33),
    'WHITE': (255, 255, 255),
    'ORANGE': (238, 129, 80),
    'GREEN': (76, 184, 143),
    'BG_BOX': (248, 249, 250),
    'BG_WHITE': (255, 255, 255),
}

results = {
    'file': '',
    'total_slides': 0,
    'issues': [],
    'summary': defaultdict(int),
    'constraint_violations': [],
}

def add_issue(severity, category, slide_idx, message, details=None):
    """이슈 기록 및 제약사항 위반 추적"""
    issue = {
        'severity': severity,
        'category': category,
        'slide_index': slide_idx,
        'message': message,
        'details': details or {}
    }
    results['issues'].append(issue)
    results['summary'][severity] += 1
    
    # CRITICAL 이슈는 제약사항 위반으로 간주
    if severity == 'CRITICAL':
        results['constraint_violations'].append(message)

def emu_to_inches(emu):
    return emu / 914400

def color_match(rgb1, rgb2, tolerance=10):
    """색상 매칭"""
    if not rgb1 or not rgb2:
        return False
    try:
        return all(abs(a - b) <= tolerance for a, b in zip(rgb1[:3], rgb2[:3]))
    except:
        return False

def get_rgb(color_obj):
    """색상 객체에서 RGB 추출"""
    try:
        return (color_obj[0], color_obj[1], color_obj[2])
    except:
        return None

def main():
    if len(sys.argv) < 2:
        print("Usage: python review_pptx_final.py <pptx_file>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    results['file'] = pptx_file
    
    print(f"=" * 80)
    print(f"PPTX 종합 리뷰 시작: {pptx_file}")
    print(f"=" * 80)
    
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        print(f"ERROR: 파일 열기 실패: {e}")
        sys.exit(1)
    
    results['total_slides'] = len(prs.slides)
    results['slide_width'] = emu_to_inches(prs.slide_width)
    results['slide_height'] = emu_to_inches(prs.slide_height)
    
    print(f"✓ Step 1: 파일 정상 열림 - {results['total_slides']} 슬라이드")
    add_issue('PASS', 'basic_info', -1, f"파일 정상 열림: {results['total_slides']} 슬라이드")
    
    # Step 2: 표지 검증
    print(f"✓ Step 2: 표지 슬라이드 검증")
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        title_found = subtitle_found = date_found = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if any(kw in text for kw in ['MSK', 'AWS', 'Kafka', '전문가']):
                    title_found = True
                if '/' in text and len(text) <= 15:
                    if re.match(r'\d{1,2}/\d{1,2}', text):
                        date_found = True
                        add_issue('PASS', 'cover', 0, f"날짜 형식 OK: {text}")
        
        if title_found:
            add_issue('PASS', 'cover', 0, "표지 제목 발견")
        else:
            add_issue('WARNING', 'cover', 0, "표지 제목을 찾을 수 없습니다")
    
    # Step 3: 목차 검증
    print(f"✓ Step 3: 목차 슬라이드 검증")
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        found = any('CONTENTS' in shape.text_frame.text.upper() 
                   for shape in slide.shapes if shape.has_text_frame)
        if found:
            add_issue('PASS', 'toc', 1, "목차 헤더 존재")
    
    # Step 4-6: 본문 슬라이드 검증
    print(f"✓ Step 4-6: 본문 슬라이드 레이아웃/중제목/타이틀 검증")
    for idx in range(2, len(prs.slides) - 1):
        slide = prs.slides[idx]
        
        # 중제목 검증 (헤더 영역의 설명 텍스트)
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'top'):
                top_in = emu_to_inches(shape.top)
                text = shape.text_frame.text.strip()
                
                # 헤더 영역 (0.5" ~ 1.5")
                if 0.5 <= top_in <= 1.5 and len(text) > 20:
                    lines = [l for l in text.split('\n') if l.strip()]
                    if len(lines) > 2:
                        add_issue('CRITICAL', 'subtitle', idx,
                                  f"중제목이 3줄 이상: {len(lines)}줄 (최대 2줄)",
                                  {'text': text[:100]})
                
                # 타이틀 줄바꿈 검증
                if re.match(r'^\d+-\d+\.', text):
                    # 단어 중간 줄바꿈 검출
                    lines = text.split('\n')
                    for i, line in enumerate(lines[:-1]):
                        if line and lines[i+1]:
                            # 영문 단어가 잘렸는지 확인
                            last_word = line.split()[-1] if line.split() else ''
                            next_first = lines[i+1].split()[0] if lines[i+1].split() else ''
                            if (last_word and next_first and 
                                last_word[-1].isalpha() and next_first[0].islower()):
                                add_issue('WARNING', 'title', idx,
                                          f"타이틀 줄바꿈이 단어 중간: '{last_word}' / '{next_first}'")
    
    # Step 7: 아이콘 검증 - 파란색 원형 fallback 검출 (CRITICAL)
    print(f"✓ Step 7: 아이콘 검증")
    for idx, slide in enumerate(prs.slides):
        blue_ovals = []
        real_icons = []
        
        for shape in slide.shapes:
            # 실제 PNG 아이콘
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width_in = emu_to_inches(shape.width)
                if 0.35 <= width_in <= 0.75:
                    real_icons.append(shape)
            
            # 파란색 원형 fallback 검출
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if hasattr(shape, 'auto_shape_type'):
                        if shape.auto_shape_type == MSO_SHAPE.OVAL:
                            # 배경 색상 확인
                            if shape.fill.type == 1:  # SOLID
                                try:
                                    rgb = shape.fill.fore_color.rgb
                                    color = (rgb[0], rgb[1], rgb[2])
                                    if color_match(color, COLORS['PRIMARY']):
                                        blue_ovals.append(shape)
                                except:
                                    pass
                except:
                    pass
        
        # CRITICAL: 파란색 원형 fallback 발견
        if blue_ovals:
            add_issue('CRITICAL', 'icon', idx,
                      f"파란색 원형 fallback 아이콘 발견: {len(blue_ovals)}개 (절대 금지)")
        elif real_icons:
            add_issue('PASS', 'icon', idx,
                      f"실제 PNG 아이콘 사용: {len(real_icons)}개")
    
    # Step 8: 배치 및 여백 검증
    print(f"✓ Step 8: 콘텐츠 배치 및 여백 검증")
    for idx in range(2, len(prs.slides) - 1):
        slide = prs.slides[idx]
        shapes_in_body = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                top_in = emu_to_inches(shape.top)
                if 2.0 <= top_in <= 7.5:
                    shapes_in_body.append(shape)
        
        if shapes_in_body:
            # 7.0" 초과 확인
            bottommost = max(s.top + s.height for s in shapes_in_body)
            bottom_in = emu_to_inches(bottommost)
            
            if bottom_in > 7.0:
                add_issue('WARNING', 'placement', idx,
                          f"콘텐츠가 7.0\" 초과: {bottom_in:.2f}\" (권장: 7.0\" 이내)")
            
            # 좌우 여백 대칭 확인
            leftmost = min(s.left for s in shapes_in_body)
            rightmost = max(s.left + s.width for s in shapes_in_body)
            left_margin = emu_to_inches(leftmost)
            right_margin = emu_to_inches(prs.slide_width - rightmost)
            
            if abs(left_margin - right_margin) > 0.2:
                add_issue('WARNING', 'placement', idx,
                          f"좌우 여백 비대칭: 좌 {left_margin:.2f}\" vs 우 {right_margin:.2f}\"")
    
    # Step 9: 텍스트 오버플로우 검증
    print(f"✓ Step 9: 텍스트 오버플로우 검증")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            text = tf.text.strip()
            if len(text) < 20:
                continue
            
            # 간단한 오버플로우 추정
            line_count = len([p for p in tf.paragraphs if p.text.strip()])
            
            # 과도한 텍스트 검출
            if line_count > 15:
                add_issue('WARNING', 'overflow', idx,
                          f"텍스트 줄 수 과다: {line_count}줄",
                          {'preview': text[:80]})
    
    # Step 11: 색상 스킴 검증 - 흰 텍스트 on 밝은 배경 (CRITICAL)
    print(f"✓ Step 11: 색상 스킴 검증")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # 배경 색상 확인
            bg_color = None
            try:
                if shape.fill.type == 1:  # SOLID
                    rgb = shape.fill.fore_color.rgb
                    bg_color = (rgb[0], rgb[1], rgb[2])
            except:
                pass
            
            # 밝은 배경 검출
            if bg_color and (color_match(bg_color, COLORS['BG_BOX']) or
                              color_match(bg_color, COLORS['BG_WHITE'])):
                # 텍스트 색상 확인
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        try:
                            if run.font.color.type == 1:  # RGB
                                text_rgb = run.font.color.rgb
                                text_color = (text_rgb[0], text_rgb[1], text_rgb[2])
                                if color_match(text_color, COLORS['WHITE']):
                                    add_issue('CRITICAL', 'color', idx,
                                              "흰 텍스트 on 밝은 배경 발견 (가독성 불가)",
                                              {'bg': bg_color, 'text': text_color})
                        except:
                            pass
    
    # Step 13: 끝맺음 검증
    print(f"✓ Step 13: 끝맺음 슬라이드 검증")
    if len(prs.slides) > 0:
        slide = prs.slides[-1]
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if 'Thank You' in text or '감사합니다' in text:
                    found = True
                    add_issue('PASS', 'ending', len(prs.slides)-1, "'Thank You' 텍스트 존재")
                    break
        if not found:
            add_issue('WARNING', 'ending', len(prs.slides)-1,
                      "끝맺음 텍스트를 찾을 수 없습니다")
    
    # Step 14: 보고서 생성
    print(f"✓ Step 14: 종합 보고서 생성")
    
    critical = [i for i in results['issues'] if i['severity'] == 'CRITICAL']
    warning = [i for i in results['issues'] if i['severity'] == 'WARNING']
    passed = [i for i in results['issues'] if i['severity'] == 'PASS']
    
    report = []
    report.append("=" * 80)
    report.append("PPTX 종합 리뷰 보고서")
    report.append("=" * 80)
    report.append(f"파일: {results['file']}")
    report.append(f"총 슬라이드: {results['total_slides']}")
    report.append(f"슬라이드 크기: {results['slide_width']:.2f}\" × {results['slide_height']:.2f}\"")
    report.append("")
    
    report.append("=" * 80)
    report.append("검증 결과 요약")
    report.append("=" * 80)
    report.append(f"✗ CRITICAL: {len(critical)} 건 (즉시 수정 필요)")
    report.append(f"⚠ WARNING:  {len(warning)} 건 (검토 권장)")
    report.append(f"✓ PASS:     {len(passed)} 건")
    report.append("")
    
    if critical:
        report.append("=" * 80)
        report.append("CRITICAL 이슈 (제약사항 위반)")
        report.append("=" * 80)
        for i, issue in enumerate(critical, 1):
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else '-'
            report.append(f"{i}. [슬라이드 {slide_num}] {issue['message']}")
            if issue['details']:
                for k, v in issue['details'].items():
                    report.append(f"   {k}: {v}")
        report.append("")
    
    if warning:
        report.append("=" * 80)
        report.append(f"WARNING 이슈 (최대 20개 표시)")
        report.append("=" * 80)
        for i, issue in enumerate(warning[:20], 1):
            slide_num = issue['slide_index'] + 1 if issue['slide_index'] >= 0 else '-'
            report.append(f"{i}. [슬라이드 {slide_num}] {issue['message']}")
        if len(warning) > 20:
            report.append(f"... 외 {len(warning) - 20}건")
        report.append("")
    
    total_checks = len(results['issues'])
    pass_count = len(passed)
    score = (pass_count / total_checks * 100) if total_checks > 0 else 0
    
    report.append("=" * 80)
    report.append(f"전체 품질 점수: {score:.1f}% ({pass_count}/{total_checks} 항목 통과)")
    report.append("=" * 80)
    report.append("")
    
    report.append("개선 권고사항:")
    if critical:
        report.append(f"1. CRITICAL 이슈 {len(critical)}건을 최우선으로 해결하세요.")
        for issue in critical[:3]:
            report.append(f"   • {issue['message']}")
    if warning:
        report.append(f"2. WARNING 이슈 {len(warning)}건을 검토하세요.")
    if not critical and not warning:
        report.append("✓ 모든 검증 항목을 통과했습니다!")
    report.append("")
    
    report_text = "\n".join(report)
    print("\n" + report_text)
    
    # 결과 파일 저장
    output_json = pptx_file.replace('.pptx', '_review.json')
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    output_txt = pptx_file.replace('.pptx', '_review.txt')
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(report_text)
    
    print(f"\n결과 파일:")
    print(f"  • JSON: {output_json}")
    print(f"  • TXT:  {output_txt}")
    
    # 종료 코드
    if critical:
        print(f"\n❌ 검증 실패: CRITICAL 이슈 {len(critical)}건")
        sys.exit(2)
    elif len(warning) > 10:
        print(f"\n⚠️  검증 경고: WARNING 이슈 {len(warning)}건")
        sys.exit(1)
    else:
        print(f"\n✅ 검증 통과")
        sys.exit(0)

if __name__ == '__main__':
    main()
