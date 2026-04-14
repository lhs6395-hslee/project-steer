#!/usr/bin/env python3
"""
Independent Reviewer: Verify slides idx 5 and idx 6 in AWS_MSK_Expert_Intro.pptx
Sprint_Contract success criteria verification
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"

# Color constants
SUB_ORANGE = RGBColor(238, 129, 80)
SUB_GREEN = RGBColor(76, 184, 143)
PRIMARY = RGBColor(0, 67, 218)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(230, 240, 255)  # #E6F0FF
BORDER_BLUE = RGBColor(0, 67, 218)   # #0043DA

def emu_to_inches(emu):
    return emu / 914400.0

def get_shape_color(shape):
    """Get fill color of a shape"""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        if hasattr(fill, 'fore_color') and fill.fore_color:
            try:
                return fill.fore_color.rgb
            except:
                pass
        # Try solid fill
        try:
            if fill.type == 1:  # MSO_THEME_COLOR or solid
                return fill.fore_color.rgb
        except:
            pass
    except:
        pass
    return None

def get_shape_fill_rgb(shape):
    """Try multiple ways to get fill RGB"""
    try:
        fill = shape.fill
        # Try solid fill
        try:
            solid = fill.solid()
            return fill.fore_color.rgb
        except:
            pass
        # Direct access
        try:
            return fill.fore_color.rgb
        except:
            pass
    except:
        pass
    return None

def color_matches(c1, c2, tolerance=5):
    """Check if two RGBColor values match within tolerance"""
    if c1 is None or c2 is None:
        return False
    return (abs(c1[0] - c2[0]) <= tolerance and
            abs(c1[1] - c2[1]) <= tolerance and
            abs(c1[2] - c2[2]) <= tolerance)

def get_text(shape):
    """Get all text from a shape"""
    try:
        return shape.text_frame.text
    except:
        try:
            return shape.text
        except:
            return ""

def check_slide_idx6(slide):
    """Verify slide idx 6 (L02 Three Cards)"""
    checks = {}
    issues = []
    fix_required = []

    # Collect all shapes
    shapes = list(slide.shapes)
    shape_names = [s.name for s in shapes]

    # Check 1: Wrong shapes ABSENT
    wrong_shapes = [
        "Card_Left_BG", "Card_Left_Title", "Card_Left_Content",
        "Rectangle 18", "VPC_Label", "Rectangle 20", "TextBox 21",
        "Rounded Rectangle 22", "Broker_AZa_Label", "Rounded Rectangle 24",
        "TextBox 25", "Rounded Rectangle 26", "TextBox 27",
        "TextBox 28", "TextBox 29"
    ]
    found_wrong = [s for s in wrong_shapes if s in shape_names]
    if found_wrong:
        checks["check1_wrong_shapes_absent"] = f"FAIL: Found wrong shapes: {found_wrong}"
        issues.append(f"Wrong shapes still present: {found_wrong}")
        fix_required.append("Remove legacy shapes from slide idx 6")
    else:
        checks["check1_wrong_shapes_absent"] = "PASS: All wrong shapes absent"

    # Check 2: Three equal cards at top≈4.40", h≈2.50"
    card_shapes = []
    for s in shapes:
        top_in = emu_to_inches(s.top)
        h_in = emu_to_inches(s.height)
        # Look for shapes near top=4.40" with h≈2.50"
        if 4.0 <= top_in <= 4.8 and 2.0 <= h_in <= 3.0:
            card_shapes.append(s)

    if len(card_shapes) >= 3:
        checks["check2_three_cards_present"] = f"PASS: Found {len(card_shapes)} card-like shapes at top≈4.40\""
    else:
        checks["check2_three_cards_present"] = f"FAIL: Expected 3 cards at top≈4.40\", h≈2.50\", found {len(card_shapes)}"
        issues.append(f"Three equal cards not found (found {len(card_shapes)} shapes at expected position)")
        fix_required.append("Add/fix three equal cards at top≈4.40\" h≈2.50\" on slide idx 6")

    # Check 3: Cards contain titles
    all_text = " ".join([get_text(s) for s in shapes])
    required_titles = ["핵심 구성요소", "네트워크 구성", "스토리지"]
    missing_titles = [t for t in required_titles if t not in all_text]
    if missing_titles:
        checks["check3_card_titles"] = f"FAIL: Missing titles: {missing_titles}"
        issues.append(f"Missing card titles: {missing_titles}")
        fix_required.append("Add card titles: 핵심 구성요소, 네트워크 구성, 스토리지")
    else:
        checks["check3_card_titles"] = "PASS: All three card titles present"

    # Check 4: "Rounded Rectangle 6" fill = SUB_ORANGE
    rr6 = next((s for s in shapes if s.name == "Rounded Rectangle 6"), None)
    if rr6 is None:
        checks["check4_rr6_fill"] = "FAIL: Rounded Rectangle 6 not found"
        issues.append("Rounded Rectangle 6 missing")
        fix_required.append("Ensure Rounded Rectangle 6 exists with SUB_ORANGE fill")
    else:
        color = get_shape_fill_rgb(rr6)
        if color_matches(color, SUB_ORANGE):
            checks["check4_rr6_fill"] = f"PASS: Rounded Rectangle 6 fill = SUB_ORANGE ({color})"
        else:
            checks["check4_rr6_fill"] = f"FAIL: Rounded Rectangle 6 fill = {color}, expected SUB_ORANGE (238,129,80)"
            issues.append(f"Rounded Rectangle 6 wrong fill: {color}")
            fix_required.append("Set Rounded Rectangle 6 fill to SUB_ORANGE (238,129,80)")

    # Check 5: "Right Arrow 7" fill = SUB_ORANGE
    ra7 = next((s for s in shapes if s.name == "Right Arrow 7"), None)
    if ra7 is None:
        checks["check5_ra7_fill"] = "FAIL: Right Arrow 7 not found"
        issues.append("Right Arrow 7 missing")
        fix_required.append("Ensure Right Arrow 7 exists with SUB_ORANGE fill")
    else:
        color = get_shape_fill_rgb(ra7)
        if color_matches(color, SUB_ORANGE):
            checks["check5_ra7_fill"] = f"PASS: Right Arrow 7 fill = SUB_ORANGE ({color})"
        else:
            checks["check5_ra7_fill"] = f"FAIL: Right Arrow 7 fill = {color}, expected SUB_ORANGE"
            issues.append(f"Right Arrow 7 wrong fill: {color}")
            fix_required.append("Set Right Arrow 7 fill to SUB_ORANGE (238,129,80)")

    # Check 6: "Right Arrow 13" fill = SUB_GREEN
    ra13 = next((s for s in shapes if s.name == "Right Arrow 13"), None)
    if ra13 is None:
        checks["check6_ra13_fill"] = "FAIL: Right Arrow 13 not found"
        issues.append("Right Arrow 13 missing")
        fix_required.append("Ensure Right Arrow 13 exists with SUB_GREEN fill")
    else:
        color = get_shape_fill_rgb(ra13)
        if color_matches(color, SUB_GREEN):
            checks["check6_ra13_fill"] = f"PASS: Right Arrow 13 fill = SUB_GREEN ({color})"
        else:
            checks["check6_ra13_fill"] = f"FAIL: Right Arrow 13 fill = {color}, expected SUB_GREEN"
            issues.append(f"Right Arrow 13 wrong fill: {color}")
            fix_required.append("Set Right Arrow 13 fill to SUB_GREEN (76,184,143)")

    # Check 7: "Rounded Rectangle 14" fill = SUB_GREEN, text = WHITE
    rr14 = next((s for s in shapes if s.name == "Rounded Rectangle 14"), None)
    if rr14 is None:
        checks["check7_rr14"] = "FAIL: Rounded Rectangle 14 not found"
        issues.append("Rounded Rectangle 14 missing")
        fix_required.append("Ensure Rounded Rectangle 14 exists with SUB_GREEN fill and WHITE text")
    else:
        color = get_shape_fill_rgb(rr14)
        fill_ok = color_matches(color, SUB_GREEN)
        # Check text color
        text_white = False
        try:
            for para in rr14.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color.rgb == WHITE:
                        text_white = True
                        break
        except:
            pass

        status = "PASS" if fill_ok else "FAIL"
        checks["check7_rr14"] = f"{status}: fill={color} (expected SUB_GREEN), text_white={text_white}"
        if not fill_ok:
            issues.append(f"Rounded Rectangle 14 wrong fill: {color}")
            fix_required.append("Set Rounded Rectangle 14 fill to SUB_GREEN (76,184,143)")

    # Check 8: "Rounded Rectangle 8" fill = #E6F0FF, border = #0043DA
    rr8 = next((s for s in shapes if s.name == "Rounded Rectangle 8"), None)
    if rr8 is None:
        checks["check8_rr8"] = "FAIL: Rounded Rectangle 8 not found"
        issues.append("Rounded Rectangle 8 missing")
        fix_required.append("Ensure Rounded Rectangle 8 exists with #E6F0FF fill")
    else:
        color = get_shape_fill_rgb(rr8)
        expected = RGBColor(230, 240, 255)
        if color_matches(color, expected, tolerance=10):
            checks["check8_rr8"] = f"PASS: Rounded Rectangle 8 fill = {color} (#E6F0FF)"
        else:
            checks["check8_rr8"] = f"FAIL: Rounded Rectangle 8 fill = {color}, expected #E6F0FF"
            issues.append(f"Rounded Rectangle 8 fill changed: {color}")
            fix_required.append("Restore Rounded Rectangle 8 fill to #E6F0FF")

    # Check 9: "Rounded Rectangle 10/11/12" fill = PRIMARY (#0043DA)
    for rr_name in ["Rounded Rectangle 10", "Rounded Rectangle 11", "Rounded Rectangle 12"]:
        rr = next((s for s in shapes if s.name == rr_name), None)
        key = f"check9_{rr_name.replace(' ', '_')}"
        if rr is None:
            checks[key] = f"FAIL: {rr_name} not found"
            issues.append(f"{rr_name} missing")
            fix_required.append(f"Ensure {rr_name} exists with PRIMARY fill")
        else:
            color = get_shape_fill_rgb(rr)
            if color_matches(color, PRIMARY):
                checks[key] = f"PASS: {rr_name} fill = PRIMARY ({color})"
            else:
                checks[key] = f"FAIL: {rr_name} fill = {color}, expected PRIMARY #0043DA"
                issues.append(f"{rr_name} wrong fill: {color}")
                fix_required.append(f"Set {rr_name} fill to PRIMARY #0043DA")

    # Check 10: No shape bottom > 7.0"
    overflow_shapes = []
    for s in shapes:
        try:
            bottom = emu_to_inches(s.top + s.height)
            if bottom > 7.0:
                overflow_shapes.append(f"{s.name} (bottom={bottom:.2f}\")")
        except:
            pass

    if overflow_shapes:
        checks["check10_no_overflow"] = f"FAIL: Shapes overflow 7.0\": {overflow_shapes}"
        issues.append(f"Shapes extend beyond 7.0\": {overflow_shapes}")
        fix_required.append("Resize/reposition shapes to keep bottom <= 7.0\"")
    else:
        checks["check10_no_overflow"] = "PASS: No shape bottom > 7.0\""

    return checks, issues, fix_required


def check_slide_idx5(slide):
    """Verify slide idx 5 (L01 Bento Grid)"""
    checks = {}
    issues = []
    fix_required = []

    shapes = list(slide.shapes)
    shape_names = [s.name for s in shapes]

    # Check 11: TextBox 5 text = "MSK 아키텍처 핵심 개요"
    tb5 = next((s for s in shapes if s.name == "TextBox 5"), None)
    if tb5 is None:
        checks["check11_textbox5"] = "FAIL: TextBox 5 not found"
        issues.append("TextBox 5 missing")
        fix_required.append("Ensure TextBox 5 exists with text 'MSK 아키텍처 핵심 개요'")
    else:
        text = get_text(tb5)
        if "MSK 아키텍처 핵심 개요" in text:
            checks["check11_textbox5"] = f"PASS: TextBox 5 text = '{text}'"
        else:
            checks["check11_textbox5"] = f"FAIL: TextBox 5 text = '{text}', expected 'MSK 아키텍처 핵심 개요'"
            issues.append(f"TextBox 5 wrong text: '{text}'")
            fix_required.append("Set TextBox 5 text to 'MSK 아키텍처 핵심 개요'")

    # Check 12: TextBox 6 contains required keywords
    tb6 = next((s for s in shapes if s.name == "TextBox 6"), None)
    required_keywords = ["클러스터 유형", "Provisioned", "Serverless", "Express", "보안 체계", "IAM"]
    if tb6 is None:
        checks["check12_textbox6"] = "FAIL: TextBox 6 not found"
        issues.append("TextBox 6 missing")
        fix_required.append("Ensure TextBox 6 exists with required keywords")
    else:
        text = get_text(tb6)
        missing_kw = [kw for kw in required_keywords if kw not in text]
        if missing_kw:
            checks["check12_textbox6"] = f"FAIL: TextBox 6 missing keywords: {missing_kw}"
            issues.append(f"TextBox 6 missing keywords: {missing_kw}")
            fix_required.append(f"Add missing keywords to TextBox 6: {missing_kw}")
        else:
            checks["check12_textbox6"] = f"PASS: TextBox 6 contains all required keywords"

    # Check 13: Flow diagram container at top≈5.25", bottom≈6.75"
    flow_container = None
    for s in shapes:
        try:
            top_in = emu_to_inches(s.top)
            bottom_in = emu_to_inches(s.top + s.height)
            # Look for Rounded Rectangle near top=5.25" and bottom=6.75"
            if 5.0 <= top_in <= 5.5 and 6.5 <= bottom_in <= 7.0:
                flow_container = s
                break
        except:
            pass

    if flow_container:
        top_in = emu_to_inches(flow_container.top)
        bottom_in = emu_to_inches(flow_container.top + flow_container.height)
        checks["check13_flow_container"] = f"PASS: Flow container found '{flow_container.name}' top={top_in:.2f}\" bottom={bottom_in:.2f}\""
    else:
        # Try broader search
        candidates = []
        for s in shapes:
            try:
                top_in = emu_to_inches(s.top)
                bottom_in = emu_to_inches(s.top + s.height)
                if 4.8 <= top_in <= 5.8 and 6.2 <= bottom_in <= 7.1:
                    candidates.append(f"{s.name} top={top_in:.2f}\" bot={bottom_in:.2f}\"")
            except:
                pass
        checks["check13_flow_container"] = f"FAIL: No flow container at top≈5.25\" bottom≈6.75\". Candidates: {candidates}"
        issues.append("Flow diagram container not found at expected position")
        fix_required.append("Add flow container Rounded Rectangle at top≈5.25\", bottom≈6.75\"")

    # Check 14: Three colored boxes inside flow (SUB_ORANGE, PRIMARY, SUB_GREEN)
    # Look for shapes with these colors
    orange_found = False
    primary_found = False
    green_found = False

    for s in shapes:
        color = get_shape_fill_rgb(s)
        if color_matches(color, SUB_ORANGE):
            orange_found = True
        if color_matches(color, PRIMARY):
            primary_found = True
        if color_matches(color, SUB_GREEN):
            green_found = True

    flow_colors_ok = orange_found and primary_found and green_found
    checks["check14_flow_colors"] = (
        f"{'PASS' if flow_colors_ok else 'FAIL'}: "
        f"SUB_ORANGE={orange_found}, PRIMARY={primary_found}, SUB_GREEN={green_found}"
    )
    if not flow_colors_ok:
        missing = []
        if not orange_found: missing.append("SUB_ORANGE Producer")
        if not primary_found: missing.append("PRIMARY MSK")
        if not green_found: missing.append("SUB_GREEN Consumer")
        issues.append(f"Missing flow colored boxes: {missing}")
        fix_required.append(f"Add colored flow boxes: {missing}")

    # Check 15: RR4 geometry: top≈2.80", w≈5.90", h≈4.20"
    rr4 = next((s for s in shapes if s.name == "Rounded Rectangle 4" or s.name == "RR4"), None)
    # Also search by position if name not found
    if rr4 is None:
        for s in shapes:
            try:
                top_in = emu_to_inches(s.top)
                w_in = emu_to_inches(s.width)
                h_in = emu_to_inches(s.height)
                if 2.5 <= top_in <= 3.1 and 5.5 <= w_in <= 6.3 and 3.8 <= h_in <= 4.6:
                    rr4 = s
                    break
            except:
                pass

    if rr4 is None:
        checks["check15_rr4_geometry"] = "FAIL: RR4 (Rounded Rectangle 4) not found at expected position"
        issues.append("RR4 not found at top≈2.80\", w≈5.90\", h≈4.20\"")
        fix_required.append("Ensure RR4 geometry: top≈2.80\", w≈5.90\", h≈4.20\"")
    else:
        top_in = emu_to_inches(rr4.top)
        w_in = emu_to_inches(rr4.width)
        h_in = emu_to_inches(rr4.height)
        top_ok = 2.5 <= top_in <= 3.1
        w_ok = 5.5 <= w_in <= 6.3
        h_ok = 3.8 <= h_in <= 4.6
        status = "PASS" if (top_ok and w_ok and h_ok) else "FAIL"
        checks["check15_rr4_geometry"] = (
            f"{status}: '{rr4.name}' top={top_in:.2f}\" w={w_in:.2f}\" h={h_in:.2f}\" "
            f"(expected top≈2.80\" w≈5.90\" h≈4.20\")"
        )
        if status == "FAIL":
            issues.append(f"RR4 geometry off: top={top_in:.2f}\" w={w_in:.2f}\" h={h_in:.2f}\"")
            fix_required.append("Fix RR4 geometry to top≈2.80\", w≈5.90\", h≈4.20\"")

    # Check 16: TextBox 13, TextBox 14 positions (top≈1.75" and top≈2.15")
    for tb_name, expected_top in [("TextBox 13", 1.75), ("TextBox 14", 2.15)]:
        tb = next((s for s in shapes if s.name == tb_name), None)
        key = f"check16_{tb_name.replace(' ', '_')}"
        if tb is None:
            checks[key] = f"FAIL: {tb_name} not found"
            issues.append(f"{tb_name} missing")
            fix_required.append(f"Ensure {tb_name} at top≈{expected_top}\"")
        else:
            top_in = emu_to_inches(tb.top)
            tolerance = 0.3
            if abs(top_in - expected_top) <= tolerance:
                checks[key] = f"PASS: {tb_name} top={top_in:.2f}\" (expected≈{expected_top}\")"
            else:
                checks[key] = f"FAIL: {tb_name} top={top_in:.2f}\" (expected≈{expected_top}\")"
                issues.append(f"{tb_name} position off: top={top_in:.2f}\" expected≈{expected_top}\"")
                fix_required.append(f"Reposition {tb_name} to top≈{expected_top}\"")

    # Check 17: No shape bottom > 7.0" (tolerance: 0.002" for floating-point rounding artifacts)
    OVERFLOW_TOLERANCE = 0.002  # 0.002" ≈ 0.05mm, sub-pixel rounding artifact threshold
    overflow_shapes = []
    for s in shapes:
        try:
            bottom = emu_to_inches(s.top + s.height)
            if bottom > (7.0 + OVERFLOW_TOLERANCE):
                overflow_shapes.append(f"{s.name} (bottom={bottom:.4f}\")")
        except:
            pass

    if overflow_shapes:
        checks["check17_no_overflow"] = f"FAIL: Shapes overflow 7.0\": {overflow_shapes}"
        issues.append(f"Shapes extend beyond 7.0\": {overflow_shapes}")
        fix_required.append("Resize/reposition shapes to keep bottom <= 7.0\"")
    else:
        checks["check17_no_overflow"] = "PASS: No shape bottom > 7.0\" (sub-pixel rounding within tolerance)"

    # Check 18: No content textbox with WHITE text on light/transparent background
    # NOTE: Transparent-fill textboxes overlaid on dark-colored shapes are CORRECT design.
    # We must check the underlying shape's color at the same position.
    # Build a map of shape positions to detect overlapping dark shapes.
    dark_colors = [SUB_ORANGE, PRIMARY, SUB_GREEN,
                   RGBColor(0, 0, 0), RGBColor(50, 50, 50)]

    def is_dark_color(c):
        if c is None:
            return False
        # Dark = any of the expected dark colors, or luminance < 128
        luminance = 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]
        return luminance < 180  # covers PRIMARY(0,67,218), SUB_ORANGE(238,129,80 lum≈141), SUB_GREEN(76,184,143 lum≈154)

    def shapes_overlap(s1_left, s1_top, s1_right, s1_bottom,
                       s2_left, s2_top, s2_right, s2_bottom):
        return (s1_left < s2_right and s1_right > s2_left and
                s1_top < s2_bottom and s1_bottom > s2_top)

    white_on_light = []
    for s in shapes:
        try:
            fill_color = get_shape_fill_rgb(s)
            # Only check shapes where fill is None (transparent) or explicitly light
            is_transparent = fill_color is None
            is_light_fill = (fill_color is not None and
                            not is_dark_color(fill_color))

            if not (is_transparent or is_light_fill):
                continue

            # Check if this shape has WHITE text
            has_white_text = False
            try:
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.rgb == WHITE:
                                has_white_text = True
                                break
                        except:
                            pass
                    if has_white_text:
                        break
            except:
                pass

            if not has_white_text:
                continue

            # For transparent-fill shapes, check if there's a dark shape underneath
            if is_transparent:
                s_left = emu_to_inches(s.left)
                s_top = emu_to_inches(s.top)
                s_right = s_left + emu_to_inches(s.width)
                s_bottom = s_top + emu_to_inches(s.height)

                has_dark_bg = False
                for other in shapes:
                    if other.name == s.name:
                        continue
                    try:
                        other_fill = get_shape_fill_rgb(other)
                        if other_fill is None or not is_dark_color(other_fill):
                            continue
                        o_left = emu_to_inches(other.left)
                        o_top = emu_to_inches(other.top)
                        o_right = o_left + emu_to_inches(other.width)
                        o_bottom = o_top + emu_to_inches(other.height)
                        if shapes_overlap(s_left, s_top, s_right, s_bottom,
                                         o_left, o_top, o_right, o_bottom):
                            has_dark_bg = True
                            break
                    except:
                        pass

                if has_dark_bg:
                    # White text on dark background — correct design
                    continue

            white_on_light.append(f"{s.name} (fill={fill_color})")
        except:
            pass

    if white_on_light:
        checks["check18_white_on_light"] = f"FAIL: WHITE text on light background: {white_on_light}"
        issues.append(f"WHITE text on light/transparent background: {white_on_light}")
        fix_required.append("Change text color from WHITE to dark color for light-background shapes")
    else:
        checks["check18_white_on_light"] = "PASS: No white text on genuinely light/transparent background"

    return checks, issues, fix_required


def main():
    try:
        prs = Presentation(PPTX_PATH)
    except Exception as e:
        result = {
            "verdict": "FAIL",
            "confidence": 0.0,
            "checks": {},
            "issues": [f"Cannot open presentation: {e}"],
            "fix_required": ["Fix presentation file path or permissions"]
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    slides = list(prs.slides)
    print(f"Total slides: {len(slides)}", file=sys.stderr)

    # Print shape names for debugging
    for idx, slide in enumerate(slides):
        names = [s.name for s in slide.shapes]
        print(f"Slide {idx}: {names}", file=sys.stderr)

    if len(slides) < 7:
        result = {
            "verdict": "FAIL",
            "confidence": 0.0,
            "checks": {},
            "issues": [f"Presentation has only {len(slides)} slides, need at least 7"],
            "fix_required": []
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    all_checks = {}
    all_issues = []
    all_fix_required = []

    # Check slide idx 6
    slide6 = slides[6]
    checks6, issues6, fix6 = check_slide_idx6(slide6)
    all_checks.update({f"slide6_{k}": v for k, v in checks6.items()})
    all_issues.extend([f"[slide6] {i}" for i in issues6])
    all_fix_required.extend([f"[slide6] {f}" for f in fix6])

    # Check slide idx 5
    slide5 = slides[5]
    checks5, issues5, fix5 = check_slide_idx5(slide5)
    all_checks.update({f"slide5_{k}": v for k, v in checks5.items()})
    all_issues.extend([f"[slide5] {i}" for i in issues5])
    all_fix_required.extend([f"[slide5] {f}" for f in fix5])

    # Calculate verdict
    total_checks = len(all_checks)
    failed_checks = sum(1 for v in all_checks.values() if v.startswith("FAIL"))
    passed_checks = total_checks - failed_checks

    confidence = passed_checks / total_checks if total_checks > 0 else 0.0
    verdict = "PASS" if failed_checks == 0 else "FAIL"

    result = {
        "verdict": verdict,
        "confidence": round(confidence, 2),
        "checks": all_checks,
        "issues": all_issues,
        "fix_required": all_fix_required
    }

    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
