#!/usr/bin/env python3
"""
Independent Reviewer: Verify fixes in AWS_MSK_Expert_Intro.pptx
Checks:
1. idx5 (L01) TextBox 5 & 6 internal padding
2. idx6 (L02) TextBox 2 line break
3. idx6 (L02) 3 card backgrounds have PRIMARY border #0043DA
4. No overflow on idx5, idx6
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Emu
from lxml import etree

PPTX_PATH = "/Users/toule/Documents/kiro/project-steer/results/pptx/AWS_MSK_Expert_Intro.pptx"
MIN_PADDING_EMU = 91440  # 0.1 inch in EMU
PRIMARY_COLOR = "0043DA"
SLIDE_HEIGHT_LIMIT = Inches(7.0)  # 7.0 inches in EMU

def emu_to_inches(emu):
    return emu / 914400.0

def get_bodyPr_padding(txBody):
    """Extract lIns, tIns from bodyPr element."""
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    bodyPr = txBody.find('.//a:bodyPr', nsmap)
    if bodyPr is None:
        return None, None
    lIns = bodyPr.get('lIns')
    tIns = bodyPr.get('tIns')
    # Default values per OOXML spec: lIns=91440 EMU, tIns=45720 EMU
    lIns_val = int(lIns) if lIns is not None else 91440
    tIns_val = int(tIns) if tIns is not None else 45720
    return lIns_val, tIns_val

def get_text_paragraphs(shape):
    """Return list of paragraph texts."""
    if not shape.has_text_frame:
        return []
    return [p.text for p in shape.text_frame.paragraphs]

def get_shape_color(shape):
    """Get line color from shape."""
    try:
        sp_tree = shape._element
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
        # spPr > ln > solidFill > srgbClr
        ln = sp_tree.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
        if ln is None:
            return None, None
        lnW = ln.get('w')
        solidFill = ln.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is None:
            return None, lnW
        srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            return srgbClr.get('val', '').upper(), lnW
        # schemeClr
        schemeClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
        if schemeClr is not None:
            return f"scheme:{schemeClr.get('val')}", lnW
        return None, lnW
    except Exception as e:
        return None, None

def main():
    prs = Presentation(PPTX_PATH)
    slides = prs.slides

    results = {
        "verdict": "PASS",
        "confidence": 1.0,
        "checks": {},
        "issues": [],
        "fix_required": []
    }

    # ---- Slide idx5 (0-based index 5, i.e. slide 6) ----
    slide5 = slides[5]  # idx5
    slide6 = slides[6]  # idx6

    # ==============================
    # CHECK 1: idx5 TextBox 5 & 6 padding
    # ==============================
    check1 = {
        "name": "idx5_textbox_padding",
        "pass": True,
        "details": {}
    }

    # Find shapes by name
    tb5 = None
    tb6 = None
    for shape in slide5.shapes:
        if shape.name == "TextBox 5":
            tb5 = shape
        elif shape.name == "TextBox 6":
            tb6 = shape

    if tb5 is None:
        check1["pass"] = False
        check1["details"]["TextBox 5"] = "NOT FOUND"
        results["issues"].append("idx5: TextBox 5 not found")
    else:
        lIns, tIns = get_bodyPr_padding(tb5._element)
        has_text = bool(tb5.text_frame.text.strip()) if tb5.has_text_frame else False
        tb5_ok = lIns >= MIN_PADDING_EMU and tIns >= MIN_PADDING_EMU and has_text
        check1["details"]["TextBox 5"] = {
            "lIns": lIns,
            "tIns": tIns,
            "lIns_ok": lIns >= MIN_PADDING_EMU,
            "tIns_ok": tIns >= MIN_PADDING_EMU,
            "has_text": has_text,
            "pass": tb5_ok
        }
        if not tb5_ok:
            check1["pass"] = False
            if lIns < MIN_PADDING_EMU:
                results["issues"].append(f"idx5: TextBox 5 lIns={lIns} < {MIN_PADDING_EMU}")
                results["fix_required"].append("idx5: TextBox 5 lIns needs >= 91440 EMU")
            if tIns < MIN_PADDING_EMU:
                results["issues"].append(f"idx5: TextBox 5 tIns={tIns} < {MIN_PADDING_EMU}")
                results["fix_required"].append("idx5: TextBox 5 tIns needs >= 91440 EMU")
            if not has_text:
                results["issues"].append("idx5: TextBox 5 is empty")

    if tb6 is None:
        check1["pass"] = False
        check1["details"]["TextBox 6"] = "NOT FOUND"
        results["issues"].append("idx5: TextBox 6 not found")
    else:
        lIns, tIns = get_bodyPr_padding(tb6._element)
        has_text = bool(tb6.text_frame.text.strip()) if tb6.has_text_frame else False
        tb6_ok = lIns >= MIN_PADDING_EMU and has_text
        check1["details"]["TextBox 6"] = {
            "lIns": lIns,
            "tIns": tIns,
            "lIns_ok": lIns >= MIN_PADDING_EMU,
            "has_text": has_text,
            "pass": tb6_ok
        }
        if not tb6_ok:
            check1["pass"] = False
            if lIns < MIN_PADDING_EMU:
                results["issues"].append(f"idx5: TextBox 6 lIns={lIns} < {MIN_PADDING_EMU}")
                results["fix_required"].append("idx5: TextBox 6 lIns needs >= 91440 EMU")
            if not has_text:
                results["issues"].append("idx5: TextBox 6 is empty")

    results["checks"]["check1_padding"] = check1

    # ==============================
    # CHECK 2: idx6 TextBox 2 line break
    # ==============================
    check2 = {
        "name": "idx6_textbox2_linebreak",
        "pass": True,
        "details": {}
    }

    tb2_slide6 = None
    for shape in slide6.shapes:
        if shape.name == "TextBox 2":
            tb2_slide6 = shape
            break

    if tb2_slide6 is None:
        check2["pass"] = False
        check2["details"]["TextBox 2"] = "NOT FOUND"
        results["issues"].append("idx6: TextBox 2 not found")
    else:
        paras = get_text_paragraphs(tb2_slide6)
        full_text = "\n".join(paras)
        expected_para1 = "L02. Three"
        expected_para2 = "Cards"

        # Check exactly 2 paragraphs with correct content
        # Filter out empty trailing paragraphs
        non_empty_paras = [p for p in paras if p.strip()]

        para1_ok = len(non_empty_paras) >= 1 and non_empty_paras[0].strip() == expected_para1
        para2_ok = len(non_empty_paras) >= 2 and non_empty_paras[1].strip() == expected_para2

        # Check for mid-word breaks (bad patterns)
        bad_break = False
        for p in paras:
            if "Card\n" in p or p == "s" or "Three Card" in p:
                bad_break = True

        tb2_ok = para1_ok and para2_ok and not bad_break
        check2["details"]["TextBox 2"] = {
            "paragraphs": paras,
            "non_empty_paragraphs": non_empty_paras,
            "para1_ok": para1_ok,
            "para2_ok": para2_ok,
            "bad_break": bad_break,
            "pass": tb2_ok
        }
        if not tb2_ok:
            check2["pass"] = False
            if not para1_ok:
                results["issues"].append(f"idx6: TextBox 2 para1='{non_empty_paras[0] if non_empty_paras else 'EMPTY'}' expected 'L02. Three'")
                results["fix_required"].append("idx6: TextBox 2 first paragraph must be 'L02. Three'")
            if not para2_ok:
                results["issues"].append(f"idx6: TextBox 2 para2='{non_empty_paras[1] if len(non_empty_paras)>1 else 'MISSING'}' expected 'Cards'")
                results["fix_required"].append("idx6: TextBox 2 second paragraph must be 'Cards'")
            if bad_break:
                results["issues"].append("idx6: TextBox 2 has mid-word break")
                results["fix_required"].append("idx6: TextBox 2 must not have mid-word break")

    results["checks"]["check2_linebreak"] = check2

    # ==============================
    # CHECK 3: idx6 3 card backgrounds have PRIMARY border
    # ==============================
    check3 = {
        "name": "idx6_card_border",
        "pass": True,
        "details": {}
    }

    # Find Rounded Rectangle 15, 18, 21 (large shapes at top > 4.0", h > 2.0")
    target_names = {"Rounded Rectangle 15", "Rounded Rectangle 18", "Rounded Rectangle 21"}
    found_shapes = {}

    for shape in slide6.shapes:
        if shape.name in target_names:
            found_shapes[shape.name] = shape

    # Also find by size criteria as fallback
    large_shapes_by_size = []
    for shape in slide6.shapes:
        top_inches = emu_to_inches(shape.top)
        height_inches = emu_to_inches(shape.height)
        if top_inches > 0.5 and height_inches > 2.0 and shape.name not in found_shapes:
            large_shapes_by_size.append(shape)

    check3["details"]["found_by_name"] = list(found_shapes.keys())
    check3["details"]["large_shapes_by_size"] = [s.name for s in large_shapes_by_size]

    if len(found_shapes) < 3:
        # Try to use size-based detection
        missing = target_names - set(found_shapes.keys())
        check3["details"]["missing_named"] = list(missing)
        results["issues"].append(f"idx6: Could not find by name: {missing}")

    border_results = {}
    shapes_to_check = list(found_shapes.values())

    for shape in shapes_to_check:
        color, lnW = get_shape_color(shape)
        lnW_val = int(lnW) if lnW is not None else 0
        color_ok = color == PRIMARY_COLOR if color else False
        width_ok = lnW_val > 0
        border_ok = color_ok and width_ok

        border_results[shape.name] = {
            "color": color,
            "line_width_emu": lnW_val,
            "color_ok": color_ok,
            "width_ok": width_ok,
            "pass": border_ok
        }

        if not border_ok:
            check3["pass"] = False
            if not color_ok:
                results["issues"].append(f"idx6: {shape.name} border color={color}, expected #{PRIMARY_COLOR}")
                results["fix_required"].append(f"idx6: {shape.name} border must be #{PRIMARY_COLOR}")
            if not width_ok:
                results["issues"].append(f"idx6: {shape.name} border width=0")
                results["fix_required"].append(f"idx6: {shape.name} border width must be > 0")

    if len(found_shapes) < 3:
        check3["pass"] = False

    check3["details"]["border_checks"] = border_results
    results["checks"]["check3_border"] = check3

    # ==============================
    # CHECK 4: No overflow on idx5, idx6
    # ==============================
    check4 = {
        "name": "no_overflow",
        "pass": True,
        "details": {}
    }

    for slide_idx, slide in [(5, slide5), (6, slide6)]:
        overflow_shapes = []
        for shape in slide.shapes:
            bottom = shape.top + shape.height
            bottom_inches = emu_to_inches(bottom)
            if bottom_inches > 7.0:
                overflow_shapes.append({
                    "name": shape.name,
                    "top_inches": round(emu_to_inches(shape.top), 3),
                    "height_inches": round(emu_to_inches(shape.height), 3),
                    "bottom_inches": round(bottom_inches, 3)
                })

        check4["details"][f"idx{slide_idx}"] = {
            "overflow_shapes": overflow_shapes,
            "pass": len(overflow_shapes) == 0
        }

        if overflow_shapes:
            check4["pass"] = False
            for s in overflow_shapes:
                results["issues"].append(f"idx{slide_idx}: {s['name']} bottom={s['bottom_inches']:.3f}\" > 7.0\"")
                results["fix_required"].append(f"idx{slide_idx}: {s['name']} must fit within 7.0\"")

    results["checks"]["check4_overflow"] = check4

    # ==============================
    # Final verdict
    # ==============================
    all_checks_pass = all(
        results["checks"][k]["pass"]
        for k in results["checks"]
    )

    if not all_checks_pass:
        results["verdict"] = "FAIL"
        results["confidence"] = 0.95
    else:
        results["verdict"] = "PASS"
        results["confidence"] = 0.97

    print(json.dumps(results, indent=2, ensure_ascii=False))
    return 0 if results["verdict"] == "PASS" else 1

if __name__ == "__main__":
    sys.exit(main())
