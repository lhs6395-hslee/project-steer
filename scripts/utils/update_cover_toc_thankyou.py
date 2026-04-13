"""템플릿의 표지, 목차, 끝맺음 페이지 텍스트를 DMS Expert Guide용으로 수정."""

from pptx import Presentation

prs = Presentation("templates/pptx_template.pptx")

# ── 슬라이드 0: 표지 ──
slide0 = prs.slides[0]

# shape[2]: 제목 "GS Neotek \nPPT 템플릿 가이드라인" → "AWS Database\nMigration Service"
s2 = slide0.shapes[2]
for para in s2.text_frame.paragraphs:
    for run in para.runs:
        run.text = ""
s2.text_frame.paragraphs[0].runs[0].text = "AWS Database"
if len(s2.text_frame.paragraphs) > 1:
    s2.text_frame.paragraphs[1].runs[0].text = "Migration Service"

# shape[3]: 회사 소개 좌 — 그대로 유지 (고객의 Digitalization...)
# shape[4]: 회사 소개 우 — 그대로 유지 (AI, Cloud, BigData...)

# shape[5]: 부제 "부제목을 적어주세요" → DMS 부제
s5 = slide0.shapes[5]
for para in s5.text_frame.paragraphs:
    for run in para.runs:
        run.text = ""
s5.text_frame.paragraphs[0].runs[0].text = "Expert Guide — 아키텍처 · 설정 · 운영 · 트러블슈팅"

# shape[7]: 연도 "2025" → "2026"
s7 = slide0.shapes[7]
for para in s7.text_frame.paragraphs:
    for run in para.runs:
        if "2025" in run.text:
            run.text = run.text.replace("2025", "2026")

# shape[8]: 날짜 "00/00" → "04/07"
s8 = slide0.shapes[8]
for para in s8.text_frame.paragraphs:
    for run in para.runs:
        if "00/00" in run.text:
            run.text = run.text.replace("00/00", "04/07")

print("✅ 표지 수정 완료")

# ── 슬라이드 1: 목차 ──
slide1 = prs.slides[1]

# shape[1]: 번호 열 — "1\n2\n3\n4\n5" 유지 (5개 섹션 그룹)
# shape[2]: 섹션명 열 — 교체
sections = [
    "아키텍처 & 핵심 구성요소",
    "Replication Instance & Endpoint",
    "Migration Task & Table Mapping",
    "CDC 운영 & 모니터링",
    "트러블슈팅 & Best Practices",
]

s_names = slide1.shapes[2]
for pi, para in enumerate(s_names.text_frame.paragraphs):
    if pi < len(sections):
        # 첫 번째 run에 텍스트 설정, 나머지 run 비우기
        for ri, run in enumerate(para.runs):
            if ri == 0:
                run.text = sections[pi]
            else:
                run.text = ""

print("✅ 목차 수정 완료")

# ── 마지막 슬라이드: 끝맺음 ──
# 템플릿에서 끝맺음 레이아웃 슬라이드 찾기 (python-pptx 내부 이름: "감사합니다")
# 주의: slide.slide_layout.name은 .pptx 파일 내부의 실제 레이아웃 이름이므로 변경 불가
thank_idx = None
for i, slide in enumerate(prs.slides):
    if slide.slide_layout.name == "감사합니다":
        thank_idx = i  # 마지막 끝맺음 슬라이드

if thank_idx is not None:
    slide_ty = prs.slides[thank_idx]
    # shape[0]: "Thank You" — 그대로 유지
    # shape[2]: 태그라인 — 그대로 유지 (GS네오텍 브랜드)
    print(f"✅ 끝맺음 슬라이드 확인 (index {thank_idx}) — 텍스트 유지")
else:
    print("⚠️ 끝맺음 슬라이드 없음")

# 저장
prs.save("results/pptx/AWS_DMS_Expert_Guide.pptx")
print("✅ 저장 완료: results/pptx/AWS_DMS_Expert_Guide.pptx")
