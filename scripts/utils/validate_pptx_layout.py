#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Layout Validation Script
Validates layout-specific structural compliance (L01-L05) and all constraints
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
import sys
import json

# Constants from style guide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BODY_START_Y = Inches(2.0)
BODY_LIMIT_Y = Inches(7.0)
MIN_BOTTOM_MARGIN = Inches(0.3)

# EMU to Inches
EMU_TO_INCHES = 914400

def emu_to_inches(emu):
    return emu / EMU_TO_INCHES

def rgb_to_tuple(rgb_color):
    """Extract RGB tuple from color object"""
    try:
        if hasattr(rgb_color, 'rgb'):
            return (rgb_color.rgb[0], rgb_color.rgb[1], rgb_color.rgb[2])
    except:
        pass
    return None

def is_bright_background(fill):
    """Check if fill is bright (#F8F9FA, #FFFFFF, etc)"""
    if fill.type != MSO_FILL_TYPE.SOLID:
        return False
    
    rgb = rgb_to_tuple(fill.fore_color)
    if not rgb:
        return False
    
    # Bright colors: F8F9FA (248,249,250), FFFFFF (255,255,255)
    r, g, b = rgb
    return r >= 240 and g >= 240 and b >= 240

def is_white_text(font):
    """Check if text is white"""
    try:
        if font.color.type == 1:  # RGB
            rgb = rgb_to_tuple(font.color)
            if rgb:
                r, g, b = rgb
                return r >= 250 and g >= 250 and b >= 250
    except:
        pass
    return False

def estimate_text_width_pt(text, font_size_pt):
    """Estimate text width in points"""
    width = 0
    for ch in text:
        if ord(ch) > 0x2E80:  # CJK characters
            width += font_size_pt
        elif ch == ' ':
            width += font_size_pt * 0.25
        elif ch in '.:/-()':
            width += font_size_pt * 0.35
        else:  # ASCII
            width += font_size_pt * 0.5
    return width

def count_lines(text):
    """Count lines in text"""
    return len(text.split('\n'))

def identify_layout(slide):
    """Identify which layout type (L01-L05) a slide uses"""
    # Count shapes and analyze structure
    shapes = [s for s in slide.shapes if s.shape_type in [
        MSO_SHAPE_TYPE.AUTO_SHAPE, 
        MSO_SHAPE_TYPE.TEXT_BOX,
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.GROUP
    ]]
    
    # Look for subtitle text to identify layout
    subtitle_text = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Subtitle is typically at top, small font, right side
            top_inches = emu_to_inches(shape.top)
            left_inches = emu_to_inches(shape.left)
            if 0.5 < top_inches < 1.5 and left_inches > 5:
                if text and len(text) < 100:
                    subtitle_text = text
                    break
    
    # Heuristic layout detection based on slide structure
    # This is simplified - real detection would need more analysis
    card_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    
    if len(card_shapes) == 3 and any('Three' in str(subtitle_text) or 'Cards' in str(subtitle_text) for _ in [1]):
        return 'L02'
    elif len(card_shapes) == 4:
        return 'L03'
    elif 'Arrow' in str(subtitle_text) or 'Process' in str(subtitle_text):
        return 'L04'
    elif 'Column' in str(subtitle_text) or 'Phased' in str(subtitle_text):
        return 'L05'
    elif 'Bento' in str(subtitle_text):
        return 'L01'
    
    return 'UNKNOWN'

def validate_slide(slide_idx, slide):
    """Validate a single slide against all constraints"""
    violations = []
    layout_type = identify_layout(slide)
    
    # Find title and subtitle shapes
    title_shape = None
    subtitle_shape = None
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        text = shape.text_frame.text.strip()
        top_inches = emu_to_inches(shape.top)
        left_inches = emu_to_inches(shape.left)
        
        # Title: left side, top
        if 0.5 < top_inches < 1.5 and left_inches < 3 and text:
            if not title_shape or emu_to_inches(title_shape.top) > top_inches:
                title_shape = shape
        
        # Subtitle: right side, top
        if 0.5 < top_inches < 1.5 and left_inches > 5 and text:
            if not subtitle_shape or emu_to_inches(subtitle_shape.top) > top_inches:
                subtitle_shape = shape
    
    # Validate title width (~340pt at 28pt font)
    if title_shape and title_shape.has_text_frame:
        title_text = title_shape.text_frame.text.strip()
        if title_text:
            # Estimate at 28pt
            width_pt = estimate_text_width_pt(title_text, 28)
            if width_pt > 340:
                violations.append(f"Title overflow: {width_pt:.0f}pt > 340pt limit | '{title_text}'")
    
    # Validate subtitle: max 2 lines, no mid-word breaks
    if subtitle_shape and subtitle_shape.has_text_frame:
        subtitle_text = subtitle_shape.text_frame.text.strip()
        if subtitle_text:
            line_count = count_lines(subtitle_text)
            if line_count > 2:
                violations.append(f"Subtitle exceeds 2 lines: {line_count} lines | '{subtitle_text}'")
            
            # Check for mid-word breaks
            lines = subtitle_text.split('\n')
            for i, line in enumerate(lines[:-1]):  # All but last line
                if line and not line[-1].isspace():
                    # Check if next line starts with lowercase (mid-word break)
                    next_line = lines[i+1].strip()
                    if next_line and next_line[0].islower():
                        violations.append(f"Mid-word line break in subtitle: '{line}' / '{next_line}'")
    
    # Validate content positioning (2.0" - 7.0" range)
    content_shapes = []
    for shape in slide.shapes:
        top_inches = emu_to_inches(shape.top)
        bottom_inches = emu_to_inches(shape.top + shape.height)
        
        # Skip header area
        if bottom_inches < 2.0:
            continue
        
        content_shapes.append(shape)
        
        # Check if content exceeds body area
        if bottom_inches > 7.0:
            violations.append(f"Content exceeds 7.0\" limit: bottom at {bottom_inches:.2f}\"")
    
    # Check bottom margin (should be at least 0.3" from 7.5")
    if content_shapes:
        max_bottom = max(emu_to_inches(s.top + s.height) for s in content_shapes)
        bottom_margin = 7.5 - max_bottom
        if bottom_margin < 0.3:
            violations.append(f"Insufficient bottom margin: {bottom_margin:.2f}\" < 0.3\" minimum")
    
    # Check for white text on bright backgrounds (CRITICAL FAIL)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # Check shape fill
        shape_fill = None
        try:
            if hasattr(shape, 'fill'):
                shape_fill = shape.fill
        except:
            pass
        
        if shape_fill and is_bright_background(shape_fill):
            # Check text color
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if is_white_text(run.font):
                        violations.append(f"CRITICAL: WHITE text on bright background in shape at ({emu_to_inches(shape.left):.1f}\", {emu_to_inches(shape.top):.1f}\")")
    
    # Check for blue oval fallback icons (CRITICAL FAIL)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's an oval
            try:
                if hasattr(shape, 'auto_shape_type') and 'OVAL' in str(shape.auto_shape_type):
                    # Check if it's blue fill
                    if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL_TYPE.SOLID:
                        rgb = rgb_to_tuple(shape.fill.fore_color)
                        if rgb and rgb[2] > 200 and rgb[0] < 100:  # Blue-ish
                            violations.append(f"CRITICAL: Blue oval fallback icon detected at ({emu_to_inches(shape.left):.1f}\", {emu_to_inches(shape.top):.1f}\")")
            except:
                pass
    
    # Layout-specific validation
    if layout_type == 'L02':
        # L02 Three Cards: Top diagram zone REQUIRED
        has_top_diagram = False
        for shape in slide.shapes:
            top_inches = emu_to_inches(shape.top)
            if 2.0 <= top_inches <= 3.0:  # Top zone after header
                # Check if it's a diagram (group, multiple connected shapes, or image)
                if shape.shape_type in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE]:
                    has_top_diagram = True
                    break
        
        if not has_top_diagram:
            violations.append(f"L02 FAIL: Top diagram zone REQUIRED but not found")
    
    elif layout_type == 'L04':
        # L04 Process Arrow: Check for bottom-right icon if no diagram
        has_diagram = any(s.shape_type in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE] 
                         and emu_to_inches(s.top) < 4.0 for s in slide.shapes)
        
        if not has_diagram:
            # Check for bottom-right icon (10.5", 6.2", 0.60" × 0.60")
            has_br_icon = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    left_inches = emu_to_inches(shape.left)
                    top_inches = emu_to_inches(shape.top)
                    width_inches = emu_to_inches(shape.width)
                    height_inches = emu_to_inches(shape.height)
                    
                    if (10.0 <= left_inches <= 11.0 and 
                        6.0 <= top_inches <= 6.5 and
                        0.5 <= width_inches <= 0.7 and
                        0.5 <= height_inches <= 0.7):
                        has_br_icon = True
                        break
            
            if not has_br_icon:
                violations.append(f"L04 FAIL: No diagram present, but bottom-right icon missing (10.5\", 6.2\", 0.60\" × 0.60\")")
    
    elif layout_type == 'L05':
        # L05 Phased Columns: Check for bottom-right icon if no diagram
        has_diagram = any(s.shape_type in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE] 
                         and emu_to_inches(s.top) < 4.0 for s in slide.shapes)
        
        if not has_diagram:
            # Check for bottom-right icon
            has_br_icon = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    left_inches = emu_to_inches(shape.left)
                    top_inches = emu_to_inches(shape.top)
                    width_inches = emu_to_inches(shape.width)
                    height_inches = emu_to_inches(shape.height)
                    
                    if (10.0 <= left_inches <= 11.0 and 
                        6.0 <= top_inches <= 6.5 and
                        0.5 <= width_inches <= 0.7 and
                        0.5 <= height_inches <= 0.7):
                        has_br_icon = True
                        break
            
            if not has_br_icon:
                violations.append(f"L05 FAIL: No diagram present, but bottom-right icon missing")
        
        # Check column color diversity
        # (This requires more complex analysis of fill colors across shapes)
    
    elif layout_type == 'L03':
        # L03 Grid 2x2: Check for actual PNG icons (not blue ovals)
        icon_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                top_inches = emu_to_inches(shape.top)
                if 2.5 <= top_inches <= 6.5:  # Body area
                    icon_count += 1
        
        if icon_count < 4:
            violations.append(f"L03 FAIL: Expected 4 PNG icons in grid, found {icon_count}")
    
    return {
        'slide_idx': slide_idx,
        'layout_type': layout_type,
        'violations': violations
    }

def main():
    pptx_path = 'results/pptx/AWS_MSK_Expert_Intro.pptx'
    
    try:
        prs = Presentation(pptx_path)
        print(f"Analyzing: {pptx_path}")
        print(f"Total slides: {len(prs.slides)}")
        print(f"Slide dimensions: {emu_to_inches(prs.slide_width):.2f}\" × {emu_to_inches(prs.slide_height):.2f}\"")
        print()
        
        all_violations = []
        results = []
        
        for idx, slide in enumerate(prs.slides):
            result = validate_slide(idx, slide)
            results.append(result)
            
            if result['violations']:
                all_violations.extend(result['violations'])
                print(f"Slide {idx + 1} [{result['layout_type']}]:")
                for violation in result['violations']:
                    print(f"  ❌ {violation}")
                print()
        
        # Summary
        print("=" * 80)
        print(f"VALIDATION SUMMARY")
        print("=" * 80)
        print(f"Total slides analyzed: {len(prs.slides)}")
        print(f"Slides with violations: {sum(1 for r in results if r['violations'])}")
        print(f"Total violations: {len(all_violations)}")
        print()
        
        if all_violations:
            print("CONSTRAINT COMPLIANCE: FAILED")
            print()
            print("All violations:")
            for i, v in enumerate(all_violations, 1):
                print(f"{i}. {v}")
        else:
            print("✅ CONSTRAINT COMPLIANCE: PASSED")
        
        # Output JSON for executor
        output = {
            'total_slides': len(prs.slides),
            'slides_with_violations': sum(1 for r in results if r['violations']),
            'total_violations': len(all_violations),
            'results': results,
            'status': 'PASS' if not all_violations else 'FAIL'
        }
        
        with open('results/pptx/validation_report.json', 'w', encoding='utf-8') as f:
            json.dump(output, f, indent=2, ensure_ascii=False)
        
        print()
        print(f"Detailed report saved to: results/pptx/validation_report.json")
        
        return 0 if not all_violations else 1
        
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
        return 2

if __name__ == '__main__':
    sys.exit(main())
